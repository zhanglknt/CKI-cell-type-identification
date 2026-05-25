"""
65_fix_s3_overlap.py — 修复S3中Exhausted状态与底部绿色色条重叠
根因: SVG viewBox=1800×560, Exhausted行底部530→与绿色条(y=500)重叠30px
修复: 扩展viewBox到600, 绿色条移至y=545, 同步调整PPTX图片高度
"""

from pptx import Presentation
from pptx.util import Inches, Emu
import re

SVG_PATH = "results/figures_final/_s2_cell_states_en_v2.svg"

# ============================================================
# Step 1: 修复 SVG
# ============================================================
with open(SVG_PATH, "r", encoding="utf-8") as f:
    svg = f.read()

# 扩展 viewBox: 560 → 600
svg = svg.replace('viewBox="0 0 1800 560"', 'viewBox="0 0 1800 600"')

# 移动底部绿色色条: y=500 → y=545
old_bar = '<rect x="20" y="500" width="1760" height="44" rx="10" fill="#0D9488"/>'
new_bar = '<rect x="20" y="545" width="1760" height="44" rx="10" fill="#0D9488"/>'
svg = svg.replace(old_bar, new_bar)

# 移动色条内文字: y=528 → y=573
old_text = '<text x="900" y="528" text-anchor="middle" font-size="14" font-weight="700" fill="#FFFFFF" font-family="Arial,Helvetica,sans-serif">'
new_text = '<text x="900" y="573" text-anchor="middle" font-size="14" font-weight="700" fill="#FFFFFF" font-family="Arial,Helvetica,sans-serif">'
svg = svg.replace(old_text, new_text)

with open(SVG_PATH, "w", encoding="utf-8") as f:
    f.write(svg)

print(f"[SVG] viewBox 560→600, green bar y=500→545, text y=528→573")
print(f"      Exhausted bottom=530, bar top=545 → gap=15px ✓")

# ============================================================
# Step 2: 修复 PPTX (EN + ZH)
# ============================================================
# 新图片高宽比: 1800/600 = 3.0, 在 9.0" 宽时自然高=3.0"
# 但 PPTX 布局空间有限，使用 2.85" 高度（轻微压缩 ~5%，可接受）
# image bottom: 1.35 + 2.85 = 4.20
# insight: T=4.40 → gap=0.20"

NEW_IMAGE_H = Inches(2.85)
NEW_INSIGHT_T = Inches(4.40)

for pptx_path in [
    "results/figures_final/CKI_Lecture_2026_v4.pptx",
    "results/figures_final/CKI_Lecture_2026_v4_ZH.pptx",
]:
    print(f"\n[PPTX] {pptx_path.split('/')[-1]}")
    prs = Presentation(pptx_path)
    slide = prs.slides[2]  # S3

    for s in slide.shapes:
        if hasattr(s, 'image'):
            old_h = s.height / 914400
            s.height = NEW_IMAGE_H
            print(f"  Image: H={old_h:.2f}→{NEW_IMAGE_H/914400:.2f}")

        elif hasattr(s, 'text_frame'):
            text = s.text_frame.text[:30]
            if "Key Insight" in text or "关键洞察" in text:
                old_t = s.top / 914400
                s.top = NEW_INSIGHT_T
                print(f"  Insight: T={old_t:.2f}→{NEW_INSIGHT_T/914400:.2f}")

    prs.save(pptx_path)

print("\nDone. Final S3 layout:")
print("  Title:    0.15 — 0.70")
print("  Subtitle: 0.85 — 1.10")
print("  Image:    1.35 — 4.20  (H=2.85, 1800×600 SVG)")
print("  Insight:  4.40 — 5.25  (gap 0.20\")")
