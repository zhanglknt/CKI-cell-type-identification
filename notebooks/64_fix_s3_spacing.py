"""
64_fix_s3_spacing.py — 修复S3元素重叠，增加各层间距
问题: 标题 H=0.70" 底部=0.85" 与副标题 T=0.75" 重叠0.10"
方案: 标题→副标题→图片→洞察 四层清晰间隔，调整尺寸
"""

from pptx import Presentation
from pptx.util import Inches, Emu
from copy import deepcopy
import copy

TARGETS = [
    "results/figures_final/CKI_Lecture_2026_v4.pptx",
    "results/figures_final/CKI_Lecture_2026_v4_ZH.pptx",
]

# 新尺寸配置
NEW_TITLE = dict(top=Inches(0.15), height=Inches(0.55))          # bottom=0.70
NEW_SUBTITLE = dict(top=Inches(0.85), height=Inches(0.25))       # bottom=1.10
NEW_IMAGE = dict(left=Inches(0.5), top=Inches(1.35), width=Inches(9.0), height=Inches(2.55))
# image bottom=3.90
NEW_INSIGHT = dict(top=Inches(4.15), height=Inches(0.85))        # bottom=5.00

for pptx_path in TARGETS:
    print(f"\n--- {pptx_path} ---")
    prs = Presentation(pptx_path)
    slide = prs.slides[2]  # S3 = index 2

    for s in slide.shapes:
        if hasattr(s, 'text_frame'):
            text = s.text_frame.text[:40]
            if "Same Cell Type" in text or "同一细胞类型" in text:
                # Title
                old_t, old_h = s.top/914400, s.height/914400
                print(f"  Title: T={old_t:.2f}→{NEW_TITLE['top']/914400:.2f}, H={old_h:.2f}→{NEW_TITLE['height']/914400:.2f}")
                s.top = NEW_TITLE['top']
                s.height = NEW_TITLE['height']

            elif "CKI captures" in text or "CKI 捕获" in text:
                # Subtitle
                old_t, old_h = s.top/914400, s.height/914400
                print(f"  Subtitle: T={old_t:.2f}→{NEW_SUBTITLE['top']/914400:.2f}, H={old_h:.2f}→{NEW_SUBTITLE['height']/914400:.2f}")
                s.top = NEW_SUBTITLE['top']
                s.height = NEW_SUBTITLE['height']

            elif "Key Insight" in text or "关键洞察" in text:
                # Insight text box
                old_t, old_h = s.top/914400, s.height/914400
                print(f"  Insight: T={old_t:.2f}→{NEW_INSIGHT['top']/914400:.2f}, H={old_h:.2f}→{NEW_INSIGHT['height']/914400:.2f}")
                s.top = NEW_INSIGHT['top']
                s.height = NEW_INSIGHT['height']

        elif hasattr(s, 'image'):
            # The cell state diagram
            old_l, old_t, old_w, old_h = s.left/914400, s.top/914400, s.width/914400, s.height/914400
            print(f"  Image: L={old_l:.1f}→{NEW_IMAGE['left']/914400:.1f}, T={old_t:.1f}→{NEW_IMAGE['top']/914400:.1f}, W={old_w:.1f}→{NEW_IMAGE['width']/914400:.1f}, H={old_h:.1f}→{NEW_IMAGE['height']/914400:.1f}")
            s.left = NEW_IMAGE['left']
            s.top = NEW_IMAGE['top']
            s.width = NEW_IMAGE['width']
            s.height = NEW_IMAGE['height']

    prs.save(pptx_path)
    print(f"  SAVED")

print("\nDone. Layout fix applied:")
print("  Title:    0.15 — 0.70  (gap 0.15\")")
print("  Subtitle: 0.85 — 1.10  (gap 0.25\")")
print("  Image:    1.35 — 3.90  (gap 0.25\")")
print("  Insight:  4.15 — 5.00  (bottom margin 0.625\")")
