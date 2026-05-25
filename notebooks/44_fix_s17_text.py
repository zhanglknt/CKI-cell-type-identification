"""
Fix Slide 17 body text for clarity (both EN and ZH).
Precise text replacement for the two text cards.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

BASE = r"C:\Users\KnightZ\Desktop\细胞受选择\results\figures_final"

EN_NEW_TEXT = {
    # Left card: Omega Gradient body
    "left": (
        "Astrocytes show omega = 107.5 in cortical regions. Bergmann glia (cerebellum) form a distinct cluster.\n"
        "Thalamus and striatum show moderate omega — hippocampus and cortex are most similar.\n"
        "The omega landscape provides a gradient-based quantification of brain cell identity."
    ),
    # Right card: Migration Hypothesis body
    "right": (
        "OPCs migrate from SVZ (omega ~3) to cortex (omega ~22).\n"
        "CKI tracks this path via the omega gradient — higher omega = stronger functional differentiation downstream.\n"
        "This gradient generates testable migration hypotheses, linking spatial location to differentiation strength."
    ),
}

ZH_NEW_TEXT = {
    "left": (
        "星形胶质细胞：皮层 omega = 107.5。Bergmann胶质（小脑）独立成簇。\n"
        "丘脑和纹状体展现中等 omega，海马体与皮层最为相似。\n"
        "omega梯度量化了脑区细胞身份——为细胞特化提供了基于梯度的视角。"
    ),
    "right": (
        "OPC从SVZ (omega ~3) 迁移至皮层 (omega ~22)。\n"
        "CKI通过omega梯度追踪此路径——omega越高 = 功能分化越强。\n"
        "该梯度生成可检验的迁移假说，将空间位置与分化强度关联起来。"
    ),
}


def set_paragraph_text(shape, lines):
    """Replace all paragraphs in a text frame with given lines."""
    tf = shape.text_frame
    # Clear all paragraphs except first
    for para in tf.paragraphs[1:]:
        pPr = para._p.getparent()
        pPr.remove(para._p)
    
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        
        # Clear existing runs
        for run in para.runs:
            run._r.getparent().remove(run._r)
        
        run = para.add_run()
        run.text = line
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        run.font.name = "Arial"


def fix_slide17_text(prs, new_texts, lang):
    """Fix body text for both cards on slide 17."""
    slide = prs.slides[16]
    
    # Find the two text body shapes (largest text shapes with multi-line content on Slide 17)
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if len(text) > 100 and ("omega" in text.lower() or "astro" in text.lower() or "OPC" in text or "星形" in text or "少突" in text):
            candidates.append(shape)
    
    if len(candidates) >= 2:
        # Left is the one mentioning astrocytes/Bergmann
        # Right is the one mentioning OPC/migration
        for shape in candidates:
            text = shape.text_frame.text
            if "astrocyte" in text or "Bergmann" in text or "星形" in text:
                lines = new_texts["left"].split('\n')
                set_paragraph_text(shape, lines)
                print(f"  [{lang}] Slide 17 left card: updated")
            elif "OPC" in text or "少突" in text:
                lines = new_texts["right"].split('\n')
                set_paragraph_text(shape, lines)
                print(f"  [{lang}] Slide 17 right card: updated")
    else:
        print(f"  [{lang}] Slide 17: WARNING - found {len(candidates)} candidate shapes, expected 2")


for lang, ppt_file, texts in [
    ("EN", "CKI_Lecture_2026_v3.pptx", EN_NEW_TEXT),
    ("ZH", "CKI_Lecture_2026_v3_ZH.pptx", ZH_NEW_TEXT),
]:
    ppt_path = os.path.join(BASE, ppt_file)
    print(f"\n=== {lang}: {ppt_file} ===")
    prs = Presentation(ppt_path)
    fix_slide17_text(prs, texts, lang)
    prs.save(ppt_path)
    print(f"  [{lang}] Saved")

print("\nDone!")
