"""
57_fix_s16_s2_data.py — Fix S16 fabricated data + S2 broken SVG
================================================================
S16: Replace fabricated omega values with real NAR Figure 6 Panel B data (10 cell classes, 7.6x gradient)
S2:  Fix SVG fill format (parentheses → #hex) and redesign cell-state diagram
"""
import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path(__file__).resolve().parent.parent / "results" / "figures_final"

# ── Colors ──
C_BLUE    = (0x1D, 0x4E, 0xD8)
C_TEAL    = (0x0D, 0x94, 0x88)
C_RED     = (0xDC, 0x26, 0x26)
C_AMBER   = (0xF5, 0x9E, 0x0B)
C_PURPLE  = (0x7C, 0x3A, 0xED)
C_GREEN   = (0x05, 0x96, 0x69)
C_DARK_BG = (0x0F, 0x17, 0x2A)
C_WHITE   = (0xFF, 0xFF, 0xFF)
C_MUTED   = (0x64, 0x74, 0x8B)
BODY_FG   = (0x1E, 0x29, 0x3B)
TITLE_FG  = (0x1E, 0x29, 0x3B)
CARD_BG   = (0xFF, 0xFF, 0xFF)


def fully_clear_slide(slide):
    """Remove ALL shapes from a slide by clearing the spTree XML."""
    from lxml import etree
    spTree = slide.shapes._spTree
    to_remove = []
    for child in spTree:
        if child.tag.endswith('}sp') or child.tag.endswith('}pic') or child.tag.endswith('}grpSp') or child.tag.endswith('}graphicFrame') or child.tag.endswith('}cxnSp'):
            to_remove.append(child)
    for child in to_remove:
        spTree.remove(child)


def set_text(shape, text, size=10, bold=False, color=BODY_FG, align=None):
    """Helper to set text on a shape's text_frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


# ═══════════════════════════════════════════════════════════════
# S16: Brain Regional CKI — CORRECT data from NAR Figure 6
# ═══════════════════════════════════════════════════════════════

# Real data from 30_nar_figures_final.py Figure 6 Panel B
NAR_BRAIN_DATA = [
    # (label, omega, color, is_highlight)
    ("Astrocyte",        121.77, C_RED,    True),   # highest
    ("OPC",               55.42, C_AMBER,  False),
    ("Choroid plexus",    45.47, C_AMBER,  False),
    ("Oligodendrocyte",   44.87, C_AMBER,  False),
    ("COP",               40.03, C_BLUE,   False),
    ("Ependymal",         39.53, C_BLUE,   False),
    ("Microglia",         35.37, C_TEAL,   False),
    ("Fibroblast",        25.41, C_TEAL,   False),
    ("Vascular",          21.54, C_GREEN,  False),
    ("Bergmann glia",     15.97, C_GREEN,  True),   # lowest
]

ZH_BRAIN_DATA = [
    ("Astrocyte",        121.77, C_RED,    True),
    ("OPC",               55.42, C_AMBER,  False),
    ("Choroid plexus",    45.47, C_AMBER,  False),
    ("Oligodendrocyte",   44.87, C_AMBER,  False),
    ("COP",               40.03, C_BLUE,   False),
    ("Ependymal",         39.53, C_BLUE,   False),
    ("Microglia",         35.37, C_TEAL,   False),
    ("Fibroblast",        25.41, C_TEAL,   False),
    ("Vascular",          21.54, C_GREEN,  False),
    ("Bergmann glia",     15.97, C_GREEN,  True),
]


def rebuild_s16(prs, slide_idx, lang="EN"):
    """Rebuild S16 with correct NAR data."""
    slide = prs.slides[slide_idx]
    fully_clear_slide(slide)

    data = NAR_BRAIN_DATA if lang == "EN" else ZH_BRAIN_DATA
    max_omega = max(d[1] for d in data)
    gradient = max_omega / min(d[1] for d in data)

    # ── Title ──
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.10), Inches(9.0), Inches(0.55))
    if lang == "EN":
        title_text = "Brain Regional CKI: A 7.6-fold omega Gradient"
        sub_text = "Cell class omega values across 108 human brain regions (CIBR dataset)"
    else:
        title_text = "脑区 CKI：7.6 倍 omega 梯度"
        sub_text = "108 个人脑区的细胞类别 omega 值（CIBR 数据集）"

    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = sub_text
    run2.font.size = Pt(10)
    run2.font.color.rgb = RGBColor(*C_MUTED)

    # ── Left column: horizontal bar chart ──
    chart_left = 0.5
    chart_top = 0.85
    label_w = 1.60
    bar_left = chart_left + label_w + 0.05
    bar_max_w = 3.3
    bar_h = 0.33
    bar_gap = 0.07
    val_w = 0.65

    for i, (label, omega, color, highlight) in enumerate(data):
        y = chart_top + i * (bar_h + bar_gap)

        # Label
        lb = slide.shapes.add_textbox(
            Inches(chart_left), Inches(y), Inches(label_w), Inches(bar_h))
        tf = lb.text_frame
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = label
        run.font.size = Pt(9)
        run.font.bold = highlight
        run.font.color.rgb = RGBColor(*BODY_FG)

        # Bar
        bar_width = bar_max_w * (omega / max_omega)
        bar = slide.shapes.add_shape(
            1, Inches(bar_left), Inches(y),
            Inches(bar_width), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*color)
        bar.line.fill.background()
        bar.shadow.inherit = False

        # Value
        vb = slide.shapes.add_textbox(
            Inches(bar_left + bar_width + 0.08), Inches(y),
            Inches(val_w), Inches(bar_h))
        tf = vb.text_frame
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{omega:.1f}"
        run.font.size = Pt(9)
        run.font.bold = highlight
        run.font.color.rgb = RGBColor(*BODY_FG)

    # Gradient summary
    sum_y = chart_top + len(data) * (bar_h + bar_gap) + 0.08
    sum_box = slide.shapes.add_textbox(Inches(chart_left), Inches(sum_y), Inches(6.4), Inches(0.35))
    tf = sum_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{max_omega:.1f} / {min(d[1] for d in data):.1f} = {gradient:.1f}x gradient across cell classes"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*C_RED)

    # ── Right column: interpretation cards ──
    card_left = 5.2
    card_w = 4.3

    # Card 1: What the Gradient Tells Us
    card1_top = 0.85
    card1_h = 1.50
    card1 = slide.shapes.add_shape(
        1, Inches(card_left), Inches(card1_top),
        Inches(card_w), Inches(card1_h))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(*CARD_BG)
    card1.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card1.line.width = Pt(1.5)
    card1.shadow.inherit = False

    # Card 1 title
    c1t = slide.shapes.add_textbox(Inches(card_left + 0.12), Inches(card1_top + 0.08), Inches(card_w - 0.24), Inches(0.30))
    set_text(c1t, "What the Gradient Tells Us" if lang == "EN" else "梯度的含义",
             size=12, bold=True, color=C_BLUE)

    # Card 1 body
    if lang == "EN":
        lines = [
            "Astrocytes show highest omega: most context-dependent identity",
            "Bergmann glia lowest: stable, region-uniform identity",
            "Non-neuronal cells span the full gradient range",
            "Gradient reflects cellular plasticity hierarchy",
        ]
    else:
        lines = [
            "Astrocyte omega 最高：环境依赖性最强",
            "Bergmann glia 最低：跨区域身份稳定",
            "非神经元细胞跨越全梯度范围",
            "梯度反映细胞可塑性层级",
        ]

    c1b = slide.shapes.add_textbox(
        Inches(card_left + 0.12), Inches(card1_top + 0.38),
        Inches(card_w - 0.24), Inches(card1_h - 0.46))
    tf = c1b.text_frame
    tf.word_wrap = True
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = "  " + line
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(*BODY_FG)

    # Card 2: Key Insight
    card2_top = card1_top + card1_h + 0.18
    card2_h = 1.70
    card2 = slide.shapes.add_shape(
        1, Inches(card_left), Inches(card2_top),
        Inches(card_w), Inches(card2_h))
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(*CARD_BG)
    card2.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card2.line.width = Pt(1.5)
    card2.shadow.inherit = False

    # Card 2 title bar
    c2bar = slide.shapes.add_shape(
        1, Inches(card_left), Inches(card2_top),
        Inches(card_w), Inches(0.35))
    c2bar.fill.solid()
    c2bar.fill.fore_color.rgb = RGBColor(*C_TEAL)
    c2bar.line.fill.background()
    c2bar.shadow.inherit = False
    tf = c2bar.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Insight" if lang == "EN" else "核心洞见"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*C_WHITE)

    # Card 2 body
    if lang == "EN":
        insight_text = (
            "Glial cells are not 'generic support cells' — "
            "they show the highest regional specialisation.\n\n"
            "Astrocytes adapt their transcriptomic identity "
            "to local brain microenvironments, far more than neurons."
        )
    else:
        insight_text = (
            "胶质细胞不是「通用支持细胞」—— "
            "它们表现出最高的区域特化。\n\n"
            "Astrocyte 会根据局部脑微环境调整 "
            "转录组身份，程度远超神经元。"
        )

    c2b = slide.shapes.add_textbox(
        Inches(card_left + 0.12), Inches(card2_top + 0.42),
        Inches(card_w - 0.24), Inches(card2_h - 0.50))
    tf = c2b.text_frame
    tf.word_wrap = True
    for j, para in enumerate(insight_text.split("\n\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(4) if j > 0 else Pt(0)
        run = p.add_run()
        run.text = para.strip()
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(*BODY_FG)

    # ── Bottom summary bar ──
    bar = slide.shapes.add_shape(
        1, Inches(0.5), Inches(5.10), Inches(9.0), Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*C_DARK_BG)
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    if lang == "EN":
        run.text = f"omega range: {min(d[1] for d in data):.1f} - {max_omega:.1f} = {gradient:.1f}x span  |  CKI quantifies this gradient from scRNA-seq alone"
    else:
        run.text = f"omega 范围：{min(d[1] for d in data):.1f} - {max_omega:.1f} = {gradient:.1f}x 跨度  |  CKI 仅从 scRNA-seq 数据即可量化此梯度"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*C_WHITE)

    # Page number
    pn = slide.shapes.add_textbox(Inches(9.2), Inches(5.25), Inches(0.5), Inches(0.25))
    set_text(pn, str(slide_idx + 1), size=8, color=C_MUTED)

    print(f"  S{slide_idx+1}: rebuilt with {len(data)} cell classes, {gradient:.1f}x gradient")


# ═══════════════════════════════════════════════════════════════
# S2: Cell States SVG — fix fill format + redesign
# ═══════════════════════════════════════════════════════════════

def generate_cell_state_svg(lang="EN"):
    """Generate SVG showing same cell type -> different states."""
    if lang == "EN":
        title = "Same Cell Type, Different Functional States"
        subtitle = "CKI captures state-specific changes within cell types"
        root_label = "Macrophage (same type)"
        root_detail = "CD45+, F4/80+"
        states = [
            ("M0 (Resting)", "Housekeeping-dominant", "#0D9488"),
            ("M1 (Pro-inflam)", "High immune-response", "#DC2626"),
            ("M2 (Anti-inflam)", "Tissue-repair program", "#F59E0B"),
            ("Exhausted", "Dysfunctional state", "#7C3AED"),
        ]
        bottom = "CKI's HK set (detection_rate + CV) automatically covers all states within a type"
    else:
        title = "同一细胞类型，不同功能状态"
        subtitle = "CKI 捕获细胞类型内的状态特异性变化"
        root_label = "巨噬细胞（同一类型）"
        root_detail = "CD45+, F4/80+"
        states = [
            ("M0（静息）", "持家基因主导 omega", "#0D9488"),
            ("M1（促炎）", "高免疫应答 omega", "#DC2626"),
            ("M2（抗炎）", "组织修复 omega", "#F59E0B"),
            ("耗竭态", "功能失调 omega", "#7C3AED"),
        ]
        bottom = "CKI 的 HK 基因集（detection_rate + CV）自动覆盖类型内所有状态"

    svg_parts = [
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1800 560">',
        '<rect width="1800" height="560" fill="#F8FAFC" rx="8"/>',
        # Title
        f'<text x="900" y="40" text-anchor="middle" font-size="24" font-weight="700" fill="#1E293B" font-family="Arial,Helvetica,sans-serif">{title}</text>',
        f'<text x="900" y="66" text-anchor="middle" font-size="13" fill="#64748B" font-family="Arial,Helvetica,sans-serif">{subtitle}</text>',
        # Arrow marker
        '<defs><marker id="arr" markerWidth="10" markerHeight="8" refX="9" refY="4" orient="auto"><path d="M0,0 L10,4 L0,8" fill="#94A3B8"/></marker></defs>',
        # Root node
        '<rect x="60" y="180" width="240" height="100" rx="14" fill="#F1F5F9" stroke="#0D9488" stroke-width="3"/>',
        f'<text x="180" y="215" text-anchor="middle" font-size="17" font-weight="700" fill="#1E293B" font-family="Arial,Helvetica,sans-serif">{root_label}</text>',
        f'<text x="180" y="242" text-anchor="middle" font-size="12" fill="#64748B" font-family="Arial,Helvetica,sans-serif">{root_detail}</text>',
    ]

    # State boxes
    state_positions = [
        (440, 100),   # M0
        (440, 220),   # M1
        (440, 340),   # M2
        (440, 460),   # Exhausted
    ]

    # Right side: omega concept
    omega_concepts = [
        ("omega_low", "omega ~ low", "#0D9488"),
        ("omega_high", "omega ~ high", "#DC2626"),
        ("omega_mid", "omega ~ medium", "#F59E0B"),
        ("omega_var", "omega ~ variable", "#7C3AED"),
    ]

    for i, ((label, desc, color), (sx, sy)) in enumerate(zip(states, state_positions)):
        # Arrow from root to state box
        arrow_y = sy + 35
        svg_parts.append(f'<line x1="300" y1="{arrow_y}" x2="430" y2="{arrow_y}" stroke="#94A3B8" stroke-width="2.5" marker-end="url(#arr)"/>')

        # State box
        svg_parts.append(f'<rect x="{sx}" y="{sy}" width="300" height="70" rx="10" fill="#FFFFFF" stroke="{color}" stroke-width="2.5"/>')
        svg_parts.append(f'<text x="{sx+150}" y="{sy+30}" text-anchor="middle" font-size="15" font-weight="700" fill="{color}" font-family="Arial,Helvetica,sans-serif">{label}</text>')
        svg_parts.append(f'<text x="{sx+150}" y="{sy+52}" text-anchor="middle" font-size="11" fill="#64748B" font-family="Arial,Helvetica,sans-serif">{desc}</text>')

        # Arrow from state box to omega concept
        omega_label, omega_desc, omega_color = omega_concepts[i]
        svg_parts.append(f'<line x1="740" y1="{arrow_y}" x2="830" y2="{arrow_y}" stroke="#94A3B8" stroke-width="2" marker-end="url(#arr)"/>')

        # Omega concept box
        svg_parts.append(f'<rect x="840" y="{sy+5}" width="200" height="60" rx="8" fill="#F1F5F9" stroke="{omega_color}" stroke-width="1.5"/>')
        svg_parts.append(f'<text x="940" y="{sy+32}" text-anchor="middle" font-size="13" font-weight="700" fill="{omega_color}" font-family="Arial,Helvetica,sans-serif">{omega_label}</text>')
        svg_parts.append(f'<text x="940" y="{sy+50}" text-anchor="middle" font-size="10" fill="#64748B" font-family="Arial,Helvetica,sans-serif">{omega_desc}</text>')

    # Right side summary box
    svg_parts.append('<rect x="1120" y="140" width="620" height="320" rx="12" fill="#FFFFFF" stroke="#E2E8F0" stroke-width="2"/>')
    svg_parts.append('<text x="1430" y="175" text-anchor="middle" font-size="16" font-weight="700" fill="#1E293B" font-family="Arial,Helvetica,sans-serif">Why This Matters for CKI</text>')

    if lang == "EN":
        right_lines = [
            "Cell types are NOT homogeneous — each contains",
            "multiple functional states with distinct gene",
            "expression programs.",
            "",
            "CKI's housekeeping gene set must be universal",
            "across ALL states within a type, because it",
            "defines the 'baseline' from which functional",
            "divergence is measured.",
            "",
            "State-specific genes drive KF (functional",
            "divergence), while universal HK genes anchor",
            "KN (conservation baseline).",
        ]
    else:
        right_lines = [
            "细胞类型不是均质的——每种类型包含",
            "多个功能状态，各有独特的基因",
            "表达程序。",
            "",
            "CKI 的持家基因集必须覆盖类型内",
            "所有状态，因为它定义了功能分化",
            "的「基线」。",
            "",
            "状态特异性基因驱动 KF（功能分化），",
            "而通用持家基因锚定 KN（保守基线）。",
        ]

    for j, line in enumerate(right_lines):
        y = 200 + j * 22
        svg_parts.append(f'<text x="1150" y="{y}" font-size="11" fill="#1E293B" font-family="Arial,Helvetica,sans-serif">{line}</text>')

    # Bottom bar
    svg_parts.append('<rect x="20" y="500" width="1760" height="44" rx="10" fill="#0D9488"/>')
    svg_parts.append(f'<text x="900" y="528" text-anchor="middle" font-size="14" font-weight="700" fill="#FFFFFF" font-family="Arial,Helvetica,sans-serif">{bottom}</text>')

    svg_parts.append('</svg>')
    return '\n'.join(svg_parts)


def render_svg_to_png(svg_path, png_path, width=1800):
    """Render SVG to PNG using cairosvg or chrome headless."""
    try:
        import cairosvg
        cairosvg.svg2png(url=str(svg_path), write_to=str(png_path),
                         output_width=width, output_height=int(width * 560 / 1800))
        print(f"  PNG: {png_path.name} ({png_path.stat().st_size // 1024}KB)")
        return True
    except ImportError:
        pass

    # Fallback: Chrome headless
    import subprocess, tempfile
    html = f'<html><body style="margin:0;padding:0"><img src="file:///{svg_path}" width="{width}"></body></html>'
    html_path = png_path.parent / "_tmp_s2_render.html"
    html_path.write_text(html, encoding='utf-8')

    try:
        result = subprocess.run([
            'C:/Program Files/Google/Chrome/Application/chrome.exe',
            '--headless', '--disable-gpu', '--screenshot',
            f'--window-size={width},560',
            '--default-background-color=0',
            f'--screenshot={png_path}',
            str(html_path)
        ], capture_output=True, timeout=30)
        if png_path.exists() and png_path.stat().st_size > 5000:
            print(f"  PNG: {png_path.name} ({png_path.stat().st_size // 1024}KB)")
            return True
    except Exception as e:
        print(f"  Chrome fallback failed: {e}")

    return False


def rebuild_s2(prs, slide_idx, lang="EN"):
    """Clear S2 and rebuild with original content + fixed cell-state diagram."""
    slide = prs.slides[slide_idx]
    fully_clear_slide(slide)

    # ── Title ──
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.10), Inches(9.0), Inches(0.60))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "What Defines Cell Type?" if lang == "EN" else "什么定义了细胞的\"身份\"？"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)

    # ── Left side: cell type circles ──
    cell_types_en = [
        ("Neuron", 1.1, 1.1, 0.65, 0.65),
        ("T cell", 2.6, 1.5, 0.65, 0.65),
        ("B cell", 1.9, 2.4, 0.65, 0.65),
        ("Macrophage", 0.6, 2.8, 0.75, 0.75),
        ("Fibroblast", 2.9, 3.0, 0.65, 0.65),
    ]
    cell_types_zh = [
        ("神经元", 1.1, 1.1, 0.65, 0.65),
        ("T细胞", 2.6, 1.5, 0.65, 0.65),
        ("B细胞", 1.9, 2.4, 0.65, 0.65),
        ("巨噬细胞", 0.6, 2.8, 0.75, 0.75),
        ("成纤维细胞", 2.9, 3.0, 0.65, 0.65),
    ]
    cell_types = cell_types_en if lang == "EN" else cell_types_zh

    colors = [C_BLUE, C_TEAL, C_AMBER, C_RED, C_PURPLE]
    for i, (name, x, y, w, h) in enumerate(cell_types):
        circle = slide.shapes.add_shape(
            9, Inches(x), Inches(y), Inches(w), Inches(h))  # oval
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*colors[i])
        circle.line.fill.background()
        circle.shadow.inherit = False

        tf = circle.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*C_WHITE)

    # Caption under left
    cap = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(3.5), Inches(0.5))
    tf = cap.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Same tissue,\ndifferent identities" if lang == "EN" else "同一组织，\n不同的身份"
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = RGBColor(*C_MUTED)

    # ── Right side: text sections ──
    # The Single-Cell Revolution
    rev_box = slide.shapes.add_textbox(Inches(5.0), Inches(0.85), Inches(4.5), Inches(1.2))
    tf = rev_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "The Single-Cell Revolution" if lang == "EN" else "单细胞组学革命"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*C_BLUE)

    bullets_rev = [
        "scRNA-seq reveals vast cellular heterogeneity",
        "Millions of cells across hundreds of cell types",
        "Each cell type has a unique gene expression program",
    ] if lang == "EN" else [
        "scRNA-seq 揭示巨大的细胞异质性",
        "数百万细胞，跨越数百种细胞类型",
        "每种细胞类型都有独特的基因表达程序",
    ]

    for j, b in enumerate(bullets_rev):
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = b
        run2.font.size = Pt(10)
        run2.font.color.rgb = RGBColor(*BODY_FG)

    # The Knowledge Gap
    gap_box = slide.shapes.add_textbox(Inches(5.0), Inches(2.3), Inches(4.5), Inches(1.5))
    tf = gap_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "The Knowledge Gap" if lang == "EN" else "知识空白"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*C_RED)

    bullets_gap = [
        "How do these identities emerge?",
        "Which cell types are truly 'differentiated'?",
        "How much of the difference is functional vs noise?",
        "Can we quantify 'cell-ness' on a continuous scale?",
        "How do cells migrate between organs?",
    ] if lang == "EN" else [
        "这些身份是如何产生的？",
        "哪些细胞类型是真正\"分化\"的？",
        "差异中多少是功能性的，多少是噪音？",
        "能否在连续尺度上量化\"细胞性\"？",
        "细胞如何在器官间迁移？",
    ]

    for j, b in enumerate(bullets_gap):
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = b
        run2.font.size = Pt(10)
        run2.font.color.rgb = RGBColor(*BODY_FG)

    # ── Bottom: Cell-state diagram as image ──
    svg_path = BASE / f"_s2_cell_states_{lang.lower()}_v2.svg"
    png_path = BASE / f"_s2_cell_states_{lang.lower()}_v2.png"

    svg_content = generate_cell_state_svg(lang)
    svg_path.write_text(svg_content, encoding='utf-8')
    print(f"  SVG: {svg_path.name}")

    if render_svg_to_png(svg_path, png_path):
        # Add image to slide at bottom
        pic = slide.shapes.add_picture(
            str(png_path),
            Inches(0.3), Inches(4.2),
            Inches(9.4), Inches(1.3))  # Scale to fit bottom strip
        print(f"  Image added: 9.4x1.3 at y=4.2")
    else:
        # Fallback: text-only cell state section
        cs_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.2), Inches(9.4), Inches(1.2))
        tf = cs_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        if lang == "EN":
            run.text = "Same cell type, different states: Macrophage → M0 (resting) / M1 (pro-inflammatory) / M2 (anti-inflammatory) / Exhausted — CKI captures state-specific changes within cell types"
        else:
            run.text = "同一细胞类型，不同状态：巨噬细胞 → M0（静息）/ M1（促炎）/ M2（抗炎）/ 耗竭态 — CKI 捕获细胞类型内的状态特异性变化"
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = RGBColor(*C_TEAL)

    # Page number
    pn = slide.shapes.add_textbox(Inches(9.2), Inches(5.25), Inches(0.5), Inches(0.25))
    set_text(pn, str(slide_idx + 1), size=8, color=C_MUTED)

    print(f"  S{slide_idx+1}: rebuilt with cell-type circles + knowledge gap + cell-state diagram")


# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════

def main():
    en_path = BASE / "CKI_Lecture_2026_v3.pptx"
    zh_path = BASE / "CKI_Lecture_2026_v3_ZH.pptx"

    # ── EN version ──
    print("[EN] Loading...")
    prs_en = Presentation(str(en_path))
    print(f"  {len(prs_en.slides)} slides")

    # S16 (index 15): Brain Regional CKI
    rebuild_s16(prs_en, 15, lang="EN")

    # S2 (index 1): What Defines Cell Type
    rebuild_s2(prs_en, 1, lang="EN")

    prs_en.save(str(en_path))
    print(f"[EN] Saved: {en_path.name}")

    # ── ZH version ──
    print("\n[ZH] Loading...")
    prs_zh = Presentation(str(zh_path))
    print(f"  {len(prs_zh.slides)} slides")

    rebuild_s16(prs_zh, 15, lang="ZH")
    rebuild_s2(prs_zh, 1, lang="ZH")

    prs_zh.save(str(zh_path))
    print(f"[ZH] Saved: {zh_path.name}")

    # ── Verify ──
    print("\n=== Verification ===")
    for tag, path in [("EN", en_path), ("ZH", zh_path)]:
        prs = Presentation(str(path))
        for idx in [1, 15]:
            slide = prs.slides[idx]
            pics = sum(1 for s in slide.shapes if s.shape_type == 13)
            texts = sum(1 for s in slide.shapes if hasattr(s, 'text_frame') and s.text.strip())
            print(f"  {tag} S{idx+1}: {texts} text shapes, {pics} pics")

    print("\nDone.")


if __name__ == "__main__":
    main()
