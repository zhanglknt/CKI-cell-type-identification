"""
63_split_s2.py - S2拆分为S2(问题引入) + S2b(细胞状态连续体)
v4 PPTX EN+ZH同步修复，页码全部+1
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pathlib import Path
import copy

BASE = Path(r"C:\Users\KnightZ\Desktop\细胞受选择\results\figures_final")
EN_PPT = BASE / "CKI_Lecture_2026_v4.pptx"
ZH_PPT = BASE / "CKI_Lecture_2026_v4_ZH.pptx"

NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(7.5)

# ============================================================
# Helper: reorder sldIdLst to insert new slide at position
# ============================================================
def insert_slide_at(prs, new_slide_idx, target_idx):
    """Move the newly added slide (at new_slide_idx) to target_idx in sldIdLst."""
    sldIdLst = prs.part._element.find(".//p:sldIdLst", NS)
    sldIds = list(sldIdLst.findall("p:sldId", NS))
    if new_slide_idx >= len(sldIds):
        return  # nothing to move
    moved = sldIds[new_slide_idx]
    sldIdLst.remove(moved)
    sldIdLst.insert(target_idx, moved)


def fix_all_page_numbers(prs):
    """Update all page numbers: 2→2, 2b→3, 3→4, ... 21→22"""
    expected = {
        # 0-indexed slide: page number (None = no number)
        0: None,   # S1 title
        1: 2,      # S2
        2: None,   # S2b (new, no page number)
        3: 4,      # old S3 → new S4
        4: 5,      # S5
        5: 6,      # S6
        6: 7,      # S7
        7: 8,      # S8
        8: 9,      # S9
        9: 10,     # S10
        10: 11,    # S11
        11: 12,    # S12
        12: 13,    # S13
        13: 14,    # S14
        14: 15,    # S15
        15: 16,    # S16
        16: 17,    # S17
        17: 18,    # S18
        18: 19,    # S19
        19: 20,    # S20
        20: 21,    # S21
        21: 22,    # S22
        22: None,  # Thank you
    }

    for idx, slide in enumerate(prs.slides):
        expected_num = expected.get(idx)
        for s in slide.shapes:
            if hasattr(s, "text_frame"):
                txt = s.text_frame.text.strip()
                if txt.isdigit() and s.left/914400 > 8.5 and s.top/914400 > 5.0:
                    for p in s.text_frame.paragraphs:
                        for r in p.runs:
                            r.text = str(expected_num) if expected_num else ""
                    if expected_num:
                        pass  # updated
                    else:
                        # Remove page number shape entirely
                        s._element.getparent().remove(s._element)
                    break


def clean_slide_text(slide):
    """Remove all content shapes from a slide, preserving required group props."""
    spTree = slide._element.find(".//p:spTree", NS)
    if spTree is None:
        return
    # Keep required group elements: nvGrpSpPr, grpSpPr
    # Also keep extLst if present
    protected = {'nvGrpSpPr', 'grpSpPr', 'extLst'}
    children = list(spTree)
    for child in children:
        tag = etree.QName(child).localname if hasattr(child, 'tag') else ''
        if tag not in protected:
            spTree.remove(child)


# ============================================================
# Process each PPTX
# ============================================================
for tag, pptx_path, cell_state_png in [
    ("EN", EN_PPT, "_s2_cell_states_en_v2.png"),
    ("ZH", ZH_PPT, "_s2_cell_states_zh_v2.png"),
]:
    print(f"\n{'='*60}")
    print(f"Processing {tag}")
    print(f"{'='*60}")

    prs = Presentation(str(pptx_path))
    print(f"  Original slides: {len(prs.slides)}")

    # ----------------------------------------------------------
    # Step 1: Add S2b slide (clone S2 layout)
    # ----------------------------------------------------------
    s2_layout = prs.slides[1].slide_layout
    s2b_slide = prs.slides.add_slide(s2_layout)
    s2b_idx = len(prs.slides) - 1  # index of newly added slide
    print(f"  Added S2b at end (idx={s2b_idx})")

    # Clear S2b content
    clean_slide_text(s2b_slide)

    # ----------------------------------------------------------
    # Step 2: Set up S2b content - Cell State Continuum
    # ----------------------------------------------------------
    # Title
    if tag == "EN":
        title_text = "Same Cell Type, Different Functional States"
        subtitle_text = "CKI captures state-specific changes within cell types"
    else:
        title_text = "同一细胞类型，不同功能状态"
        subtitle_text = "CKI 捕获细胞类型内的状态特异性变化"

    tb = s2b_slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9.0), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    # Subtitle
    tb2 = s2b_slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(9.0), Inches(0.35))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle_text
    run2.font.size = Pt(13)
    run2.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
    run2.font.italic = True

    # Cell state diagram (large, centered)
    png_path = BASE / cell_state_png
    if png_path.exists():
        # Image is 1800x560, we want to use ~9" width
        img_w = Inches(9.2)
        img_h = Inches(2.86)  # 560/1800 * 9.2
        s2b_slide.shapes.add_picture(
            str(png_path),
            Inches(0.4), Inches(1.3),
            img_w, img_h
        )
        print(f"  S2b: added {cell_state_png}")
    else:
        print(f"  S2b: WARNING - {cell_state_png} not found")

    # Key point callout box
    if tag == "EN":
        key_text = "Key Insight: CKI's HK gene set (detection_rate + CV) automatically covers all states within a cell type, providing a stable baseline for quantifying functional divergence."
    else:
        key_text = "关键洞察：CKI 的 HK 基因集（检出率 + CV）自动覆盖细胞类型内的所有状态，为量化功能分化为提供稳定基线。"

    tb3 = s2b_slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9.0), Inches(0.8))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = key_text
    run3.font.size = Pt(13)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x0D, 0x94, 0x88)

    print(f"  S2b: content set up")

    # ----------------------------------------------------------
    # Step 3: Reorder slides - move S2b to position 2 (after S2)
    # ----------------------------------------------------------
    insert_slide_at(prs, s2b_idx, 2)
    print(f"  S2b: moved to position 2 (after S2)")

    # ----------------------------------------------------------
    # Step 4: Clean up S2 (remove cell state image, adjust layout)
    # ----------------------------------------------------------
    s2_slide = prs.slides[1]  # S2 is now at index 1 (unchanged)

    # Remove the cell state image (shape[10] from earlier analysis)
    # Find and remove all picture shapes
    removed_pics = 0
    for s in list(s2_slide.shapes):
        if hasattr(s, "image"):
            s._element.getparent().remove(s._element)
            removed_pics += 1
    print(f"  S2: removed {removed_pics} image(s)")

    # Adjust text boxes - expand right-side boxes since bottom image is gone
    for s in s2_slide.shapes:
        if hasattr(s, "text_frame"):
            txt = s.text_frame.text.strip()
            # Expand the Knowledge Gap box
            if "Knowledge Gap" in txt or "知识空白" in txt:
                s.height = Inches(3.5)  # was 1.5, now much taller
                s.top = Inches(2.3)
            # Move the single-cell revolution box if needed
            if "Single-Cell Revolution" in txt or "单细胞革命" in txt:
                s.height = Inches(1.5)
                s.top = Inches(0.8)

    # Rearrange cell type circles - better spacing
    # Current circles: Neuron(1.1,1.1), Tcell(2.6,1.5), Bcell(1.9,2.4), Macrophage(0.6,2.8), Fibroblast(2.9,3.0)
    # New arrangement: 2 rows, evenly spaced
    circle_specs = []  # collect circle shapes
    for s in s2_slide.shapes:
        if hasattr(s, "text_frame"):
            txt = s.text_frame.text.strip()
            if txt in ("Neuron", "T cell", "B cell", "Macrophage", "Fibroblast",
                       "神经元", "T 细胞", "B 细胞", "巨噬细胞", "成纤维细胞"):
                circle_specs.append((txt, s))

    if len(circle_specs) == 5 and tag == "EN":
        # Arrange in 2 rows: 3 top + 2 bottom
        positions = [
            (1.0, 1.1),   # Neuron
            (2.4, 1.1),   # T cell
            (3.8, 1.1),   # B cell
            (1.7, 2.6),   # Macrophage
            (3.1, 2.6),   # Fibroblast
        ]
        # Sort circles by name to match positions
        order = ["Neuron", "T cell", "B cell", "Macrophage", "Fibroblast"]
        for name, shape in circle_specs:
            idx = order.index(name)
            shape.left = Inches(positions[idx][0])
            shape.top = Inches(positions[idx][1])
        print(f"  S2: rearranged cell type circles")

        # Move "Same tissue, different identities" text
        for s in s2_slide.shapes:
            if hasattr(s, "text_frame") and "Same tissue" in s.text_frame.text:
                s.top = Inches(3.8)
                s.left = Inches(0.8)

    elif len(circle_specs) == 5 and tag == "ZH":
        # Similar arrangement for Chinese version
        positions = [
            (1.0, 1.1),   # 神经元
            (2.4, 1.1),   # T 细胞
            (3.8, 1.1),   # B 细胞
            (1.7, 2.6),   # 巨噬细胞
            (3.1, 2.6),   # 成纤维细胞
        ]
        order = ["神经元", "T 细胞", "B 细胞", "巨噬细胞", "成纤维细胞"]
        for name, shape in circle_specs:
            if name in order:
                idx = order.index(name)
                shape.left = Inches(positions[idx][0])
                shape.top = Inches(positions[idx][1])
        print(f"  S2: rearranged cell type circles")

        # Move Chinese subtitle
        for s in s2_slide.shapes:
            if hasattr(s, "text_frame") and ("同一组织" in s.text_frame.text or "不同身份" in s.text_frame.text):
                s.top = Inches(3.8)
                s.left = Inches(0.8)

    # ----------------------------------------------------------
    # Step 5: Fix all page numbers
    # ----------------------------------------------------------
    fix_all_page_numbers(prs)
    print(f"  Page numbers: fixed (S2b=no#, S3→S4, ... S21→S22)")

    # ----------------------------------------------------------
    # Step 6: Save
    # ----------------------------------------------------------
    prs.save(str(pptx_path))
    new_count = len(prs.slides)
    print(f"  {tag}: saved ({pptx_path.stat().st_size//1024} KB, {new_count} slides)")

# ============================================================
# Final verification
# ============================================================
print(f"\n{'='*60}")
print("VERIFICATION")
print(f"{'='*60}")

for tag, pptx_path in [("EN", EN_PPT), ("ZH", ZH_PPT)]:
    prs = Presentation(str(pptx_path))
    print(f"\n  {tag}: {len(prs.slides)} slides")
    
    # Check S2 has no images
    s2 = prs.slides[1]
    s2_pics = sum(1 for s in s2.shapes if hasattr(s, "image"))
    print(f"    S2: {s2_pics} images (expect 0)")
    
    # Check S2b has image
    s2b = prs.slides[2]
    s2b_pics = sum(1 for s in s2b.shapes if hasattr(s, "image"))
    print(f"    S2b: {s2b_pics} images (expect 1)")
    
    # Check page numbers
    page_issues = []
    expected = {0: None, 1: 2, 2: None, 3: 4, 4: 5, 5: 6, 6: 7, 7: 8,
                8: 9, 9: 10, 10: 11, 11: 12, 12: 13, 13: 14, 14: 15,
                15: 16, 16: 17, 17: 18, 18: 19, 19: 20, 20: 21,
                21: 22, 22: None}
    for idx, slide in enumerate(prs.slides):
        exp = expected.get(idx)
        found = None
        for s in slide.shapes:
            if hasattr(s, "text_frame"):
                txt = s.text_frame.text.strip()
                if txt.isdigit() and s.left/914400 > 8.5 and s.top/914400 > 5.0:
                    found = int(txt)
        if exp is not None and found != exp:
            page_issues.append(f"S{idx+1}: page={found}, expect={exp}")
    
    if page_issues:
        for iss in page_issues:
            print(f"    ⚠ {iss}")
    else:
        print(f"    Page numbers: OK")

print("\nDone.")
