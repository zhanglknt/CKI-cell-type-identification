"""
notebooks/53_redraw_s15_17_wide.py

Redesign S15-17 SVGs to match the 8.5" × 2.0" PPT layout (4.25:1 aspect ratio).
Current SVGs (1200×460, ~2.6:1) get stretched and look blurry/harsh.

Strategy:
  - SVG viewBox: 1700 × 400 (= 4.25:1, exact match to PPT)
  - Remove titles (PPT already has slide titles)
  - Redesign content for wide-short horizontal format
  - Render at high-res via Chrome headless (window 2800,700)
  - Update both EN and ZH PPTs with new images
"""
import sys, http.server, threading, time, os, subprocess, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE   = Path(r"C:\Users\KnightZ\Desktop\细胞受选择")
ASSETS = BASE / "results/figures_final"
ASSETS.mkdir(parents=True, exist_ok=True)

PPT_EN = ASSETS / "CKI_Lecture_2026_v3.pptx"
PPT_ZH = ASSETS / "CKI_Lecture_2026_v3_ZH.pptx"

# ── SVG constants ──
VW, VH = 1700, 400  # 4.25:1 exactly
BG = "#F8FAFC"
CARD_BG = "#FFFFFF"
CARD_STROKE = "#E2E8F0"

# Colors
TEAL   = "#0D9488"
BLUE   = "#2563EB"
RED    = "#DC2626"
PURPLE = "#7C3AED"
AMBER  = "#F59E0B"
DARK   = "#0F172A"
MUTED  = "#64748B"
GREEN  = "#166534"
RED_DK = "#991B1B"
WHITE  = "#FFFFFF"

# Font
FONT = 'font-family="Arial,Helvetica,sans-serif"'


def _esc(s: str) -> str:
    """Escape XML special chars in text content."""
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def _t(x, y, text, size=11, fill=DARK, bold=False, anchor="middle", italic=False):
    """Shortcut: SVG <text> element."""
    fs = f' font-style="italic"' if italic else ""
    fw = ' font-weight="700"' if bold else ""
    return f'<text x="{x}" y="{y}" text-anchor="{anchor}" font-size="{size}"{fw}{fs} fill="{fill}" {FONT}>{_esc(text)}</text>'


# ═══════════════════════════════════════════════════════
#  SVG A — Slide 15: Baseline (Intra vs Inter ω)
# ═══════════════════════════════════════════════════════

def make_s15() -> str:
    """Wide horizontal bar comparison: Intra-region vs Inter-region."""
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW} {VH}">']

    # ── Background ──
    out.append(f'<rect width="{VW}" height="{VH}" fill="{BG}" rx="6"/>')

    # ── LEFT PANEL: Intra-region ──
    lx, lw = 20, 800
    out.append(f'<rect x="{lx}" y="55" width="{lw}" height="290" rx="8" fill="{CARD_BG}" stroke="{CARD_STROKE}" stroke-width="1.5"/>')
    out.append(_t(lx+lw//2, 80, "Intra-region ω", 14, DARK, bold=True))
    out.append(_t(lx+lw//2, 100, "Same brain region · Different cell types", 10, MUTED))

    # Bar chart: Astrocytes ω=28, OPCs ω=22, Neurons ω=18
    bar_max_w = 380  # max bar width in px
    bar_h = 30
    bar_x0 = lx + 140
    items_intra = [
        ("Astrocytes", 28, TEAL),
        ("OPCs", 22, BLUE),
        ("Neurons(ex)", 18, PURPLE),
    ]
    for i, (label, omega, color) in enumerate(items_intra):
        y = 128 + i * 48
        bw = int(bar_max_w * omega / 35.0)
        out.append(f'<rect x="{bar_x0}" y="{y}" width="{bw}" height="{bar_h}" rx="4" fill="{color}" opacity="0.85"/>')
        out.append(_t(bar_x0 + bw + 12, y + bar_h//2 + 4, str(omega), 13, color, bold=True, anchor="start"))
        out.append(_t(bar_x0 - 12, y + bar_h//2 + 4, label, 12, DARK, anchor="end"))

    # Insight box
    out.append(f'<rect x="{lx+15}" y="310" width="{lw-30}" height="26" rx="6" fill="#F0FDF4" stroke="#86EFAC" stroke-width="1"/>')
    out.append(_t(lx+lw//2, 328, "Intra-region ω = LOW  —  Same microenvironment → similar expression", 10, GREEN, bold=True))

    # ── CENTER: "<" divider ──
    cx = lx + lw + 30
    out.append(f'<line x1="{cx}" y1="60" x2="{cx}" y2="340" stroke="#CBD5E1" stroke-width="2" stroke-dasharray="6,4"/>')
    out.append(_t(cx, 180, "<", 36, "#94A3B8", bold=True))
    out.append(_t(cx, 205, "BASELINE", 9, "#94A3B8", bold=True))

    # ── RIGHT PANEL: Inter-region ──
    rx = cx + 30
    rw = VW - rx - 20
    out.append(f'<rect x="{rx}" y="55" width="{rw}" height="290" rx="8" fill="{CARD_BG}" stroke="{CARD_STROKE}" stroke-width="1.5"/>')
    out.append(_t(rx+rw//2, 80, "Inter-region ω", 14, DARK, bold=True))
    out.append(_t(rx+rw//2, 100, "Same cell type · Different brain regions", 10, MUTED))

    items_inter = [
        ("OPC · SVZ", 3.2, TEAL),
        ("OPC · Hipp", 18, BLUE),
        ("OPC · Cortex", 22, PURPLE),
    ]
    bar_x1 = rx + 140
    for i, (label, omega, color) in enumerate(items_inter):
        y = 128 + i * 48
        bw = int(bar_max_w * omega / 35.0)
        out.append(f'<rect x="{bar_x1}" y="{y}" width="{bw}" height="{bar_h}" rx="4" fill="{color}" opacity="0.85"/>')
        out.append(_t(bar_x1 + bw + 12, y + bar_h//2 + 4, str(omega), 13, color, bold=True, anchor="start"))
        out.append(_t(bar_x1 - 12, y + bar_h//2 + 4, label, 12, DARK, anchor="end"))

    out.append(f'<rect x="{rx+15}" y="310" width="{rw-30}" height="26" rx="6" fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1"/>')
    out.append(_t(rx+rw//2, 328, "Inter-region ω = HIGH  —  Different microenvironments → distinct expression", 10, RED_DK, bold=True))

    # ── Bottom conclusion bar ──
    out.append(f'<rect x="20" y="358" width="{VW-40}" height="34" rx="8" fill="{DARK}"/>')
    out.append(_t(VW//2, 381, "Baseline Rule: Intra-region ω < Inter-region ω  →  Cells differentiate LOCALLY, not by migrating", 12, WHITE, bold=True))

    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════
#  SVG B — Slide 16: Anomaly (OPC breaks baseline)
# ═══════════════════════════════════════════════════════

def make_s16() -> str:
    """Two-panel bar comparison: Intra-type OPC (tiny diff) vs Inter-type (huge diff)."""
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW} {VH}">']
    out.append(f'<rect width="{VW}" height="{VH}" fill="{BG}" rx="6"/>')

    # ── LEFT: Intra-type ω (OPC SVZ vs OPC Cortex) ──
    lx, lw = 20, 800
    out.append(f'<rect x="{lx}" y="55" width="{lw}" height="290" rx="8" fill="{CARD_BG}" stroke="{CARD_STROKE}" stroke-width="1.5"/>')
    out.append(_t(lx+lw//2, 78, "Intra-type ω (OPC across regions)", 14, DARK, bold=True))
    out.append(_t(lx+lw//2, 98, "OPC in SVZ vs OPC in Cortex  —  expected LOW divergence", 10, MUTED))

    # Bar chart: two bars — SVZ ω=3.2 (tiny), Cortex ω=22 (small)
    bar_max = 220
    bar_w = 90
    bar_base_y = 280
    # SVZ
    svz_h = int(bar_max * 3.2 / 110.0)  # 6 px
    svz_h = max(svz_h, 10)  # min visible
    out.append(f'<rect x="{lx+160}" y="{bar_base_y - svz_h}" width="{bar_w}" height="{svz_h}" rx="4" fill="{TEAL}"/>')
    out.append(_t(lx+160+bar_w//2, bar_base_y-svz_h-10, "ω=3.2", 12, TEAL, bold=True))
    out.append(_t(lx+160+bar_w//2, bar_base_y+18, "SVZ", 11, DARK))

    # Cortex
    ctx_h = int(bar_max * 22.0 / 110.0)  # 44 px
    out.append(f'<rect x="{lx+340}" y="{bar_base_y - ctx_h}" width="{bar_w}" height="{ctx_h}" rx="4" fill="{BLUE}"/>')
    out.append(_t(lx+340+bar_w//2, bar_base_y-ctx_h-10, "ω=22", 12, BLUE, bold=True))
    out.append(_t(lx+340+bar_w//2, bar_base_y+18, "Cortex", 11, DARK))

    # Difference annotation
    out.append(_t(lx+600, 180, "Δω = 18.8", 12, MUTED))
    out.append(_t(lx+600, 198, "(surprisingly SMALL)", 10, MUTED, italic=True))

    # Insight box
    out.append(f'<rect x="{lx+15}" y="308" width="{lw-30}" height="28" rx="6" fill="#F0FDF4" stroke="#86EFAC" stroke-width="1"/>')
    out.append(_t(lx+lw//2, 327, "Intra-type ω = LOW  →  OPCs in SVZ and Cortex are HIGHLY similar → likely shared origin", 10, GREEN, bold=True))

    # ── CENTER divider ──
    cx = lx + lw + 30
    out.append(f'<line x1="{cx}" y1="60" x2="{cx}" y2="340" stroke="#CBD5E1" stroke-width="2" stroke-dasharray="6,4"/>')
    out.append(_t(cx, 170, "≪", 36, "#94A3B8", bold=True))
    out.append(_t(cx, 198, "vs", 10, "#94A3B8"))

    # ── RIGHT: Inter-type ω (Astrocyte vs OPC in Cortex) ──
    rx, rw = cx + 30, VW - cx - 50
    out.append(f'<rect x="{rx}" y="55" width="{rw}" height="290" rx="8" fill="{CARD_BG}" stroke="{CARD_STROKE}" stroke-width="1.5"/>')
    out.append(_t(rx+rw//2, 78, "Inter-type ω (within Cortex)", 14, DARK, bold=True))
    out.append(_t(rx+rw//2, 98, "Astrocytes vs OPCs in SAME region  —  expected HIGH divergence", 10, MUTED))

    # Astrocyte ω=107.5 (very tall)
    ast_h = int(bar_max * 107.5 / 110.0)
    out.append(f'<rect x="{rx+160}" y="{bar_base_y - ast_h}" width="{bar_w}" height="{ast_h}" rx="4" fill="{RED}"/>')
    out.append(_t(rx+160+bar_w//2, bar_base_y-ast_h-10, "ω=107.5", 12, RED, bold=True))
    out.append(_t(rx+160+bar_w//2, bar_base_y+18, "Astrocytes", 11, DARK))

    # OPC ω=22 (small)
    opc_h = int(bar_max * 22.0 / 110.0)
    out.append(f'<rect x="{rx+340}" y="{bar_base_y - opc_h}" width="{bar_w}" height="{opc_h}" rx="4" fill="{BLUE}"/>')
    out.append(_t(rx+340+bar_w//2, bar_base_y-opc_h-10, "ω=22", 12, BLUE, bold=True))
    out.append(_t(rx+340+bar_w//2, bar_base_y+18, "OPCs", 11, DARK))

    out.append(_t(rx+600, 180, "Δω = 85.5", 12, MUTED))
    out.append(_t(rx+600, 198, "(extremely LARGE)", 10, MUTED, italic=True))

    out.append(f'<rect x="{rx+15}" y="308" width="{rw-30}" height="28" rx="6" fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1"/>')
    out.append(_t(rx+rw//2, 327, "Inter-type ω = HIGH  →  radically different in SAME region → independent lineages", 10, RED_DK, bold=True))

    # ── Bottom: CONTRADICTION ──
    out.append(f'<rect x="20" y="358" width="{VW-40}" height="34" rx="8" fill="{PURPLE}"/>')
    out.append(_t(VW//2, 381, "CONTRADICTION: Intra-type ω (22) ≪ Inter-type ω (107.5)  →  OPC similarity across regions EXCEEDS cell-type difference within region", 11, WHITE, bold=True))

    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════
#  SVG C — Slide 17: Brain migration hypothesis
# ═══════════════════════════════════════════════════════

def make_s17() -> str:
    """Wide brain diagram with regions, ω values, and migration arrow."""
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW} {VH}">']
    out.append(f'<rect width="{VW}" height="{VH}" fill="{BG}" rx="6"/>')

    # ── Brain outline (simplified silhouette) ──
    # Draw as a large bean-shaped path
    brain_path = ("M 180 120 "
                  "C 350 60, 700 50, 950 70 "
                  "C 1200 90, 1400 130, 1520 200 "
                  "C 1560 230, 1560 280, 1520 310 "
                  "C 1400 370, 1200 360, 1000 350 "
                  "C 800 340, 500 350, 300 330 "
                  "C 150 315, 120 280, 130 230 "
                  "C 140 170, 150 140, 180 120 Z")
    out.append(f'<path d="{brain_path}" fill="#F1F5F9" stroke="#CBD5E1" stroke-width="3" opacity="0.7"/>')

    # ── Region ellipses ──
    # Cortex (top-right)
    out.append('<ellipse cx="1150" cy="110" rx="160" ry="65" fill="#FEF2F2" stroke="#EF4444" stroke-width="2.5" opacity="0.55"/>')
    out.append(_t(1150, 92, "Cortex", 12, RED_DK, bold=True))
    out.append(_t(1150, 108, "ω=107.5 Astrocytes  |  ω=22 OPCs", 9, RED_DK))

    # SVZ (below cortex)
    out.append('<ellipse cx="870" cy="280" rx="70" ry="45" fill="#F0FDF4" stroke="#16A34A" stroke-width="2.5" opacity="0.6"/>')
    out.append(_t(870, 268, "SVZ", 11, GREEN, bold=True))
    out.append(_t(870, 284, "ω=3.2  OPCs", 9, GREEN))

    # Hippocampus (upper-center)
    out.append('<ellipse cx="780" cy="170" rx="65" ry="38" fill="#FEFCE8" stroke="#F59E0B" stroke-width="2" opacity="0.5"/>')
    out.append(_t(780, 162, "Hipp", 10, "#92400E", bold=True))
    out.append(_t(780, 176, "ω=45", 9, "#92400E"))

    # Striatum
    out.append('<ellipse cx="600" cy="190" rx="55" ry="32" fill="#F5F3FF" stroke="#7C3AED" stroke-width="2" opacity="0.5"/>')
    out.append(_t(600, 183, "Striatum", 10, "#5B21B6", bold=True))
    out.append(_t(600, 197, "ω=38", 9, "#5B21B6"))

    # Cerebellum (bottom-left)
    out.append('<ellipse cx="320" cy="310" rx="100" ry="55" fill="#E0F2FE" stroke="#0EA5E9" stroke-width="2.5" opacity="0.55"/>')
    out.append(_t(320, 300, "Cerebellum", 11, "#0369A1", bold=True))
    out.append(_t(320, 316, "ω=85  Bergmann glia", 9, "#0369A1"))

    # Thalamus
    out.append('<ellipse cx="650" cy="260" rx="40" ry="25" fill="#FFF7ED" stroke="#F97316" stroke-width="1.5" opacity="0.45"/>')
    out.append(_t(650, 263, "Thal ω=40", 8, "#9A3412"))

    # ── MIGRATION ARROW: SVZ → Cortex ──
    out.append('<defs>'
               '<filter id="glow"><feGaussianBlur stdDeviation="3" result="blur"/>'
               '<feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge></filter>'
               '<marker id="arrRed" markerWidth="14" markerHeight="10" refX="12" refY="5" orient="auto">'
               '<path d="M0,0 L14,5 L0,10" fill="#DC2626"/></marker>'
               '</defs>')
    out.append(f'<path d="M 890 250 C 970 190, 1050 150, 1120 120" '
               f'stroke="{RED}" stroke-width="4.5" fill="none" '
               f'stroke-dasharray="12,6" opacity="0.9" filter="url(#glow)" '
               f'marker-end="url(#arrRed)"/>')

    # Migration label pill
    out.append(f'<rect x="930" y="182" width="130" height="30" rx="15" fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(_t(995, 196, "OPC Migration", 12, RED, bold=True))
    out.append(_t(995, 210, "ω 3.2 → 22", 9, MUTED, italic=True))

    # ── Bottom conclusion ──
    out.append(f'<rect x="20" y="358" width="{VW-40}" height="34" rx="8" fill="{TEAL}"/>')
    out.append(_t(VW//2, 381, "CKI identifies migration via ω deviation: OPC SVZ→Cortex similarity exceeds inter-type divergence within Cortex", 11, WHITE))

    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════
#  SVG → PNG via Chrome headless (high-res)
# ═══════════════════════════════════════════════════════

def svg_to_png(svg_path: Path, png_path: Path, view_w=2800, view_h=700) -> Path:
    """Render SVG to high-res PNG using Chrome headless."""
    chrome_candidates = [
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    chrome = None
    for c in chrome_candidates:
        if Path(c).exists():
            chrome = c
            break
    if not chrome:
        print("  ERROR: No Chrome/Edge found")
        return None

    os.chdir(str(svg_path.parent))
    server = http.server.HTTPServer(("127.0.0.1", 18901), http.server.SimpleHTTPRequestHandler)
    t = threading.Thread(target=server.serve_forever, daemon=True)
    t.start()
    time.sleep(0.5)

    url = f"http://127.0.0.1:18901/{svg_path.name}"
    cmd = [
        chrome, "--headless", "--disable-gpu",
        "--force-device-scale-factor=2",  # 2x retina for crisp rendering
        f"--screenshot={png_path.resolve()}",
        f"--window-size={view_w},{view_h}",
        url
    ]
    subprocess.run(cmd, capture_output=True, timeout=30)
    server.shutdown()

    if png_path.exists():
        sz = png_path.stat().st_size // 1024
        print(f"    PNG: {png_path.name} ({sz} KB)")
        return png_path
    else:
        print(f"    ERROR: PNG not created")
        return None


# ═══════════════════════════════════════════════════════
#  PPT helpers
# ═══════════════════════════════════════════════════════

def remove_all_pics(slide):
    spTree = slide.shapes._spTree
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    pics = spTree.findall('.//p:pic', ns)
    for pic in list(pics):
        spTree.remove(pic)
    return len(pics)

def replace_image_on_slide(slide, png_path):
    n = remove_all_pics(slide)
    slide.shapes.add_picture(str(png_path),
        Inches(0.5), Inches(1.05), Inches(8.5), Inches(2.0))
    return n


# ═══════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("Redesign S15-17 SVGs for 8.5\"×2.0\" wide-strip layout")
    print(f"  Aspect ratio: {VW}:{VH} = {VW/VH:.2f}:1")
    print("=" * 60)

    # ── Step 1: Generate SVGs ──
    print("\n[1/4] Generating wide-format SVGs...")

    svg_files = {
        "s15": (ASSETS / "_s15_baseline_wide.svg", make_s15()),
        "s16": (ASSETS / "_s16_anomaly_wide.svg", make_s16()),
        "s17": (ASSETS / "_s17_migration_wide.svg", make_s17()),
    }
    for key, (path, svg) in svg_files.items():
        path.write_text(svg, encoding="utf-8")
        print(f"  {key}: {path.name} ({len(svg)//1024} KB, {VW}×{VH})")

    # ── Step 2: SVG → PNG ──
    print("\n[2/4] Rendering high-res PNGs (Chrome 2800×700, 2x scale)...")

    png_files = {}
    for key, (svg_path, _) in svg_files.items():
        png_path = ASSETS / f"_{key}_wide.png"
        result = svg_to_png(svg_path, png_path)
        if result:
            png_files[key] = result

    if len(png_files) < 3:
        print("  ERROR: Not all PNGs rendered. Aborting.")
        return

    # ── Step 3: Update EN PPT ──
    print("\n[3/4] Updating EN PPT...")
    prs = Presentation(str(PPT_EN))

    slide_map = {"s15": 14, "s16": 15, "s17": 16}
    for key, idx in slide_map.items():
        slide = prs.slides[idx]
        n = replace_image_on_slide(slide, png_files[key])
        print(f"  Slide {idx+1}: removed {n} old pic(s), inserted {png_files[key].name}")

    prs.save(str(PPT_EN))
    print(f"  Saved: {PPT_EN.name}")

    # ── Step 4: Update ZH PPT ──
    print("\n[4/4] Updating ZH PPT...")
    prs_zh = Presentation(str(PPT_ZH))
    for key, idx in slide_map.items():
        slide = prs_zh.slides[idx]
        n = replace_image_on_slide(slide, png_files[key])
        print(f"  Slide {idx+1}: removed {n} old pic(s), inserted {png_files[key].name}")

    prs_zh.save(str(PPT_ZH))
    print(f"  Saved: {PPT_ZH.name}")

    print("\n" + "=" * 60)
    print("ALL DONE.")
    print(f"  New SVGs: viewBox 0 0 {VW} {VH} (4.25:1, matches PPT exactly)")
    print(f"  PNGs rendered at 2800×700 with 2x device scale factor")
    print(f"  Images embedded at 8.5\"×2.0\" in both EN and ZH PPTs")
    print("=" * 60)


if __name__ == "__main__":
    main()
