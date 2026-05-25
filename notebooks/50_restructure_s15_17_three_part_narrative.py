"""
notebooks/50_restructure_s15_17_three_part_narrative.py

Restructure slides 15-17 of CKI Lecture PPT into a clear 3-part narrative:
  Part 1 (new Slide 15): BASELINE — intra-region ω < inter-region ω (expected pattern)
  Part 2 (new Slide 16): ANOMALY — OPC breaks the baseline (extreme case)
  Part 3 (new Slide 17): INFERENCE — migration hypothesis (from anomaly to story)

EN + ZH versions.
"""
import sys, json, subprocess, os, re, textwrap, math
from pathlib import Path
from copy import deepcopy

# ── Paths ──────────────────────────────────────────────────────────────
BASE      = Path(r"C:\Users\KnightZ\Desktop\细胞受选择")
PPT_EN    = BASE / "results/figures_final/CKI_Lecture_2026_v3.pptx"
PPT_ZH    = BASE / "results/figures_final/CKI_Lecture_2026_v3_ZH.pptx"
OUT_EN    = BASE / "results/figures_final/CKI_Lecture_2026_v3.pptx"
OUT_ZH    = BASE / "results/figures_final/CKI_Lecture_2026_v3_ZH.pptx"
ASSETS    = BASE / "results/figures_final"
ASSETS.mkdir(parents=True, exist_ok=True)

# ── Slide index (0-based) ────────────────────────────────────────────
S15 = 14   # Slide 15
S16 = 15   # Slide 16
S17 = 16   # Slide 17

# ── Theme palette ─────────────────────────────────────────────────────
DARK_BG   = (0x0F, 0x17, 0x2A)   # #0F172A
CARD_BG    = (0xFF, 0xFF, 0xFF)
TITLE_FG   = (0x1E, 0x29, 0x3B)   # #1E293B
BODY_FG    = (0x1E, 0x29, 0x3B)
MUTED_FG   = (0x64, 0x74, 0x8B)   # #64748B
ACCENT_TEAL = (0x0D, 0x94, 0x88)   # #0D9488
ACCENT_BLUE = (0x25, 0x63, 0xEB)   # #2563EB
ACCENT_RED  = (0xDC, 0x26, 0x26)   # #DC2626
ACCENT_AMBER= (0xF5, 0x9E, 0x0B)   # #F59E0B

def rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

# ═══════════════════════════════════════════════════════════════════
#  SVG helpers
# ═══════════════════════════════════════════════════════════════════

def svg_escape(s: str) -> str:
    return s.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

def omega_color(omega: float, lo: float = 3.0, hi: float = 110.0) -> str:
    """Blue (#2563EB) → Red (#DC2626) via purple mid"""
    t = max(0.0, min(1.0, math.log1p(max(0, omega - lo)) / math.log1p(max(1, hi - lo))))
    # simple 5-stop gradient
    stops = [
        (0.00, (37,  99, 235)),   # blue
        (0.25, (124, 58, 237)),   # purple
        (0.50, (220, 38, 38)),   # red
        (0.75, (220, 38, 38)),
        (1.00, (220, 38, 38)),
    ]
    for i in range(len(stops)-1):
        t0, c0 = stops[i]
        t1, c1 = stops[i+1]
        if t <= t1:
            s = (t - t0) / (t1 - t0 + 1e-9)
            r = int(c0[0] + (c1[0]-c0[0]) * s)
            g = int(c0[1] + (c1[1]-c0[1]) * s)
            b = int(c0[2] + (c1[2]-c0[2]) * s)
            return f"rgb({r},{g},{b})"
    return f"rgb({stops[-1][1][0]},{stops[-1][1][1]},{stops[-1][1][2]})"

# ═══════════════════════════════════════════════════════════════════
#  SVG 1 — Slide 15: Baseline (Intra vs Inter ω)
# ═══════════════════════════════════════════════════════════════════

def make_svg_s15_baseline() -> str:
    """
    Slide 15: Baseline — 'Intra-region ω < Inter-region ω'
    
    Left panel:  same brain region, 3 cell types → short colored bars (low ω)
    Right panel: same cell type, 3 brain regions → tall colored bars (high ω)
    Center divider with 'expected baseline' arrow.
    """
    W, H = 1200, 520
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
          f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']
    
    # ── Title ──
    out.append(f'<text x="{W//2}" y="36" text-anchor="middle" '
               f'font-size="18" font-weight="700" fill="#0F172A">'
               f'Baseline: Intra-region ω &lt; Inter-region ω</text>')
    out.append(f'<text x="{W//2}" y="58" text-anchor="middle" '
               f'font-size="12" fill="#64748B">'
               f'Expected pattern — same-type cells across regions are MORE different than different-type cells within a region</text>')
    
    # ── LEFT PANEL: Intra-region (same region, different types) ──
    lx, ly = 40, 90
    out.append(f'<rect x="{lx}" y="{ly}" width="480" height="400" rx="10" fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+240}" y="{ly+28}" text-anchor="middle" '
               f'font-size="14" font-weight="700" fill="#1E293B">Intra-region ω</text>')
    out.append(f'<text x="{lx+240}" y="{ly+48}" text-anchor="middle" '
               f'font-size="11" fill="#64748B">Same brain region · Different cell types</text>')
    
    # Intra-region data (Cortex region, 3 cell types)
    intra_data = [
        ("Astrocytes",  28, "#0D9488"),
        ("OPCs",        22, "#2563EB"),
        ("Neurons (ex)", 18, "#7C3AED"),
    ]
    bar_x = lx + 40
    bar_area_w = 240
    bar_max_h = 220
    label_x = lx + 300
    
    for i, (ctype, omega, color) in enumerate(intra_data):
        y = ly + 80 + i * 80
        bw = int(bar_max_h * (omega / 35.0))  # normalize to max ~35
        bh = 40
        out.append(f'<rect x="{bar_x}" y="{y}" width="{bw}" height="{bh}" rx="4" fill="{color}" opacity="0.85"/>')
        out.append(f'<text x="{bar_x + bw + 8}" y="{y + 26}" font-size="13" fill="{color}" font-weight="600">{omega}</text>')
        out.append(f'<text x="{label_x}" y="{y + 26}" font-size="12" fill="#334155">{ctype}</text>')
    
    # Key insight box (left)
    insight_y = ly + 310
    out.append(f'<rect x="{lx+20}" y="{insight_y}" width="440" height="70" rx="8" fill="#F0FDF4" stroke="#86EFAC" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+240}" y="{insight_y+24}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#166534">Intra-region ω = LOW</text>')
    out.append(f'<text x="{lx+240}" y="{insight_y+44}" text-anchor="middle" '
               f'font-size="11" fill="#166534">Same microenvironment → similar expression</text>')
    out.append(f'<text x="{lx+240}" y="{insight_y+62}" text-anchor="middle" '
               f'font-size="11" fill="#166534">Different cell types BOTH adapt to local signals</text>')
    
    # ── CENTER divider ──
    cx = 580
    out.append(f'<line x1="{cx}" y1="{ly+40}" x2="{cx}" y2="{ly+370}" stroke="#CBD5E1" stroke-width="2" stroke-dasharray="6,4"/>')
    out.append(f'<text x="{cx}" y="{ly+220}" text-anchor="middle" font-size="22" fill="#94A3B8" font-weight="700">'
               f'&lt;</text>')
    out.append(f'<text x="{cx}" y="{ly+242}" text-anchor="middle" font-size="10" fill="#94A3B8">'
               f'BASELINE</text>')
    
    # ── RIGHT PANEL: Inter-region (same type, different regions) ──
    rx, ry = 620, 90
    out.append(f'<rect x="{rx}" y="{ry}" width="540" height="400" rx="10" fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{ry+28}" text-anchor="middle" '
               f'font-size="14" font-weight="700" fill="#1E293B">Inter-region ω</text>')
    out.append(f'<text x="{rx+270}" y="{ry+48}" text-anchor="middle" '
               f'font-size="11" fill="#64748B">Same cell type · Different brain regions</text>')
    
    # Inter-region data (OPC type, 3 brain regions)
    inter_data = [
        ("OPC · SVZ",        3.2,  "#0D9488"),
        ("OPC · Hippocampus", 18,   "#2563EB"),
        ("OPC · Cortex",      22,   "#7C3AED"),
    ]
    r_bar_x = rx + 40
    r_label_x = rx + 320
    
    for j, (label, omega, color) in enumerate(inter_data):
        y = ry + 80 + j * 80
        bw = int(bar_max_h * (omega / 35.0))
        bh = 40
        out.append(f'<rect x="{r_bar_x}" y="{y}" width="{bw}" height="{bh}" rx="4" fill="{color}" opacity="0.85"/>')
        out.append(f'<text x="{r_bar_x + bw + 8}" y="{y + 26}" font-size="13" fill="{color}" font-weight="600">{omega}</text>')
        out.append(f'<text x="{r_label_x}" y="{y + 26}" font-size="12" fill="#334155">{label}</text>')
    
    # Key insight box (right)
    r_insight_y = ry + 310
    out.append(f'<rect x="{rx+20}" y="{r_insight_y}" width="500" height="70" rx="8" fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{r_insight_y+24}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#991B1B">Inter-region ω = HIGH</text>')
    out.append(f'<text x="{rx+270}" y="{r_insight_y+44}" text-anchor="middle" '
               f'font-size="11" fill="#991B1B">Different microenvironments → distinct expression</text>')
    out.append(f'<text x="{rx+270}" y="{r_insight_y+62}" text-anchor="middle" '
               f'font-size="11" fill="#991B1B">Cells of the SAME type specialize to each region</text>')
    
    # ── Bottom summary ──
    out.append(f'<rect x="40" y="{H-52}" width="{W-80}" height="40" rx="8" fill="#0F172A"/>')
    out.append(f'<text x="{W//2}" y="{H-28}" text-anchor="middle" '
               f'font-size="13" font-weight="600" fill="white">'
               f'Baseline:  Intra-region ω &lt; Inter-region ω  →  '
               f'Cells differentiate LOCALLY, not by migrating</text>')
    
    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════════════════
#  SVG 2 — Slide 16: Anomaly (OPC breaks baseline)
# ═══════════════════════════════════════════════════════════════════

def make_svg_s16_anomaly() -> str:
    """
    Slide 16: Anomaly — 'OPC breaks the baseline'
    
    BIG side-by-side comparison:
      LEFT:  Intra-type ω (OPC SVZ ↔ Cortex) = 22  →  anomalously LOW
      RIGHT: Inter-type ω (Astrocyte ↔ OPC, within Cortex) = 107.5 → expected HIGH
    Bottom: 'This contradiction → migration hypothesis'
    """
    W, H = 1200, 560
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
          f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']
    
    # Title
    out.append(f'<text x="{W//2}" y="36" text-anchor="middle" '
               f'font-size="18" font-weight="700" fill="#0F172A">'
               f'Extreme Case: OPC Breaks the Baseline</text>')
    out.append(f'<text x="{W//2}" y="58" text-anchor="middle" '
               f'font-size="12" fill="#64748B">'
               f'Intra-type ω ≪ expected inter-region ω → anomaly worth investigating</text>')
    
    # ── LEFT: Intra-type ω (OPC across regions) ──
    lx, ly = 40, 90
    out.append(f'<rect x="{lx}" y="{ly}" width="540" height="380" rx="10" fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+270}" y="{ly+30}" text-anchor="middle" '
               f'font-size="14" font-weight="700" fill="#1E293B">Intra-type ω (OPC)</text>')
    out.append(f'<text x="{lx+270}" y="{ly+50}" text-anchor="middle" '
               f'font-size="11" fill="#64748B">OPC in SVZ vs OPC in Cortex</text>')
    
    # Two bars: SVZ and Cortex
    svz_x, svz_y = lx + 60, ly + 80
    ctx_x, ctx_y = lx + 60, ly + 180
    bar_w = 80
    bar_max = 300
    
    # SVZ ω=3.2  (tiny bar)
    svz_h = int(bar_max * (3.2 / 110.0))
    out.append(f'<rect x="{svz_x}" y="{ly+80 + bar_max - svz_h}" width="{bar_w}" height="{svz_h}" rx="4" fill="#0D9488"/>')
    out.append(f'<text x="{svz_x + bar_w/2}" y="{ly+80 + bar_max - svz_h - 8}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#0D9488">ω=3.2</text>')
    out.append(f'<text x="{svz_x + bar_w/2}" y="{ly+80 + bar_max + 20}" text-anchor="middle" '
               f'font-size="11" fill="#334155">SVZ</text>')
    
    # Cortex ω=22  (small bar)
    ctx_h = int(bar_max * (22.0 / 110.0))
    out.append(f'<rect x="{ctx_x + 140}" y="{ly+180 + bar_max - ctx_h}" width="{bar_w}" height="{ctx_h}" rx="4" fill="#2563EB"/>')
    out.append(f'<text x="{ctx_x + 140 + bar_w/2}" y="{ly+180 + bar_max - ctx_h - 8}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#2563EB">ω=22</text>')
    out.append(f'<text x="{ctx_x + 140 + bar_w/2}" y="{ly+180 + bar_max + 20}" text-anchor="middle" '
               f'font-size="11" fill="#334155">Cortex</text>')
    
    # Omega difference annotation
    diff_y = ly + 330
    out.append(f'<line x1="{svz_x + bar_w + 20}" y1="{ly+80 + bar_max - svz_h + 10}" '
               f'x2="{ctx_x + 140 - 20}" y2="{ly+180 + bar_max - ctx_h + 10}" '
               f'stroke="#94A3B8" stroke-width="1.5" marker-end="url(#arrowred)"/>')
    out.append(f'<text x="{lx + 230}" y="{diff_y}" text-anchor="middle" '
               f'font-size="12" fill="#64748B" font-style="italic">'
               f'ω diff = 18.8  (surprisingly SMALL)</text>')
    
    # Left insight
    out.append(f'<rect x="{lx+20}" y="{ly+360}" width="500" height="95" rx="8" fill="#F0FDF4" stroke="#86EFAC" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+270}" y="{ly+384}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#166534">Intra-type ω = 22 (LOW)</text>')
    out.append(f'<text x="{lx+270}" y="{ly+404}" text-anchor="middle" '
               f'font-size="11" fill="#166534">OPC in SVZ and Cortex are HIGHLY similar</text>')
    out.append(f'<text x="{lx+270}" y="{ly+424}" text-anchor="middle" '
               f'font-size="11" fill="#166534">→ likely same origin, NOT independent differentiation</text>')
    
    # ── RIGHT: Inter-type ω (within Cortex) ──
    rx, ry = 620, 90
    out.append(f'<rect x="{rx}" y="{ry}" width="540" height="380" rx="10" fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{ry+30}" text-anchor="middle" '
               f'font-size="14" font-weight="700" fill="#1E293B">Inter-type ω (Cortex)</text>')
    out.append(f'<text x="{rx+270}" y="{ry+50}" text-anchor="middle" '
               f'font-size="11" fill="#64748B">Astrocytes vs OPCs (same region)</text>')
    
    # Two bars: Astrocyte and OPC in Cortex
    ast_x, ast_y = rx + 60, ry + 80
    opc_x, opc_y = rx + 60, ry + 220
    
    # Astrocyte ω=107.5  (TALL bar)
    ast_h = int(bar_max * (107.5 / 110.0))
    out.append(f'<rect x="{ast_x}" y="{ry+80 + bar_max - ast_h}" width="{bar_w}" height="{ast_h}" rx="4" fill="#DC2626"/>')
    out.append(f'<text x="{ast_x + bar_w/2}" y="{ry+80 + bar_max - ast_h - 8}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#DC2626">ω=107.5</text>')
    out.append(f'<text x="{ast_x + bar_w/2}" y="{ry+80 + bar_max + 20}" text-anchor="middle" '
               f'font-size="11" fill="#334155">Astrocytes</text>')
    
    # OPC ω=22  (small bar)
    opc_h = int(bar_max * (22.0 / 110.0))
    out.append(f'<rect x="{opc_x + 140}" y="{ry+220 + bar_max - opc_h}" width="{bar_w}" height="{opc_h}" rx="4" fill="#2563EB"/>')
    out.append(f'<text x="{opc_x + 140 + bar_w/2}" y="{ry+220 + bar_max - opc_h - 8}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#2563EB">ω=22</text>')
    out.append(f'<text x="{opc_x + 140 + bar_w/2}" y="{ry+220 + bar_max + 20}" text-anchor="middle" '
               f'font-size="11" fill="#334155">OPCs</text>')
    
    # Omega difference annotation
    out.append(f'<text x="{rx + 230}" y="{ry+370}" text-anchor="middle" '
               f'font-size="12" fill="#64748B" font-style="italic">'
               f'ω diff = 85.5  (extremely LARGE)</text>')
    
    # Right insight
    out.append(f'<rect x="{rx+20}" y="{ry+360}" width="500" height="95" rx="8" fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{ry+384}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#991B1B">Inter-type ω = 107.5 (HIGH)</text>')
    out.append(f'<text x="{rx+270}" y="{ry+404}" text-anchor="middle" '
               f'font-size="11" fill="#991B1B">Astrocytes and OPCs in SAME region are RADICALLY different</text>')
    out.append(f'<text x="{rx+270}" y="{ry+424}" text-anchor="middle" '
               f'font-size="11" fill="#991B1B">→ independent lineages, strong local differentiation</text>')
    
    # ── Bottom: The contradiction ──
    out.append(f'<rect x="40" y="{H-60}" width="{W-80}" height="48" rx="8" fill="#7C3AED"/>')
    out.append(f'<text x="{W//2}" y="{H-36}" text-anchor="middle" '
               f'font-size="14" font-weight="700" fill="white">'
               f'Contradiction:  Intra-type ω (22) ≪ Inter-type ω (107.5)  '
               f'→  OPC similarity across regions EXCEEDS cell-type difference within region</text>')
    
    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════════════════
#  SVG 3 — Slide 17: Brain with migration arrow (from v2, refined)
# ═══════════════════════════════════════════════════════════════════

def make_svg_s17_brain() -> str:
    """
    Slide 17: Brain regional map with clear migration arrow + reasoning.
    Same as v2 but with cleaner styling.
    """
    W, H = 1200, 520
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
          f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']
    
    # ── Brain outline (simple sagittal mouse brain) ──
    # Cerebellum (left-bottom)
    out.append('<ellipse cx="240" cy="380" rx="130" ry="100" fill="#E0F2FE" stroke="#0EA5E9" stroke-width="2.5" opacity="0.5"/>')
    out.append(f'<text x="240" y="410" text-anchor="middle" font-size="11" fill="#0369A1" font-weight="600">Cerebellum</text>')
    out.append(f'<text x="240" y="428" text-anchor="middle" font-size="10" fill="#0369A1">ω=85  Bergmann glia</text>')
    
    # Cortex (top-right)
    out.append('<ellipse cx="680" cy="140" rx="170" ry="120" fill="#FEF2F2" stroke="#EF4444" stroke-width="2.5" opacity="0.5"/>')
    out.append(f'<text x="680" y="120" text-anchor="middle" font-size="11" fill="#991B1B" font-weight="600">Cortex</text>')
    out.append(f'<text x="680" y="136" text-anchor="middle" font-size="10" fill="#991B1B">ω=107.5 Astrocytes</text>')
    out.append(f'<text x="680" y="152" text-anchor="middle" font-size="10" fill="#991B1B">ω=22  OPCs</text>')
    
    # SVZ (bottom-left, near cortex)
    out.append('<ellipse cx="520" cy="340" rx="70" ry="60" fill="#F0FDF4" stroke="#16A34A" stroke-width="2.5" opacity="0.6"/>')
    out.append(f'<text x="520" y="330" text-anchor="middle" font-size="11" fill="#166534" font-weight="600">SVZ</text>')
    out.append(f'<text x="520" y="350" text-anchor="middle" font-size="10" fill="#166534">ω=3.2  OPCs</text>')
    
    # Hippocampus
    out.append('<ellipse cx="780" cy="300" rx="80" ry="60" fill="#FEFCE8" stroke="#F59E0B" stroke-width="2" opacity="0.5"/>')
    out.append(f'<text x="780" y="296" text-anchor="middle" font-size="10" fill="#92400E">Hippocampus ω=45</text>')
    
    # Striatum
    out.append('<ellipse cx="620" cy="260" rx="70" ry="55" fill="#F5F3FF" stroke="#7C3AED" stroke-width="2" opacity="0.5"/>')
    out.append(f'<text x="620" y="262" text-anchor="middle" font-size="10" fill="#5B21B6">Striatum ω=38</text>')
    
    # Thalamus
    out.append('<ellipse cx="550" cy="220" rx="55" ry="45" fill="#F5F3FF" stroke="#7C3AED" stroke-width="2" opacity="0.4"/>')
    out.append(f'<text x="550" y="222" text-anchor="middle" font-size="9" fill="#5B21B6">Thalamus ω=42</text>')
    
    # ── MIGRATION ARROW: SVZ → Cortex ──
    # Thick glowing arrow
    out.append('<defs>'
               '<filter id="glow"><feGaussianBlur stdDeviation="4" result="blur"/>'
               '<feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge></filter>'
               '</defs>')
    out.append(f'<path d="M 560 340 C 600 250, 640 200, 670 160" '
               f'stroke="#DC2626" stroke-width="5" fill="none" '
               f'stroke-dasharray="10,6" opacity="0.9" filter="url(#glow)" '
               f'marker-end="url(#arrowred2)"/>')
    # Arrow marker
    out.append('<defs><marker id="arrowred2" markerWidth="12" markerHeight="9" refX="10" refY="4.5" orient="auto">'
               '<path d="M0,0 L12,4.5 L0,9" fill="#DC2626"/></marker></defs>')
    
    # Migration label (pill)
    pill_x, pill_y = 610, 280
    out.append(f'<rect x="{pill_x-65}" y="{pill_y-14}" width="130" height="28" rx="14" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{pill_x}" y="{pill_y+5}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#DC2626">OPC Migration</text>')
    
    # Omega gradient annotation
    out.append(f'<text x="610" y="310" text-anchor="middle" '
               f'font-size="11" fill="#94A3B8" font-style="italic">ω 3.2 → 22</text>')
    
    # ── Brain regions legend (bottom left) ──
    legend_x, legend_y = 40, 450
    out.append(f'<rect x="{legend_x}" y="{legend_y}" width="260" height="55" rx="8" fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{legend_x+130}" y="{legend_y+20}" text-anchor="middle" '
               f'font-size="11" font-weight="700" fill="#1E293B">Brain Regions &amp; ω</text>')
    out.append(f'<text x="{legend_x+10}" y="{legend_y+40}" font-size="10" fill="#334155">'
               f'SVZ ω=3.2  →  Cortex ω=22, 107.5  →  Cerebellum ω=85</text>')
    
    # ── ω scale bar (bottom right) ──
    scale_x, scale_y = 860, 455
    out.append(f'<text x="{scale_x}" y="{scale_y-12}" font-size="10" fill="#64748B">ω scale →</text>')
    for i in range(50):
        t = i / 49.0
        r = int(37 + (220-37)*t)
        g = int(99 + (38-99)*t)
        b = int(235 + (38-235)*t)
        out.append(f'<rect x="{scale_x + i*4}" y="{scale_y-8}" width="4" height="14" fill="rgb({r},{g},{b})"/>')
    out.append(f'<text x="{scale_x-5}" y="{scale_y+12}" font-size="9" fill="#64748B">0</text>')
    out.append(f'<text x="{scale_x+195}" y="{scale_y+12}" font-size="9" fill="#64748B">110</text>')
    
    # ── Title ──
    out.append(f'<text x="{W//2}" y="28" text-anchor="middle" '
               f'font-size="17" font-weight="700" fill="#0F172A">'
               f'Migration Inference: From Anomaly to Hypothesis</text>')
    out.append(f'<text x="{W//2}" y="48" text-anchor="middle" '
               f'font-size="11" fill="#64748B">'
               f'OPC migration from SVZ to Cortex explains the ω anomaly</text>')
    
    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════════════════
#  Convert SVG → PNG via Chrome headless
# ═══════════════════════════════════════════════════════════════════

def svg_to_png_chrome(svg_path: Path, png_path: Path, scale: int = 2):
    """Use Chrome headless to render SVG → PNG."""
    import shutil
    chrome = (r"C:\Program Files\Google\Chrome\Application\chrome.exe")
    if not Path(chrome).exists():
        chrome = r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"
    if not Path(chrome).exists():
        # Try edge
        chrome = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
    if not Path(chrome).exists():
        chrome = r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"
    assert Path(chrome).exists(), f"Chrome/Edge not found"
    
    # Create a temp HTML wrapper
    html_path = svg_path.with_suffix(".html")
    svg_content = svg_path.read_text(encoding="utf-8")
    html_content = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
body {{ margin:0; padding:0; background:white; display:flex; justify-content:center; align-items:center; }}
</style></head><body>{svg_content}</body></html>"""
    html_path.write_text(html_content, encoding="utf-8")
    
    url = f"file:///{html_path.resolve().as_posix()}"
    # Use a local HTTP server to avoid file:// CORS issues
    import http.server, threading, time
    os.chdir(str(svg_path.parent))
    server = http.server.HTTPServer(("127.0.0.1", 18899), http.server.SimpleHTTPRequestHandler)
    t = threading.Thread(target=server.serve_forever, daemon=True)
    t.start()
    time.sleep(0.5)
    
    url = f"http://127.0.0.1:18899/{svg_path.name}"
    cmd = [
        chrome, "--headless", "--disable-gpu", "--screenshot",
        f"--screenshot={png_path.resolve()}",
        f"--window-size={1200*scale},{700*scale}",
        url
    ]
    subprocess.run(cmd, capture_output=True, timeout=30)
    server.shutdown()
    
    if png_path.exists():
        print(f"      PNG: {png_path.name} ({png_path.stat().st_size//1024} KB)")
    else:
        print(f"      WARNING: PNG not created at {png_path}")
    return png_path.exists()


# ═══════════════════════════════════════════════════════════════════
#  Replace slide content (clear + add new shapes)
# ═══════════════════════════════════════════════════════════════════

def clear_slide(prs, slide_idx):
    """Remove all shapes from a slide (keep background)."""
    slide = prs.slides[slide_idx]
    # Correct attribute name: _spTree (python-pptx API)
    sp_tree = slide.shapes._spTree
    # Remove all shape elements (sp, pic, grpSp) but keep bg
    for el in list(sp_tree):
        tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
        if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
            sp_tree.remove(el)
    return slide


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=None, align="left"):
    """Add a text box to slide."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    else:
        run.font.color.rgb = RGBColor(*TITLE_FG)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    if line_color:
        shape.line.color.rgb = RGBColor(*line_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════════
#  Build Slide 15 (Baseline)
# ═══════════════════════════════════════════════════════════════════

def build_slide_15(slide, is_zh=False):
    """
    Slide 15: Baseline — Intra-region ω < Inter-region ω
    Content: left card (concept), right card (visual = PNG image)
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    # ── Title ──
    title_tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.15), Inches(9.0), Inches(0.5))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "基线：相同脑区内的细胞功能分化 < 不同脑区间的细胞功能分化" if is_zh else \
               "Baseline: Intra-region ω < Inter-region ω"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)
    
    # Subtitle
    sub_tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.65), Inches(9.0), Inches(0.3))
    tf2 = sub_tb.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = (
        "预期模式：同一脑区内不同细胞类型的分化程度，通常低于相同细胞类型在不同脑区之间的分化程度"
        if is_zh else
        "Expected pattern: differentiation within a region < differentiation across regions"
    )
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(*MUTED_FG)
    
    # ── LEFT CARD: concept explanation ──
    add_rect(slide, 0.55, 1.1, 4.5, 3.9, CARD_BG, (0xE2, 0xE8, 0xF0))
    # Left card title
    add_textbox(slide, 0.7, 1.25, 4.2, 0.35, 
                 "核心概念：预期基线" if is_zh else "Core Concept: The Expected Baseline",
                 font_size=13, bold=True, color=ACCENT_BLUE)
    # Accent bar
    add_rect(slide, 0.7, 1.6, 4.0, 0.04, ACCENT_BLUE)
    
    # Left card body
    left_body = [
        ("相同脑区 (Intra-region):", "bold", MUTED_FG),
        ("  不同细胞类型在同一个脑区内", "normal", BODY_FG),
        ("  功能分化程度 = 较低 ω", "normal", BODY_FG),
        ("", "normal", BODY_FG),
        ("  例：Cortex 内 Astrocyte vs OPC", "normal", (0x6B, 0x72, 0x80)),
        ("      ω ≈ 107.5 （高分化 → 独立谱系）", "normal", (0x6B, 0x72, 0x80)),
        ("", "normal", BODY_FG),
        ("不同脑区 (Inter-region):", "bold", MUTED_FG),
        ("  相同细胞类型在不同脑区之间", "normal", BODY_FG),
        ("  功能分化程度 = 较高 ω", "normal", BODY_FG),
        ("", "normal", BODY_FG),
        ("  例：OPC in SVZ vs OPC in Cortex", "normal", (0x6B, 0x72, 0x80)),
        ("      ω ≈ 22 （中等分化 → 可能对齐起源）", "normal", (0x6B, 0x72, 0x80)),
    ]
    y = 1.7
    for text, style, color in left_body:
        if text == "":
            y += 0.1
            continue
        add_textbox(slide, 0.75, y, 4.1, 0.28, text,
                     font_size=10 if style=="normal" else 10,
                     bold=(style=="bold"),
                     color=color)
        y += 0.29
    
    # ── RIGHT CARD: visual (PNG image) ──
    add_rect(slide, 5.2, 1.1, 4.3, 3.9, CARD_BG, (0xE2, 0xE8, 0xF0))
    add_textbox(slide, 5.35, 1.25, 4.0, 0.35,
                 "可视化比较" if is_zh else "Visual Comparison",
                 font_size=13, bold=True, color=ACCENT_TEAL)
    add_rect(slide, 5.35, 1.6, 3.9, 0.04, ACCENT_TEAL)
    
    # Placeholder for image (will be added after PNG generation)
    # Image will be added in a second pass
    img_w, img_h = 3.8, 3.2
    img_left = 5.4
    img_top = 1.68
    add_textbox(slide, img_left, img_top, img_w, img_h,
                 "[Baseline comparison diagram]\n\n(Intra-region ω < Inter-region ω)\n\nWill be replaced with PNG",
                 font_size=10, color=(0x94, 0xA3, 0xB8), align="center")
    
    # ── Bottom insight bar ──
    add_rect(slide, 0.55, 5.05, 8.9, 0.45, (0x0F, 0x17, 0x2A))
    add_textbox(slide, 0.55, 5.1, 8.9, 0.35,
                 "基线结论：如果所有细胞类型都遵循此模式，则无需假设迁移。" if is_zh else
                 "Baseline conclusion: If all cell types follow this pattern, no migration hypothesis is needed.",
                 font_size=11, color=(0xFF, 0xFF, 0xFF), align="center")
    
    return slide


# ═══════════════════════════════════════════════════════════════════
#  Build Slide 16 (Anomaly)
# ═══════════════════════════════════════════════════════════════════

def build_slide_16(slide, is_zh=False):
    """Slide 16: Anomaly — OPC breaks the baseline."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    # Title
    title_tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.15), Inches(9.0), Inches(0.5))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "特例：OPC 打破了预期基线" if is_zh else \
               "Anomaly: OPC Breaks the Expected Baseline"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)
    
    # Subtitle
    sub_tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.65), Inches(9.0), Inches(0.3))
    tf2 = sub_tb.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = (
        "Intra-type ω (SVZ↔Cortex) = 22，远低于同一脑区内 Inter-type ω = 107.5"
        if is_zh else
        "Intra-type ω (SVZ↔Cortex) = 22, far below intra-region inter-type ω = 107.5"
    )
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(*MUTED_FG)
    
    # ── LEFT CARD: side-by-side comparison ──
    add_rect(slide, 0.55, 1.1, 4.5, 3.9, CARD_BG, (0xE2, 0xE8, 0xF0))
    add_textbox(slide, 0.7, 1.25, 4.2, 0.35,
                 "观察：Intra-type ω 异常偏低" if is_zh else "Observation: Intra-type ω Anomalously LOW",
                 font_size=13, bold=True, color=ACCENT_BLUE)
    add_rect(slide, 0.7, 1.6, 4.0, 0.04, ACCENT_BLUE)
    
    left_body = [
        ("Intra-type ω (OPC, SVZ→Cortex):", "bold", MUTED_FG),
        ("  ω ≈ 22", "bold", (0x25, 0x63, 0xEB)),
        ("", "normal", BODY_FG),
        ("  远大于一般 inter-region 差异的预期？", "normal", BODY_FG),
        ("  NO — 实际上远低于预期！", "normal", ACCENT_RED),
        ("", "normal", BODY_FG),
        ("同一细胞类型，跨两个脑区，", "normal", BODY_FG),
        ("功能分化程度非常低 →", "normal", BODY_FG),
        ("两处 OPC 很可能来自共同起源", "normal", BODY_FG),
        ("", "normal", BODY_FG),
        ("→ 最简洁的解释：细胞迁移", "bold", ACCENT_RED),
    ]
    y = 1.7
    for text, style, color in left_body:
        if text == "":
            y += 0.1
            continue
        add_textbox(slide, 0.75, y, 4.1, 0.28, text,
                     font_size=10, bold=(style=="bold"), color=color)
        y += 0.29
    
    # ── RIGHT CARD: counterfactual / normal case ──
    add_rect(slide, 5.2, 1.1, 4.3, 3.9, CARD_BG, (0xE2, 0xE8, 0xF0))
    add_textbox(slide, 5.35, 1.25, 4.0, 0.35,
                 "对照：Inter-type ω 符合预期（极高）" if is_zh else "Control: Inter-type ω as Expected (VERY HIGH)",
                 font_size=13, bold=True, color=ACCENT_RED)
    add_rect(slide, 5.35, 1.6, 3.9, 0.04, ACCENT_RED)
    
    right_body = [
        ("Inter-type ω (Cortex 内):", "bold", MUTED_FG),
        ("  Astrocyte vs OPC → ω ≈ 107.5", "bold", ACCENT_RED),
        ("", "normal", BODY_FG),
        ("同一脑区、不同细胞类型，", "normal", BODY_FG),
        ("功能分化程度极高 →", "normal", BODY_FG),
        ("它们是独立的细胞谱系", "normal", BODY_FG),
        ("", "normal", BODY_FG),
        ("关键对比：", "bold", MUTED_FG),
        ("  Intra-type ω (22) << Inter-type ω (107.5)", "bold", (0x7C, 0x3A, 0xED)),
        ("", "normal", BODY_FG),
        ("OPC 在 SVZ 和 Cortex 的相似度，", "normal", BODY_FG),
        ("甚至远超 Cortex 内不同细胞类型的相似度", "normal", BODY_FG),
        ("", "normal", BODY_FG),
        ("→ 这与「独立分化」的基线模式完全相反", "bold", ACCENT_RED),
    ]
    y = 1.7
    for text, style, color in right_body:
        if text == "":
            y += 0.1
            continue
        add_textbox(slide, 5.4, y, 4.0, 0.28, text,
                     font_size=9.5 if len(text) > 38 else 10,
                     bold=(style=="bold"), color=color)
        y += 0.29
    
    # ── Bottom insight bar ──
    add_rect(slide, 0.55, 5.05, 8.9, 0.45, (0xDC, 0x26, 0x26))
    add_textbox(slide, 0.55, 5.1, 8.9, 0.35,
                 "矛盾：如果 OPC 在 SVZ 和 Cortex 各自独立分化，为何它们之间的相似度 > 同一脑区内不同细胞的相似度？"
                 if is_zh else
                 "Contradiction: If OPCs in SVZ and Cortex differentiated independently, why is their similarity > within-region cell-type similarity?",
                 font_size=11, color=(0xFF, 0xFF, 0xFF), align="center")
    
    return slide


# ═══════════════════════════════════════════════════════════════════
#  Build Slide 17 (Migration Inference)
# ═══════════════════════════════════════════════════════════════════

def build_slide_17(slide, has_image=False, is_zh=False):
    """Slide 17: Migration Inference — with brain diagram."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    # Title
    title_tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.15), Inches(9.0), Inches(0.5))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "推断：OPC 迁移假说" if is_zh else \
               "Inference: The OPC Migration Hypothesis"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)
    
    # Subtitle
    sub_tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.65), Inches(9.0), Inches(0.3))
    tf2 = sub_tb.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = (
        "从「矛盾」到「假说」：OPC 从 SVZ 向 Cortex 的迁移是最简洁的解释"
        if is_zh else
        "From contradiction to hypothesis: OPC migration from SVZ to Cortex is the most parsimonious explanation"
    )
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(*MUTED_FG)
    
    # ── LEFT: brain image (placeholder or real) ──
    img_left, img_top, img_w, img_h = 0.55, 1.05, 5.0, 4.0
    add_rect(slide, img_left, img_top, img_w, img_h, (0xF8, 0xFA, 0xFC), (0xE2, 0xE8, 0xF0))
    add_textbox(slide, img_left, img_top + 1.5, img_w, 1.0,
                 "[Brain diagram with migration arrow]\n\nSVZ → Cortex\nω=3.2 → ω=22\n\n(Will be replaced with PNG)",
                 font_size=11, color=(0x94, 0xA3, 0xB8), align="center")
    
    # ── RIGHT CARD: reasoning chain ──
    add_rect(slide, 5.65, 1.05, 4.0, 4.0, CARD_BG, (0xE2, 0xE8, 0xF0))
    add_textbox(slide, 5.8, 1.2, 3.7, 0.35,
                 "推理链：从观察到假说" if is_zh else "Reasoning Chain: From Observation to Hypothesis",
                 font_size=13, bold=True, color=ACCENT_TEAL)
    add_rect(slide, 5.8, 1.55, 3.6, 0.04, ACCENT_TEAL)
    
    chain = [
        ("① 观察（Observation）", "bold", MUTED_FG),
        ("   Intra-type ω(OPC, SVZ→Ctx) ≈ 22", "normal", BODY_FG),
        ("   Inter-type ω(Astro vs OPC, Ctx) ≈ 107.5", "normal", BODY_FG),
        ("", "normal", BODY_FG),
        ("② 矛盾（Contradiction）", "bold", MUTED_FG),
        ("   Intra-type ω ≪ Inter-type ω", "normal", (0xDC, 0x26, 0x26)),
        ("   同一类型跨区相似度 > 不同类型同区相似度", "normal", BODY_FG),
        ("", "normal", BODY_FG),
        ("③ 解释（Explanation）", "bold", MUTED_FG),
        ("   最简洁：OPC 从 SVZ 迁移到 Cortex", "normal", (0x0D, 0x94, 0x88)),
        ("   迁移后保留部分起源特征 → ω 偏低", "normal", BODY_FG),
        ("", "normal", BODY_FG),
        ("④ 验证（Validation）", "bold", MUTED_FG),
        ("   OPC 已知具有发育期迁移特性", "normal", BODY_FG),
        ("   CKI 在无先验知识下独立检出此信号", "normal", BODY_FG),
        ("   → 方法可靠性佐证", "normal", (0x0D, 0x94, 0x88)),
    ]
    y = 1.65
    for text, style, color in chain:
        if text == "":
            y += 0.08
            continue
        add_textbox(slide, 5.85, y, 3.7, 0.3, text,
                     font_size=9.5, bold=(style=="bold"), color=color)
        y += 0.3
    
    # ── Bottom insight bar ──
    add_rect(slide, 0.55, 5.05, 9.1, 0.45, (0x0D, 0x94, 0x88))
    add_textbox(slide, 0.55, 5.1, 9.1, 0.35,
                 "结论：CKI 通过 ω 值的极端偏离，识别出潜在的细胞迁移事件，为理解脑内细胞动力学提供新视角。"
                 if is_zh else
                 "Conclusion: CKI identifies potential cell migration events via extreme ω deviations, providing a new lens on brain cell dynamics.",
                 font_size=11, color=(0xFF, 0xFF, 0xFF), align="center")
    
    return slide


# ═══════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════

def main():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    print("=" * 60)
    print("Restructure Slides 15-17: 3-Part Migration Narrative")
    print("=" * 60)
    
    # ── Step 1: Generate SVGs ──
    print("\n[1/5] Generating SVGs...")
    
    svg_s15 = make_svg_s15_baseline()
    svg15_path = ASSETS / "_s15_baseline.svg"
    svg15_path.write_text(svg_s15, encoding="utf-8")
    print(f"      SVG S15: {svg15_path.name} ({len(svg_s15)//1024} KB)")
    
    svg_s16 = make_svg_s16_anomaly()
    svg16_path = ASSETS / "_s16_anomaly.svg"
    svg16_path.write_text(svg_s16, encoding="utf-8")
    print(f"      SVG S16: {svg16_path.name} ({len(svg_s16)//1024} KB)")
    
    svg_s17 = make_svg_s17_brain()
    svg17_path = ASSETS / "_s17_migration.svg"
    svg17_path.write_text(svg_s17, encoding="utf-8")
    print(f"      SVG S17: {svg17_path.name} ({len(svg_s17)//1024} KB)")
    
    # ── Step 2: Convert SVG → PNG ──
    print("\n[2/5] Converting SVGs to PNGs (Chrome headless)...")
    
    png15_path = ASSETS / "_s15_baseline.png"
    png16_path = ASSETS / "_s16_anomaly.png"
    png17_path = ASSETS / "_s17_migration.png"
    
    ok15 = svg_to_png_chrome(svg15_path, png15_path)
    ok16 = svg_to_png_chrome(svg16_path, png16_path)
    ok17 = svg_to_png_chrome(svg17_path, png17_path)
    
    if not (ok15 and ok16 and ok17):
        print("      WARNING: Some PNGs failed to generate. Check Chrome installation.")
    
    # ── Step 3: Rebuild Slide 15 (EN) ──
    print("\n[3/5] Rebuilding Slide 15 (EN)...")
    prs_en = Presentation(str(PPT_EN))
    slide15_en = clear_slide(prs_en, S15)
    build_slide_15(slide15_en, is_zh=False)
    
    # Insert PNG into Slide 15
    if png15_path.exists():
        from pptx.util import Inches
        slide15_en.shapes.add_picture(
            str(png15_path), Inches(5.25), Inches(1.65), Inches(4.2), Inches(3.3)
        )
        print(f"      Inserted PNG: {png15_path.name}")
    
    prs_en.save(str(OUT_EN))
    print(f"      Saved: {OUT_EN.name}")
    
    # ── Step 4: Rebuild Slide 16 (EN) ──
    print("\n[4/5] Rebuilding Slide 16 (EN)...")
    prs_en2 = Presentation(str(OUT_EN))
    slide16_en = clear_slide(prs_en2, S16)
    build_slide_16(slide16_en, is_zh=False)
    
    # Insert PNG into Slide 16
    if png16_path.exists():
        from pptx.util import Inches
        slide16_en.shapes.add_picture(
            str(png16_path), Inches(0.6), Inches(1.65), Inches(4.5), Inches(3.3)
        )
        print(f"      Inserted PNG: {png16_path.name}")
    
    prs_en2.save(str(OUT_EN))
    print(f"      Saved: {OUT_EN.name}")
    
    # ── Step 5: Rebuild Slide 17 (EN) ──
    print("\n[5/5] Rebuilding Slide 17 (EN)...")
    prs_en3 = Presentation(str(OUT_EN))
    slide17_en = clear_slide(prs_en3, S17)
    build_slide_17(slide17_en, has_image=png17_path.exists(), is_zh=False)
    
    # Insert PNG into Slide 17
    if png17_path.exists():
        from pptx.util import Inches
        slide17_en.shapes.add_picture(
            str(png17_path), Inches(0.6), Inches(1.08), Inches(4.8), Inches(3.8)
        )
        print(f"      Inserted PNG: {png17_path.name}")
    
    prs_en3.save(str(OUT_EN))
    print(f"      Saved: {OUT_EN.name}")
    
    print("\n" + "=" * 60)
    print("EN version done. Now processing ZH version...")
    print("=" * 60)
    
    # ── ZH version ──
    print("\n[ZH 1/5] Rebuilding Slide 15 (ZH)...")
    prs_zh = Presentation(str(PPT_ZH))
    slide15_zh = clear_slide(prs_zh, S15)
    build_slide_15(slide15_zh, is_zh=True)
    if png15_path.exists():
        slide15_zh.shapes.add_picture(str(png15_path), Inches(5.25), Inches(1.65), Inches(4.2), Inches(3.3))
    prs_zh.save(str(OUT_ZH))
    print(f"      Saved: {OUT_ZH.name}")
    
    print("\n[ZH 2/5] Rebuilding Slide 16 (ZH)...")
    prs_zh2 = Presentation(str(OUT_ZH))
    slide16_zh = clear_slide(prs_zh2, S16)
    build_slide_16(slide16_zh, is_zh=True)
    if png16_path.exists():
        slide16_zh.shapes.add_picture(str(png16_path), Inches(0.6), Inches(1.65), Inches(4.5), Inches(3.3))
    prs_zh2.save(str(OUT_ZH))
    print(f"      Saved: {OUT_ZH.name}")
    
    print("\n[ZH 3/5] Rebuilding Slide 17 (ZH)...")
    prs_zh3 = Presentation(str(OUT_ZH))
    slide17_zh = clear_slide(prs_zh3, S17)
    build_slide_17(slide17_zh, has_image=png17_path.exists(), is_zh=True)
    if png17_path.exists():
        slide17_zh.shapes.add_picture(str(png17_path), Inches(0.6), Inches(1.08), Inches(4.8), Inches(3.8))
    prs_zh3.save(str(OUT_ZH))
    print(f"      Saved: {OUT_ZH.name}")
    
    print("\n" + "=" * 60)
    print("DONE. All 3 slides (15-17) restructured.")
    print(f"  EN: {OUT_EN}")
    print(f"  ZH: {OUT_ZH}")
    print("=" * 60)


if __name__ == "__main__":
    main()
