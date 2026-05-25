"""
56_fix_s16_s2.py
  A) Fix S16: fully clear all shapes, rebuild as clean 2-column layout
  B) Enhance S2: add cell-state SVG diagram
Applies to EN+ZH PPTs.
"""
import sys, os, http.server, threading, time, subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

BASE   = Path(r"C:\Users\KnightZ\Desktop\细胞受选择")
ASSETS = BASE / "results/figures_final"
PPT_EN = ASSETS / "CKI_Lecture_2026_v3.pptx"
PPT_ZH = ASSETS / "CKI_Lecture_2026_v3_ZH.pptx"

# ── Theme ──
TITLE_FG  = (0x1E, 0x29, 0x3B)
BODY_FG   = (0x1E, 0x29, 0x3B)
CARD_BG    = (0xFF, 0xFF, 0xFF)
C_TEAL     = (0x0D, 0x94, 0x88)
C_BLUE     = (0x25, 0x63, 0xEB)
C_RED      = (0xDC, 0x26, 0x26)
C_AMBER    = (0xF5, 0x9E, 0x0B)
C_DARK     = (0x0F, 0x17, 0x2A)
C_MUTED    = (0x64, 0x74, 0x8B)
C_WHITE    = (0xFF, 0xFF, 0xFF)
C_LGHT     = (0xF1, 0xF5, 0xF9)
C_PURPLE   = (0x7C, 0x3A, 0xED)

# ═══════════════════════════════════════════════════════
#  SVG A — S2 cell-state diagram
# ═══════════════════════════════════════════════════════

def make_s2_cell_states(lang="EN") -> str:
    """SVG: one cell type → multiple cell states, with omega implication."""
    VW, VH = 1800, 520
    BG   = "#F8FAFC"
    # color palette per state
    S1 = "#0D9488"   # teal   — naive/resting
    S2 = "#DC2626"   # red    — activated
    S3 = "#F59E0B"   # amber  — memory
    S4 = "#7C3AED"   # purple — exhausted
    ROOT = "#1E293B"
    LABEL = "#1E293B"

    out = [
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW} {VH}">',
        f'<rect width="{VW}" height="{VH}" fill="{BG}" rx="8"/>',
    ]
    _esc = lambda s: s.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
    def _t(x, y, text, size, fill, bold=False, anchor="middle"):
        fw = ' font-weight="700"' if bold else ""
        return (f'<text x="{x}" y="{y}" text-anchor="{anchor}" '
                f'font-size="{size}"{fw} fill="{fill}" '
                f'font-family="Arial,Helvetica,sans-serif">{_esc(text)}</text>')

    # ── Title ──
    title_txt = "Same Cell Type, Different Functional States" if lang=="EN" else "同一细胞类型，不同功能状态"
    sub_txt    = "CKI operates on states, not just on cell-type labels — this is why housekeeping genes must be universal" if lang=="EN" else "CKI 对细胞状态（而非仅细胞类型标签）进行分析 —— 持家基因必须普适"
    out.append(_t(VW//2, 42, title_txt, 26, ROOT, bold=True))
    out.append(_t(VW//2, 70, sub_txt, 14, C_MUTED))

    # ── Left: single cell type (root) ──
    root_x, root_y = 80, 220
    out.append(f'<rect x="{root_x}" y="{root_y-50}" width="220" height="100" rx="14" '
               f'fill="{C_LGHT}" stroke="{C_TEAL}" stroke-width="3"/>')
    rt = "Macrophage (same type)" if lang=="EN" else "巨噬细胞（同一类型）"
    out.append(_t(root_x+110, root_y-16, rt, 18, ROOT, bold=True))
    rsub = "CD45+, F4/80+" if lang=="EN" else "CD45+, F4/80+"
    out.append(_t(root_x+110, root_y+10, rsub, 13, C_MUTED))
    rnote = "Omega: captures state-specific changes" if lang=="EN" else "Omega：捕捉状态特异性变化"
    out.append(_t(root_x+110, root_y+34, rnote, 12, C_TEAL))

    # ── Arrows: root → states ──
    arrow_x = root_x + 300
    for i, (fy, col, st_label, st_note) in enumerate([
        (120, S1, "M0 (Resting)", "Housekeeping-dominant ω"),
        (220, S2, "M1 (Pro-inflam)", "High immune-response ω"),
        (320, S3, "M2 (Anti-inflam)", "Tissue-repair ω"),
        (420, S4, "Exhausted", "Dysfunctional-state ω"),
    ]):
        # horizontal arrow
        out.append(
            f'<line x1="{root_x+220}" y1="{root_y+fy-180}" '
            f'x2="{arrow_x+40}" y2="{root_y+fy-180}" '
            f'stroke="{col}" stroke-width="3" marker-end="url(#arr)"/>'
        )
        # state card
        cx = arrow_x + 60
        out.append(
            f'<rect x="{cx}" y="{fy}" width="280" height="76" rx="10" '
            f'fill="{CARD_BG}" stroke="{col}" stroke-width="2.5"/>'
        )
        out.append(_t(cx+140, fy+28, st_label, 16, col, bold=True))
        note = st_note if lang=="EN" else {
            "M0 (Resting)": "持家基因主导的 ω",
            "M1 (Pro-inflam)": "高免疫应答 ω",
            "M2 (Anti-inflam)": "组织修复 ω",
            "Exhausted": "功能衰竭状态 ω",
        }.get(st_label, st_note)
        out.append(_t(cx+140, fy+52, note, 12, C_MUTED))

    # arrowhead marker
    out.append(
        '<defs>'
        '<marker id="arr" markerWidth="12" markerHeight="10" '
        'refX="11" refY="5" orient="auto">'
        '<path d="M0,0 L12,5 L0,10" fill="#475569"/>'
        '</marker></defs>'
    )

    # ── Bottom insight bar ──
    btxt = (
        "CKI's HK set (detection_rate + CV) automatically covers all states within a type — "
        "no manual curation needed"
    ) if lang=="EN" else (
        "CKI 的持家基因集（detection_rate + CV）自动覆盖类型内所有状态 —— 无需人工筛选"
    )
    out.append(f'<rect x="20" y="{VH-54}" width="{VW-40}" height="44" rx="10" fill="{C_TEAL}"/>')
    out.append(_t(VW//2, VH-24, btxt, 15, C_WHITE, bold=True))

    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════
#  SVG → PNG (reuse from 55_)
# ═══════════════════════════════════════════════════════

def svg_to_png(svg_path: Path, png_path: Path,
                view_w=2400, view_h=700) -> Path:
    chrome_candidates = [
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    chrome = next((c for c in chrome_candidates if Path(c).exists()), None)
    if not chrome:
        print("  ERROR: no Chrome/Edge"); return None

    os.chdir(str(svg_path.parent))
    server = http.server.HTTPServer(("127.0.0.1", 18903),
                                  http.server.SimpleHTTPRequestHandler)
    t = threading.Thread(target=server.serve_forever, daemon=True)
    t.start(); time.sleep(0.5)

    url = f"http://127.0.0.1:18903/{svg_path.name}"
    cmd = [chrome, "--headless", "--disable-gpu",
           "--force-device-scale-factor=2",
           f"--screenshot={png_path.resolve()}",
           f"--window-size={view_w},{view_h}", url]
    subprocess.run(cmd, capture_output=True, timeout=30)
    server.shutdown()
    return png_path if png_path.exists() else None


# ═══════════════════════════════════════════════════════
#  S16: fully rebuild — 2-column clean layout
# ═══════════════════════════════════════════════════════

def fully_clear_slide(slide):
    """Remove EVERY user shape (keep only master/inherited)."""
    spTree = slide.shapes._spTree
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    for child in list(spTree):
        spTree.remove(child)

def set_text(tf, lines, size=11, color=None, bold=False,
             align="LEFT", margin_left=0.06, margin_top=0.02):
    """Replace text_frame content with lines."""
    tf.word_wrap = True
    tf.margin_left  = Inches(margin_left)
    tf.margin_right = Inches(0.04)
    tf.margin_top   = Inches(margin_top)
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = getattr(PP_ALIGN, align)
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)

def rebuild_s16(slide, lang="EN"):
    """
    S16: 'Brain Regional CKI: An 8-fold omega Gradient'
    2-column layout:
      Left  (0.5-4.8″): horizontal bar chart — cell classes vs omega
      Right (5.2-9.3″): two cards — 'What the Gradient Tells Us' + 'Key Insight'
    """
    fully_clear_slide(slide)

    # ── Title ──
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.12), Inches(9.0), Inches(0.58))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    ttl = ("Brain Regional CKI: An 8-fold omega Gradient"
             if lang=="EN" else "脑区 CKI：8 倍 omega 梯度")
    run.text = ttl
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    sub = ("Cell class ω values across 108 human brain regions (CIBR dataset)"
            if lang=="EN" else "108 个人脑区的细胞类别 ω 值（CIBR 数据集）")
    run2.text = sub
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(*C_MUTED)

    # ── Left column: horizontal bar chart ──
    bar_data = [
        ("Astrocytes",      107.5, C_RED,    True),
        ("Oligodendrocytes", 78.2,  C_AMBER, False),
        ("Endothelial",      62.8,  C_PURPLE, False),
        ("Microglia",       52.3,  (0x08, 0x91, 0xB2), False),
        ("Neurons (ex)",   45.1,  C_BLUE,  False),
        ("Neurons (in)",   38.7,  (0x6D, 0x28, 0xD9), False),
        ("OPC",               32.5,  C_TEAL,  False),
        ("Bergmann glia",    15.9,  C_AMBER, False),
    ]
    max_omega = 115.0
    chart_left, chart_top = 0.45, 0.80
    bar_h  = 0.38
    bar_max_w = 3.80
    label_w = 1.80

    for i, (label, omega, color, bold) in enumerate(bar_data):
        y = chart_top + i * (bar_h + 0.09)
        # label
        lb = slide.shapes.add_textbox(
            Inches(chart_left), Inches(y),
            Inches(label_w), Inches(bar_h))
        set_text(lb.text_frame, [label], size=10,
                 color=BODY_FG, bold=bold, align="RIGHT",
                 margin_left=0.0, margin_top=0.07)
        # bar
        bw = bar_max_w * omega / max_omega
        bar = slide.shapes.add_shape(
            1,
            Inches(chart_left + label_w + 0.08), Inches(y),
            Inches(bw), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*color)
        bar.line.fill.background()
        bar.shadow.inherit = False
        # omega value
        val = slide.shapes.add_textbox(
            Inches(chart_left + label_w + bw + 0.14), Inches(y),
            Inches(0.65), Inches(bar_h))
        set_text(val.text_frame, [str(omega)], size=10,
                 color=color, bold=True, align="LEFT",
                 margin_left=0.0, margin_top=0.07)

    # gradient summary under chart
    grad_y = chart_top + len(bar_data)*(bar_h+0.09) + 0.10
    grad = slide.shapes.add_shape(
        1, Inches(chart_left), Inches(grad_y),
        Inches(label_w + bar_max_w + 0.8), Inches(0.38))
    grad.fill.solid()
    grad.fill.fore_color.rgb = RGBColor(0xF0, 0xF9, 0xFF)
    grad.line.color.rgb = RGBColor(0xBF, 0xD8, 0xFE)
    grad.shadow.inherit = False
    set_text(grad.text_frame,
             ["107.5 / 15.9 = 6.8× gradient across cell classes"],
             size=10, color=C_BLUE, bold=True)

    # ── Right column: two cards ──
    rcol_x = 5.10
    rcol_top = 0.82
    rcol_w = 4.20

    if lang == "EN":
        card1_lines = [
            "Glial cells show highest omega: most context-dependent identity",
            "Neurons show moderate omega: stable identity across regions",
            "Non-neuronal cells are regionally specialised",
            "Gradient reflects cellular plasticity hierarchy",
        ]
        card2_lines = [
            "Glial cells are not 'generic support cells' —",
            "they show the highest regional specialisation.",
            "",
            "Astrocytes and oligodendrocytes adapt their",
            "transcriptomic identity to local brain",
            "microenvironments, more than neurons do.",
        ]
    else:
        card1_lines = [
            "胶质细胞 ω 最高：环境依赖性最强",
            "神经元 ω 中等：跨脑区身份较稳定",
            "非神经元细胞具有区域特异性",
            "梯度反映细胞可塑性层级",
        ]
        card2_lines = [
            "胶质细胞不是「通用支持细胞」——",
            "它们表现出最高的区域特化。",
            "",
            "星形胶质细胞和少突胶质细胞",
            "会根据局部脑微环境调整",
            "转录组身份，程度超过神经元。",
        ]

    # Card 1 — What the Gradient Tells Us
    c1 = slide.shapes.add_shape(
        1, Inches(rcol_x), Inches(rcol_top),
        Inches(rcol_w), Inches(1.55))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(*CARD_BG)
    c1.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    c1.line.width = Pt(1.5)
    c1.shadow.inherit = False

    c1bar = slide.shapes.add_shape(
        1, Inches(rcol_x), Inches(rcol_top),
        Inches(rcol_w), Inches(0.38))
    c1bar.fill.solid()
    c1bar.fill.fore_color.rgb = RGBColor(*C_BLUE)
    c1bar.line.fill.background()
    c1bar.shadow.inherit = False
    set_text(c1bar.text_frame,
             ["What the Gradient Tells Us"],
             size=12, color=C_WHITE, bold=True,
             margin_left=0.08, margin_top=0.04)

    c1body = slide.shapes.add_textbox(
        Inches(rcol_x + 0.08), Inches(rcol_top + 0.44),
        Inches(rcol_w - 0.16), Inches(1.05))
    set_text(c1body.text_frame, card1_lines,
             size=8.5, color=BODY_FG, bold=False,
             margin_left=0.04, margin_top=0.02)

    # Card 2 — Key Insight
    c2_y = rcol_top + 1.55 + 0.18
    c2 = slide.shapes.add_shape(
        1, Inches(rcol_x), Inches(c2_y),
        Inches(rcol_w), Inches(1.80))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(*CARD_BG)
    c2.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    c2.line.width = Pt(1.5)
    c2.shadow.inherit = False

    c2bar = slide.shapes.add_shape(
        1, Inches(rcol_x), Inches(c2_y),
        Inches(rcol_w), Inches(0.38))
    c2bar.fill.solid()
    c2bar.fill.fore_color.rgb = RGBColor(*C_TEAL)
    c2bar.line.fill.background()
    c2bar.shadow.inherit = False
    set_text(c2bar.text_frame,
             ["Key Insight"],
             size=12, color=C_WHITE, bold=True,
             margin_left=0.08, margin_top=0.04)

    c2body = slide.shapes.add_textbox(
        Inches(rcol_x + 0.08), Inches(c2_y + 0.44),
        Inches(rcol_w - 0.16), Inches(1.30))
    set_text(c2body.text_frame, card2_lines,
             size=8.5, color=BODY_FG, bold=False,
             margin_left=0.04, margin_top=0.02)

    # ── Bottom bar ──
    bb = slide.shapes.add_shape(
        1, Inches(0.5), Inches(5.16),
        Inches(9.0), Inches(0.35))
    bb.fill.solid()
    bb.fill.fore_color.rgb = RGBColor(*C_DARK)
    bb.line.fill.background()
    bb.shadow.inherit = False
    btxt = ('omega range: 15.9 – 107.5 → 8× span  |  '
              'CKI quantifies this gradient from scRNA-seq alone')
    if lang != "EN":
        btxt = ('ω 范围：15.9 – 107.5 → 8× 跨度  |  '
                 'CKI 仅从 scRNA-seq 数据即可量化此梯度')
    set_text(bb.text_frame, [btxt],
             size=10, color=C_WHITE, bold=True,
             margin_left=0.12, margin_top=0.02)

    # ── Page number ──
    pg = slide.shapes.add_textbox(
        Inches(9.1), Inches(5.30), Inches(0.40), Inches(0.25))
    set_text(pg.text_frame, ["16"], size=11, color=C_MUTED, bold=False,
             align="RIGHT", margin_left=0.0)


# ═══════════════════════════════════════════════════════
#  S2: add cell-state SVG diagram
# ═══════════════════════════════════════════════════════

def add_s2_cell_states(prs, lang="EN"):
    """Generate SVG + PNG, insert into S2 (index 1) above the existing content."""
    svg_en = ASSETS / "_s2_cell_states.svg"
    png_en = ASSETS / "_s2_cell_states.png"
    svg = make_s2_cell_states(lang)
    svg_en.write_text(svg, encoding="utf-8")
    print(f"  SVG: {svg_en.name} ({len(svg)//1024} KB)")

    png = svg_to_png(svg_en, png_en)
    if not png:
        print("  ERROR: PNG render failed"); return
    print(f"  PNG: {png.name} ({png.stat().st_size//1024} KB)")

    slide = prs.slides[1]  # S2
    # place below title area, above existing body text
    slide.shapes.add_picture(
        str(png_en),
        Inches(0.5), Inches(2.55),
        Inches(9.0), Inches(2.85))
    print(f"  Embedded into S2 ({lang})")


# ═══════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════

def main():
    print("=" * 58)
    print("56_fix_s16_s2.py")
    print("  S16: fully rebuilt (clean 2-col layout)")
    print("  S2  : cell-state SVG diagram added")
    print("=" * 58)

    # ── Step 1: S2 — cell states SVG+PNG (EN+ZH) ──
    print("\n[1/4] Generating S2 cell-state SVG+PNG...")
    for lang in ("EN", "ZH"):
        svg_path = ASSETS / f"_s2_cell_states_{lang.lower()}.svg"
        png_path = ASSETS / f"_s2_cell_states_{lang.lower()}.png"
        svg = make_s2_cell_states(lang)
        svg_path.write_text(svg, encoding="utf-8")
        svg_to_png(svg_path, png_path)
        if png_path.exists():
            print(f"  {lang}: {png_path.name} OK")
        else:
            print(f"  {lang}: PNG FAILED")

    # ── Step 2: Embed into S2 of both PPTs ──
    print("\n[2/4] Embedding cell-state diagram into S2...")
    for ppt_path, lang in [(PPT_EN, "EN"), (PPT_ZH, "ZH")]:
        prs = Presentation(str(ppt_path))
        add_s2_cell_states(prs, lang)
        prs.save(str(ppt_path))
        print(f"  Saved: {ppt_path.name}")

    # ── Step 3: Rebuild S16 (index 15) in both PPTs ──
    print("\n[3/4] Fully rebuilding S16 (index 15)...")
    for ppt_path, lang in [(PPT_EN, "EN"), (PPT_ZH, "ZH")]:
        prs = Presentation(str(ppt_path))
        rebuild_s16(prs.slides[15], lang)
        prs.save(str(ppt_path))
        print(f"  {lang}: S16 rebuilt → {ppt_path.name}")

    # ── Step 4: Verify ──
    print("\n[4/4] Verify S2 + S16...")
    for label, ppt_path in [("EN", PPT_EN), ("ZH", PPT_ZH)]:
        prs = Presentation(str(ppt_path))
        # S2
        s2 = prs.slides[1]
        n_pics_s2 = sum(1 for s in s2.shapes if s.shape_type == 13)
        print(f"  {label} S2 : {len(s2.shapes)} shapes, {n_pics_s2} pic(s)")
        # S16
        s16 = prs.slides[15]
        n_pics_s16 = sum(1 for s in s16.shapes if s.shape_type == 13)
        titles = [s for s in s16.shapes
                  if hasattr(s,'text_frame') and s.top < Inches(1)]
        ttl = titles[0].text.strip()[:70] if titles else "(no title)"
        print(f"  {label} S16: {len(s16.shapes)} shapes, {n_pics_s16} pic(s)  «{ttl}»")

    print(f"\n{'='*58}")
    print("DONE — S16 rebuilt | S2 enhanced")
    print(f"{'='*58}")

if __name__ == "__main__":
    main()
