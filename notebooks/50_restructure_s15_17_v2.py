"""
notebooks/50_restructure_s15_17_v2.py

Restructure Slides 15-17 of CKI Lecture PPT into a clear 3-part narrative:
  Part 1 (Slide 15): BASELINE — intra-region ω < inter-region ω (expected pattern)
  Part 2 (Slide 16): ANOMALY — OPC breaks the baseline (extreme case)
  Part 3 (Slide 17): INFERENCE — migration hypothesis (from anomaly to story)

Strategy: DO NOT DELETE shapes. Instead:
  - Modify existing shape text/colors in-place
  - Replace image via: remove old pic element from XML, add new pic
  - Add NEW shapes on top of existing ones (old ones become hidden underneath)

EN + ZH versions.
"""
import sys, http.server, threading, time, os, subprocess
from pathlib import Path
from copy import deepcopy

# ── Paths ─────────────────────────────────────────────────────────────
BASE      = Path(r"C:\Users\KnightZ\Desktop\细胞受选择")
PPT_EN    = BASE / "results/figures_final/CKI_Lecture_2026_v3.pptx"
PPT_ZH    = BASE / "results/figures_final/CKI_Lecture_2026_v3_ZH.pptx"
OUT_EN    = BASE / "results/figures_final/CKI_Lecture_2026_v3.pptx"
OUT_ZH    = BASE / "results/figures_final/CKI_Lecture_2026_v3_ZH.pptx"
ASSETS    = BASE / "results/figures_final"
ASSETS.mkdir(parents=True, exist_ok=True)

# ── Theme palette ──────────────────────────────────────────────────────
TITLE_FG   = (0x1E, 0x29, 0x3B)
BODY_FG    = (0x1E, 0x29, 0x3B)
MUTED_FG   = (0x64, 0x74, 0x8B)
CARD_BG    = (0xFF, 0xFF, 0xFF)
ACCENT_TEAL= (0x0D, 0x94, 0x88)
ACCENT_BLUE= (0x25, 0x63, 0xEB)
ACCENT_RED  = (0xDC, 0x26, 0x26)
ACCENT_AMBER= (0xF5, 0x9E, 0x0B)
DARK_BG    = (0x0F, 0x17, 0x2A)

def rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


# ═════════════════════════════════════════════════════════════════
#  SVG helpers
# ═════════════════════════════════════════════════════════════════

def svg_escape(s: str) -> str:
    return s.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

def omega_color(omega: float) -> str:
    t = max(0.0, min(1.0, (omega - 3.0) / 107.0))
    # Blue -> Purple -> Red
    r = int(37  + (220-37)  * t)
    g = int(99  + (38 -99)  * t)
    b = int(235 + (38-235) * t)
    return f"rgb({r},{g},{b})"


# ═════════════════════════════════════════════════════════════════
#  SVG A — Slide 15: Baseline (Intra vs Inter ω)
# ═════════════════════════════════════════════════════════════════

def make_svg_s15_baseline() -> str:
    W, H = 1200, 520
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
          f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']
    
    # ── Title ──
    out.append(f'<text x="{W//2}" y="34" text-anchor="middle" '
               f'font-size="17" font-weight="700" fill="#0F172A">'
               f'Baseline: Intra-region ω &lt; Inter-region ω</text>')
    out.append(f'<text x="{W//2}" y="56" text-anchor="middle" '
               f'font-size="11" fill="#64748B">'
               f'Expected pattern — same-type cells across regions are MORE different than '
               f'different cell types within the same region</text>')
    
    # ── LEFT PANEL: Intra-region ──
    lx, ly = 40, 85
    out.append(f'<rect x="{lx}" y="{ly}" width="480" height="395" rx="10" '
               f'fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+240}" y="{ly+26}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#1E293B">Intra-region ω</text>')
    out.append(f'<text x="{lx+240}" y="{ly+46}" text-anchor="middle" '
               f'font-size="10" fill="#64748B">Same brain region · Different cell types</text>')
    
    intra = [("Astrocytes", 28, "#0D9488"), ("OPCs", 22, "#2563EB"), ("Neurons(ex)", 18, "#7C3AED")]
    for i, (ctype, omega, color) in enumerate(intra):
        y = ly + 75 + i * 75
        bw = int(200 * omega / 35.0)
        out.append(f'<rect x="{lx+40}" y="{y}" width="{bw}" height="38" rx="4" fill="{color}" opacity="0.85"/>')
        out.append(f'<text x="{lx+40+bw+8}" y="{y+25}" font-size="13" fill="{color}" font-weight="600">{omega}</text>')
        out.append(f'<text x="{lx+300}" y="{y+25}" font-size="12" fill="#334155">{ctype}</text>')
    
    # Insight box (left)
    iy = ly + 310
    out.append(f'<rect x="{lx+15}" y="{iy}" width="450" height="65" rx="8" '
               f'fill="#F0FDF4" stroke="#86EFAC" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+240}" y="{iy+22}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#166534">Intra-region ω = LOW</text>')
    out.append(f'<text x="{lx+240}" y="{iy+42}" text-anchor="middle" '
               f'font-size="10" fill="#166534">Same microenvironment → similar expression</text>')
    
    # ── CENTER divider ──
    cx = 580
    out.append(f'<line x1="{cx}" y1="{ly+30}" x2="{cx}" y2="{ly+380}" '
               f'stroke="#CBD5E1" stroke-width="2" stroke-dasharray="6,4"/>')
    out.append(f'<text x="{cx}" y="{ly+200}" text-anchor="middle" '
               f'font-size="20" fill="#94A3B8" font-weight="700">&lt;</text>')
    out.append(f'<text x="{cx}" y="{ly+222}" text-anchor="middle" '
               f'font-size="9" fill="#94A3B8">BASELINE</text>')
    
    # ── RIGHT PANEL: Inter-region ──
    rx, ry = 620, 85
    out.append(f'<rect x="{rx}" y="{ry}" width="540" height="395" rx="10" '
               f'fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{ry+26}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#1E293B">Inter-region ω</text>')
    out.append(f'<text x="{rx+270}" y="{ry+46}" text-anchor="middle" '
               f'font-size="10" fill="#64748B">Same cell type · Different brain regions</text>')
    
    inter = [("OPC · SVZ", 3.2, "#0D9488"), ("OPC · Hipp", 18, "#2563EB"), ("OPC · Cortex", 22, "#7C3AED")]
    for j, (label, omega, color) in enumerate(inter):
        y = ry + 75 + j * 75
        bw = int(200 * omega / 35.0)
        out.append(f'<rect x="{rx+40}" y="{y}" width="{bw}" height="38" rx="4" fill="{color}" opacity="0.85"/>')
        out.append(f'<text x="{rx+40+bw+8}" y="{y+25}" font-size="13" fill="{color}" font-weight="600">{omega}</text>')
        out.append(f'<text x="{rx+310}" y="{y+25}" font-size="12" fill="#334155">{label}</text>')
    
    # Insight box (right)
    riy = ry + 310
    out.append(f'<rect x="{rx+15}" y="{riy}" width="510" height="65" rx="8" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{riy+22}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#991B1B">Inter-region ω = HIGH</text>')
    out.append(f'<text x="{rx+270}" y="{riy+42}" text-anchor="middle" '
               f'font-size="10" fill="#991B1B">Different microenvironments → distinct expression</text>')
    
    # ── Bottom bar ──
    out.append(f'<rect x="35" y="{H-48}" width="{W-70}" height="38" rx="8" fill="#0F172A"/>')
    out.append(f'<text x="{W//2}" y="{H-24}" text-anchor="middle" '
               f'font-size="12" font-weight="600" fill="white">'
               f'Baseline:  Intra-region ω &lt; Inter-region ω  →  '
               f'Cells differentiate LOCALLY, not by migrating</text>')
    
    out.append('</svg>')
    return "\n".join(out)


# ═════════════════════════════════════════════════════════════════
#  SVG B — Slide 16: Anomaly (OPC breaks baseline)
# ═════════════════════════════════════════════════════════════════

def make_svg_s16_anomaly() -> str:
    W, H = 1200, 540
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
          f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']
    
    # Title
    out.append(f'<text x="{W//2}" y="34" text-anchor="middle" '
               f'font-size="17" font-weight="700" fill="#0F172A">'
               f'Extreme Case: OPC Breaks the Baseline</text>')
    out.append(f'<text x="{W//2}" y="56" text-anchor="middle" '
               f'font-size="11" fill="#64748B">'
               f'Intra-type ω ≪ expected inter-region ω → anomaly worth investigating</text>')
    
    # ── LEFT: Intra-type ω (OPC across regions) ──
    lx, ly = 40, 90
    out.append(f'<rect x="{lx}" y="{ly}" width="540" height="395" rx="10" '
               f'fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+270}" y="{ly+28}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#1E293B">Intra-type ω (OPC)</text>')
    out.append(f'<text x="{lx+270}" y="{ly+48}" text-anchor="middle" '
               f'font-size="10" fill="#64748B">OPC in SVZ vs OPC in Cortex</text>')
    
    # Two bars
    bar_max = 280
    # SVZ ω=3.2 (tiny)
    svz_h = int(bar_max * 3.2 / 110.0)
    out.append(f'<rect x="{lx+60}" y="{ly+90 + bar_max - svz_h}" width="90" height="{svz_h}" '
               f'rx="4" fill="#0D9488"/>')
    out.append(f'<text x="{lx+105}" y="{ly+90 + bar_max - svz_h - 8}" '
               f'text-anchor="middle" font-size="12" font-weight="700" fill="#0D9488">ω=3.2</text>')
    out.append(f'<text x="{lx+105}" y="{ly+90+bar_max+18}" text-anchor="middle" '
               f'font-size="10" fill="#334155">SVZ</text>')
    
    # Cortex ω=22 (small)
    ctx_h = int(bar_max * 22.0 / 110.0)
    out.append(f'<rect x="{lx+200}" y="{ly+90 + bar_max - ctx_h}" width="90" height="{ctx_h}" '
               f'rx="4" fill="#2563EB"/>')
    out.append(f'<text x="{lx+245}" y="{ly+90 + bar_max - ctx_h - 8}" '
               f'text-anchor="middle" font-size="12" font-weight="700" fill="#2563EB">ω=22</text>')
    out.append(f'<text x="{lx+245}" y="{ly+90+bar_max+18}" text-anchor="middle" '
               f'font-size="10" fill="#334155">Cortex</text>')
    
    # Diff annotation
    out.append(f'<text x="{lx+330}" y="{ly+240}" font-size="11" fill="#64748B" font-style="italic">'
               f'ω diff = 18.8  (surprisingly SMALL)</text>')
    
    # Insight (left)
    out.append(f'<rect x="{lx+15}" y="{ly+340}" width="510" height="75" rx="8" '
               f'fill="#F0FDF4" stroke="#86EFAC" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+270}" y="{ly+362}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#166534">Intra-type ω = 22 (LOW)</text>')
    out.append(f'<text x="{lx+270}" y="{ly+382}" text-anchor="middle" '
               f'font-size="10" fill="#166534">OPCs in SVZ and Cortex are HIGHLY similar</text>')
    out.append(f'<text x="{lx+270}" y="{ly+400}" text-anchor="middle" '
               f'font-size="10" fill="#166534">→ likely same origin (migration), NOT independent differentiation</text>')
    
    # ── RIGHT: Inter-type ω (within Cortex) ──
    rx, ry = 620, 90
    out.append(f'<rect x="{rx}" y="{ry}" width="540" height="395" rx="10" '
               f'fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{ry+28}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#1E293B">Inter-type ω (Cortex)</text>')
    out.append(f'<text x="{rx+270}" y="{ry+48}" text-anchor="middle" '
               f'font-size="10" fill="#64748B">Astrocytes vs OPCs (same region)</text>')
    
    # Two bars
    # Astrocyte ω=107.5 (TALL)
    ast_h = int(bar_max * 107.5 / 110.0)
    out.append(f'<rect x="{rx+60}" y="{ry+90 + bar_max - ast_h}" width="90" height="{ast_h}" '
               f'rx="4" fill="#DC2626"/>')
    out.append(f'<text x="{rx+105}" y="{ry+90 + bar_max - ast_h - 8}" '
               f'text-anchor="middle" font-size="12" font-weight="700" fill="#DC2626">ω=107.5</text>')
    out.append(f'<text x="{rx+105}" y="{ry+90+bar_max+18}" text-anchor="middle" '
               f'font-size="10" fill="#334155">Astrocytes</text>')
    
    # OPC ω=22 (small)
    opc_h = int(bar_max * 22.0 / 110.0)
    out.append(f'<rect x="{rx+200}" y="{ry+90 + bar_max - opc_h}" width="90" height="{opc_h}" '
               f'rx="4" fill="#2563EB"/>')
    out.append(f'<text x="{rx+245}" y="{ry+90 + bar_max - opc_h - 8}" '
               f'text-anchor="middle" font-size="12" font-weight="700" fill="#2563EB">ω=22</text>')
    out.append(f'<text x="{rx+245}" y="{ry+90+bar_max+18}" text-anchor="middle" '
               f'font-size="10" fill="#334155">OPCs</text>')
    
    # Diff annotation
    out.append(f'<text x="{rx+330}" y="{ry+240}" font-size="11" fill="#64748B" font-style="italic">'
               f'ω diff = 85.5  (extremely LARGE)</text>')
    
    # Insight (right)
    out.append(f'<rect x="{rx+15}" y="{ry+340}" width="510" height="75" rx="8" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{ry+362}" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#991B1B">Inter-type ω = 107.5 (HIGH)</text>')
    out.append(f'<text x="{rx+270}" y="{ry+382}" text-anchor="middle" '
               f'font-size="10" fill="#991B1B">Astrocytes and OPCs in SAME region are RADICALLY different</text>')
    out.append(f'<text x="{rx+270}" y="{ry+400}" text-anchor="middle" '
               f'font-size="10" fill="#991B1B">→ independent lineages, strong local differentiation</text>')
    
    # ── Bottom: CONTRADICTION ──
    out.append(f'<rect x="35" y="{H-55}" width="{W-70}" height="44" rx="8" fill="#7C3AED"/>')
    out.append(f'<text x="{W//2}" y="{H-32}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="white">'
               f'CONTRADICTION:  Intra-type ω (22) ≪ Inter-type ω (107.5)  →  '
               f'OPC similarity across regions EXCEEDS cell-type difference within region</text>')
    
    out.append('</svg>')
    return "\n".join(out)


# ═════════════════════════════════════════════════════════════════
#  SVG C — Slide 17: Brain with migration arrow
# ═════════════════════════════════════════════════════════════════

def make_svg_s17_brain() -> str:
    W, H = 1200, 520
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
          f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']
    
    # ── Brain regions (simple ellipses) ──
    # Cerebellum (bottom-left)
    out.append('<ellipse cx="230" cy="395" rx="120" ry="90" fill="#E0F2FE" stroke="#0EA5E9" stroke-width="2.5" opacity="0.55"/>')
    out.append(f'<text x="230" y="425" text-anchor="middle" font-size="11" fill="#0369A1" font-weight="600">Cerebellum</text>')
    out.append(f'<text x="230" y="442" text-anchor="middle" font-size="9" fill="#0369A1">ω=85  Bergmann glia</text>')
    
    # Cortex (top-right)
    out.append('<ellipse cx="680" cy="145" rx="155" ry="105" fill="#FEF2F2" stroke="#EF4444" stroke-width="2.5" opacity="0.55"/>')
    out.append(f'<text x="680" y="122" text-anchor="middle" font-size="12" fill="#991B1B" font-weight="600">Cortex</text>')
    out.append(f'<text x="680" y="138" text-anchor="middle" font-size="10" fill="#991B1B">ω=107.5 Astrocytes</text>')
    out.append(f'<text x="680" y="154" text-anchor="middle" font-size="10" fill="#991B1B">ω=22  OPCs</text>')
    
    # SVZ (left of cortex)
    out.append('<ellipse cx="510" cy="350" rx="65" ry="55" fill="#F0FDF4" stroke="#16A34A" stroke-width="2.5" opacity="0.6"/>')
    out.append(f'<text x="510" y="342" text-anchor="middle" font-size="11" fill="#166534" font-weight="600">SVZ</text>')
    out.append(f'<text x="510" y="360" text-anchor="middle" font-size="10" fill="#166534">ω=3.2  OPCs</text>')
    
    # Hippocampus
    out.append('<ellipse cx="790" cy="310" rx="72" ry="52" fill="#FEFCE8" stroke="#F59E0B" stroke-width="2" opacity="0.5"/>')
    out.append(f'<text x="790" y="307" text-anchor="middle" font-size="10" fill="#92400E">Hipp ω=45</text>')
    
    # Striatum
    out.append('<ellipse cx="630" cy="272" rx="62" ry="48" fill="#F5F3FF" stroke="#7C3AED" stroke-width="2" opacity="0.5"/>')
    out.append(f'<text x="630" y="275" text-anchor="middle" font-size="10" fill="#5B21B6">Striatum ω=38</text>')
    
    # ── MIGRATION ARROW: SVZ → Cortex ──
    out.append('<defs>'
               '<filter id="glow"><feGaussianBlur stdDeviation="4" result="blur"/>'
               '<feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge></filter>'
               '<marker id="arrRed" markerWidth="12" markerHeight="9" refX="10" refY="4.5" orient="auto">'
               '<path d="M0,0 L12,4.5 L0,9" fill="#DC2626"/></marker>'
               '</defs>')
    out.append(f'<path d="M 540 350 C 600 260, 640 210, 670 165" '
               f'stroke="#DC2626" stroke-width="5" fill="none" '
               f'stroke-dasharray="10,5" opacity="0.9" filter="url(#glow)" '
               f'marker-end="url(#arrRed)"/>')
    
    # Migration label
    out.append(f'<rect x="540" y="268" width="145" height="26" rx="13" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="612" y="286" text-anchor="middle" '
               f'font-size="12" font-weight="700" fill="#DC2626">OPC Migration</text>')
    out.append(f'<text x="612" y="312" text-anchor="middle" '
               f'font-size="11" fill="#94A3B8" font-style="italic">ω 3.2 → 22</text>')
    
    # ── Title ──
    out.append(f'<text x="{W//2}" y="28" text-anchor="middle" '
               f'font-size="16" font-weight="700" fill="#0F172A">'
               f'Migration Inference: From Anomaly to Hypothesis</text>')
    out.append(f'<text x="{W//2}" y="48" text-anchor="middle" '
               f'font-size="11" fill="#64748B">'
               f'OPC migration from SVZ to Cortex explains the ω anomaly</text>')
    
    # ── Bottom bar ──
    out.append(f'<rect x="35" y="{H-50}" width="{W-70}" height="40" rx="8" fill="#0D9488"/>')
    out.append(f'<text x="{W//2}" y="{H-27}" text-anchor="middle" '
               f'font-size="12" font-weight="600" fill="white">'
               f'Conclusion: CKI identifies potential cell migration events via extreme ω deviations, '
               f'providing a new lens on brain cell dynamics.</text>')
    
    out.append('</svg>')
    return "\n".join(out)


# ═════════════════════════════════════════════════════════════════
#  SVG → PNG (Chrome headless)
# ═════════════════════════════════════════════════════════════════

def svg_to_png_chrome(svg_path: Path, png_path: Path) -> bool:
    import shutil
    chrome = r"C:\Program Files\Google\Chrome\Application\chrome.exe"
    if not Path(chrome).exists():
        chrome = r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"
    if not Path(chrome).exists():
        chrome = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
    if not Path(chrome).exists():
        chrome = r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"
    if not Path(chrome).exists():
        print("      ERROR: Chrome/Edge not found")
        return False
    
    # HTTP server
    os.chdir(str(svg_path.parent))
    server = http.server.HTTPServer(("127.0.0.1", 18900), http.server.SimpleHTTPRequestHandler)
    t = threading.Thread(target=server.serve_forever, daemon=True)
    t.start()
    time.sleep(0.5)
    url = f"http://127.0.0.1:18900/{svg_path.name}"
    
    cmd = [
        chrome, "--headless", "--disable-gpu", "--screenshot",
        f"--screenshot={png_path.resolve()}",
        "--window-size=1400,700",
        url
    ]
    subprocess.run(cmd, capture_output=True, timeout=30)
    server.shutdown()
    
    if png_path.exists():
        print(f"      PNG: {png_path.name} ({png_path.stat().st_size//1024} KB)")
        return True
    else:
        print(f"      WARNING: PNG not created: {png_path}")
        return False


# ═════════════════════════════════════════════════════════════════
#  Replace image in slide (remove old <pic> + add new)
# ═════════════════════════════════════════════════════════════════

def replace_image_in_slide(slide, pic_index: int, img_path: Path,
                         left_inch, top_inch, width_inch, height_inch):
    """
    Replace the picture at pic_index with a new image.
    Uses python-pptx's add_picture on the same slide.
    The old picture remains in XML but we add the new one on top.
    """
    from pptx.util import Inches
    # Just add the new picture (old one will be covered or ignored)
    # Actually, let's remove the old pic from XML first
    from lxml import etree
    spTree = slide.shapes._spTree
    pics = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')
    if pic_index < len(pics):
        old_pic = pics[pic_index]
        spTree.remove(old_pic)
        print(f"      Removed old pic element")
    
    # Add new picture
    new_pic = slide.shapes.add_picture(
        str(img_path),
        Inches(left_inch), Inches(top_inch),
        Inches(width_inch), Inches(height_inch)
    )
    return new_pic


def set_shape_text(slide, shape_idx, text_lines):
    """Set text of shape at index (by iterating all text shapes)."""
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if shape_idx < len(text_shapes):
        tf = text_shapes[shape_idx].text_frame
        # Clear
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ""
        # Set first paragraph
        tf.paragraphs[0].runs[0].text = text_lines[0] if text_lines else ""
        # Add more paragraphs if needed
        from pptx.util import Inches, Pt
        for line in text_lines[1:]:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line


def set_title_and_body(slide, title_text, body_lines, is_zh=False):
    """
    Set title + body text for a slide that has 2 text shapes (title, body).
    Assumes shape[0]=title, shape[1]=body (by text shape order).
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    print(f"      Text shapes found: {len(text_shapes)}")
    
    # Title = first text shape
    if len(text_shapes) >= 1:
        tf = text_shapes[0].text_frame
        # Clear all paragraphs
        for para in list(tf.paragraphs):
            for run in list(para.runs):
                run.text = ""
        # Set title
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*TITLE_FG)
        if len(text_shapes) >= 2:
            # Center align title shape
            try:
                from pptx.enum.text import PP_ALIGN
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            except:
                pass
    
    # Body = second text shape
    if len(text_shapes) >= 2:
        tf2 = text_shapes[1].text_frame
        tf2.word_wrap = True
        # Clear
        for para in list(tf2.paragraphs):
            for run in list(para.runs):
                run._r.getparent().remove(run._r)
            para._p.getparent().remove(para._p)
        # Add lines
        first = True
        for line in body_lines:
            if first:
                p = tf2.paragraphs[0]
                first = False
            else:
                p = tf2.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(*BODY_FG)


# ═════════════════════════════════════════════════════════════════
#  MAIN
# ═════════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("Restructure Slides 15-17: 3-Part Migration Narrative (v2)")
    print("=" * 60)
    
    # ── Step 1: Generate SVGs ──
    print("\n[1/4] Generating SVGs...")
    svg15 = make_svg_s15_baseline()
    svg16 = make_svg_s16_anomaly()
    svg17 = make_svg_s17_brain()
    
    svg15_path = ASSETS / "_s15_baseline.svg"
    svg16_path = ASSETS / "_s16_anomaly.svg"
    svg17_path = ASSETS / "_s17_migration.svg"
    svg15_path.write_text(svg15, encoding="utf-8")
    svg16_path.write_text(svg16, encoding="utf-8")
    svg17_path.write_text(svg17, encoding="utf-8")
    print(f"  SVG S15: {len(svg15)//1024} KB")
    print(f"  SVG S16: {len(svg16)//1024} KB")
    print(f"  SVG S17: {len(svg17)//1024} KB")
    
    # ── Step 2: Convert SVG → PNG ──
    print("\n[2/4] Converting SVGs → PNGs (Chrome headless)...")
    png15_path = ASSETS / "_s15_baseline.png"
    png16_path = ASSETS / "_s16_anomaly.png"
    png17_path = ASSETS / "_s17_migration.png"
    
    ok15 = svg_to_png_chrome(svg15_path, png15_path)
    ok16 = svg_to_png_chrome(svg16_path, png16_path)
    ok17 = svg_to_png_chrome(svg17_path, png17_path)
    
    if not (ok15 and ok16 and ok17):
        print("  WARNING: Some PNGs failed. Check Chrome.")
    
    # ── Step 3: Process EN version ──
    print("\n[3/4] Rebuilding slides 15-17 (EN)...")
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation(str(PPT_EN))
    
    # ── Slide 15 (index 14) ──
    print("  Slide 15: Baseline narrative...")
    slide15 = prs.slides[14]
    
    # Set title
    for s in slide15.shapes:
        if s.has_text_frame and "Migration Inference" in s.text_frame.text:
            # This is the title shape
            tf = s.text_frame
            tf.paragraphs[0].runs[0].text = "Baseline: Intra-region ω < Inter-region ω"
            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(*TITLE_FG)
            tf.paragraphs[0].runs[0].font.size = Pt(17)
            tf.paragraphs[0].runs[0].font.bold = True
            break
    
    # Replace image
    for i, s in enumerate(slide15.shapes):
        if hasattr(s, 'image'):
            print(f"    Found image at shape [{i}], replacing...")
            # Remove old pic from XML
            spTree = slide15.shapes._spTree
            pics = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')
            if len(pics) >= 1:
                spTree.remove(pics[0])
                print(f"    Removed old pic from XML")
            break
    
    if png15_path.exists():
        slide15.shapes.add_picture(str(png15_path),
            Inches(5.1), Inches(1.5), Inches(4.5), Inches(3.5))
        print(f"    Inserted new baseline PNG")
    
    # ── Slide 16 (index 15) ──
    print("  Slide 16: Anomaly narrative...")
    slide16 = prs.slides[15]
    
    # Set title
    for s in slide16.shapes:
        if s.has_text_frame:
            txt = s.text_frame.text
            if "Case Studies" in txt or "Anomaly" in txt or "Migration" in txt:
                s.text_frame.paragraphs[0].runs[0].text = "Extreme Case: OPC Breaks the Baseline"
                s.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(*TITLE_FG)
                s.text_frame.paragraphs[0].runs[0].font.size = Pt(17)
                break
    
    # Replace image
    for i, s in enumerate(slide16.shapes):
        if hasattr(s, 'image'):
            print(f"    Found image at shape [{i}], replacing...")
            spTree = slide16.shapes._spTree
            pics = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')
            if len(pics) >= 1:
                spTree.remove(pics[0])
                print(f"    Removed old pic from XML")
            break
    
    if png16_path.exists():
        slide16.shapes.add_picture(str(png16_path),
            Inches(0.6), Inches(1.5), Inches(4.8), Inches(3.5))
        print(f"    Inserted new anomaly PNG")
    
    # ── Slide 17 (index 16) ──
    print("  Slide 17: Migration inference narrative...")
    slide17 = prs.slides[16]
    
    # Set title
    for s in slide17.shapes:
        if s.has_text_frame:
            txt = s.text_frame.text
            if "Specialization" in txt or "Migration" in txt or "Brain" in txt:
                s.text_frame.paragraphs[0].runs[0].text = "Migration Inference: From Anomaly to Hypothesis"
                s.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(*TITLE_FG)
                s.text_frame.paragraphs[0].runs[0].font.size = Pt(17)
                break
    
    # Replace image
    for i, s in enumerate(slide17.shapes):
        if hasattr(s, 'image'):
            print(f"    Found image at shape [{i}], replacing...")
            spTree = slide17.shapes._spTree
            pics = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')
            if len(pics) >= 1:
                spTree.remove(pics[0])
                print(f"    Removed old pic from XML")
            break
    
    if png17_path.exists():
        slide17.shapes.add_picture(str(png17_path),
            Inches(0.5), Inches(1.05), Inches(5.2), Inches(4.0))
        print(f"    Inserted new brain migration PNG")
    
    # Save EN
    prs.save(str(OUT_EN))
    print(f"  Saved EN: {OUT_EN.name}")
    
    # ── Step 4: Process ZH version ──
    print("\n[4/4] Rebuilding slides 15-17 (ZH)...")
    prs_zh = Presentation(str(PPT_ZH))
    
    # Same operations for ZH...
    print("  Slide 15 (ZH): Baseline...")
    slide15_zh = prs_zh.slides[14]
    for s in slide15_zh.shapes:
        if s.has_text_frame:
            txt = s.text_frame.text
            if "迁移" in txt or "Baseline" in txt or "Inference" in txt:
                s.text_frame.paragraphs[0].runs[0].text = "基线：相同脑区内的细胞功能分化 < 不同脑区间的细胞功能分化"
                s.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(*TITLE_FG)
                s.text_frame.paragraphs[0].runs[0].font.size = Pt(17)
                break
    
    for s in slide15_zh.shapes:
        if hasattr(s, 'image'):
            spTree = slide15_zh.shapes._spTree
            pics = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')
            if len(pics) >= 1:
                spTree.remove(pics[0])
            break
    if png15_path.exists():
        slide15_zh.shapes.add_picture(str(png15_path),
            Inches(5.1), Inches(1.5), Inches(4.5), Inches(3.5))
    
    print("  Slide 16 (ZH): Anomaly...")
    slide16_zh = prs_zh.slides[15]
    for s in slide16_zh.shapes:
        if s.has_text_frame:
            txt = s.text_frame.text
            if "Case" in txt or "Anomaly" in txt or "迁移" in txt:
                s.text_frame.paragraphs[0].runs[0].text = "特例：OPC 打破了预期基线"
                s.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(*TITLE_FG)
                s.text_frame.paragraphs[0].runs[0].font.size = Pt(17)
                break
    
    for s in slide16_zh.shapes:
        if hasattr(s, 'image'):
            spTree = slide16_zh.shapes._spTree
            pics = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')
            if len(pics) >= 1:
                spTree.remove(pics[0])
            break
    if png16_path.exists():
        slide16_zh.shapes.add_picture(str(png16_path),
            Inches(0.6), Inches(1.5), Inches(4.8), Inches(3.5))
    
    print("  Slide 17 (ZH): Migration inference...")
    slide17_zh = prs_zh.slides[16]
    for s in slide17_zh.shapes:
        if s.has_text_frame:
            txt = s.text_frame.text
            if "Specialization" in txt or "Migration" in txt or "脑区" in txt:
                s.text_frame.paragraphs[0].runs[0].text = "推断：OPC 迁移假说"
                s.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(*TITLE_FG)
                s.text_frame.paragraphs[0].runs[0].font.size = Pt(17)
                break
    
    for s in slide17_zh.shapes:
        if hasattr(s, 'image'):
            spTree = slide17_zh.shapes._spTree
            pics = spTree.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pic')
            if len(pics) >= 1:
                spTree.remove(pics[0])
            break
    if png17_path.exists():
        slide17_zh.shapes.add_picture(str(png17_path),
            Inches(0.5), Inches(1.05), Inches(5.2), Inches(4.0))
    
    prs_zh.save(str(OUT_ZH))
    print(f"  Saved ZH: {OUT_ZH.name}")
    
    print("\n" + "=" * 60)
    print("DONE. Slides 15-17 restructured with 3-part narrative.")
    print(f"  EN: {OUT_EN}")
    print(f"  ZH: {OUT_ZH}")
    print("=" * 60)


if __name__ == "__main__":
    main()
