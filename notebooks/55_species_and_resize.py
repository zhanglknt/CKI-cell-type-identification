"""
notebooks/55_species_and_resize.py

Three tasks:
  A) Redraw S15-17 SVGs with larger fonts + taller viewBox (400→640)
  B) Adjust S15-17 PPT layout: plot H 2.5→3.2″, cards H 1.50→1.05″
  C) Insert new slide: Generalizing to Any Species (after S7)

Applies to both EN and ZH PPTs.
"""
import sys, http.server, threading, time, os, subprocess, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

BASE   = Path(r"C:\Users\KnightZ\Desktop\细胞受选择")
ASSETS = BASE / "results/figures_final"
ASSETS.mkdir(parents=True, exist_ok=True)

PPT_EN = ASSETS / "CKI_Lecture_2026_v3.pptx"
PPT_ZH = ASSETS / "CKI_Lecture_2026_v3_ZH.pptx"

# ═══════════════════════════════════════════════════════
#  SVG constants — taller viewBox (640 vs 400, 1.6×)
# ═══════════════════════════════════════════════════════
VW, VH = 1700, 640

BG          = "#F8FAFC"
CARD_BG     = "#FFFFFF"
CARD_STROKE = "#E2E8F0"
TEAL        = "#0D9488"
BLUE        = "#2563EB"
RED         = "#DC2626"
PURPLE      = "#7C3AED"
AMBER       = "#F59E0B"
DARK        = "#0F172A"
MUTED       = "#64748B"
GREEN       = "#166534"
RED_DK      = "#991B1B"
WHITE       = "#FFFFFF"
FONT        = 'font-family="Arial,Helvetica,sans-serif"'

# Font sizes scaled 1.6× from original (400→640)
FS_TITLE    = 22   # was 14
FS_SUBTITLE = 16   # was 10
FS_LABEL    = 18   # was 12
FS_VALUE    = 20   # was 13
FS_INSIGHT  = 16   # was 10
FS_BOTTOM   = 19   # was 12
FS_DIVIDER  = 54   # was 36
FS_DIV_LBL  = 14   # was 9
FS_SMALL    = 15   # was 9-10

def _esc(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def _t(x, y, text, size=FS_LABEL, fill=DARK, bold=False, anchor="middle", italic=False):
    fs = f' font-style="italic"' if italic else ""
    fw = ' font-weight="700"' if bold else ""
    return f'<text x="{x}" y="{y}" text-anchor="{anchor}" font-size="{size}"{fw}{fs} fill="{fill}" {FONT}>{_esc(text)}</text>'


# ═══════════════════════════════════════════════════════
#  SVG A — Slide 15: Baseline (Intra vs Inter ω)
# ═══════════════════════════════════════════════════════

def make_s15() -> str:
    """Wide comparison: Intra-region vs Inter-region ω — baseline rule."""
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW} {VH}">']
    out.append(f'<rect width="{VW}" height="{VH}" fill="{BG}" rx="6"/>')

    # ── LEFT: Intra-region ──
    lx, lw = 20, 820
    panel_h = 510
    out.append(f'<rect x="{lx}" y="90" width="{lw}" height="{panel_h}" rx="10" fill="{CARD_BG}" stroke="{CARD_STROKE}" stroke-width="2"/>')
    out.append(_t(lx+lw//2, 128, "Intra-region ω", FS_TITLE, DARK, bold=True))
    out.append(_t(lx+lw//2, 158, "Same brain region · Different cell types", FS_SUBTITLE, MUTED))

    bar_max_w = 380
    bar_h = 52
    bar_x0 = lx + 160
    items_intra = [
        ("Astrocytes", 28, TEAL),
        ("OPCs", 22, BLUE),
        ("Neurons(ex)", 18, PURPLE),
    ]
    for i, (label, omega, color) in enumerate(items_intra):
        y = 195 + i * 85
        bw = int(bar_max_w * omega / 35.0)
        out.append(f'<rect x="{bar_x0}" y="{y}" width="{bw}" height="{bar_h}" rx="5" fill="{color}" opacity="0.85"/>')
        out.append(_t(bar_x0 + bw + 16, y + bar_h//2 + 6, str(omega), FS_VALUE, color, bold=True, anchor="start"))
        out.append(_t(bar_x0 - 16, y + bar_h//2 + 6, label, FS_LABEL, DARK, anchor="end"))

    # Insight strip
    out.append(f'<rect x="{lx+20}" y="505" width="{lw-40}" height="42" rx="8" fill="#F0FDF4" stroke="#86EFAC" stroke-width="1.5"/>')
    out.append(_t(lx+lw//2, 534, "Intra-region ω = LOW — Similar microenvironment", FS_INSIGHT, GREEN, bold=True))

    # ── CENTER: "<" divider ──
    cx = lx + lw + 20
    out.append(f'<line x1="{cx}" y1="95" x2="{cx}" y2="560" stroke="#CBD5E1" stroke-width="2.5" stroke-dasharray="8,6"/>')
    out.append(_t(cx, 305, "<", FS_DIVIDER, "#94A3B8", bold=True))
    out.append(_t(cx, 340, "BASELINE", FS_DIV_LBL, "#94A3B8", bold=True))

    # ── RIGHT: Inter-region ──
    rx = cx + 20
    rw = VW - rx - 20
    out.append(f'<rect x="{rx}" y="90" width="{rw}" height="{panel_h}" rx="10" fill="{CARD_BG}" stroke="{CARD_STROKE}" stroke-width="2"/>')
    out.append(_t(rx+rw//2, 128, "Inter-region ω", FS_TITLE, DARK, bold=True))
    out.append(_t(rx+rw//2, 158, "Same cell type · Different brain regions", FS_SUBTITLE, MUTED))

    items_inter = [
        ("OPC · SVZ", 3.2, TEAL),
        ("OPC · Hipp", 18, BLUE),
        ("OPC · Cortex", 22, PURPLE),
    ]
    bar_x1 = rx + 160
    for i, (label, omega, color) in enumerate(items_inter):
        y = 195 + i * 85
        bw = int(bar_max_w * omega / 35.0)
        out.append(f'<rect x="{bar_x1}" y="{y}" width="{bw}" height="{bar_h}" rx="5" fill="{color}" opacity="0.85"/>')
        out.append(_t(bar_x1 + bw + 16, y + bar_h//2 + 6, str(omega), FS_VALUE, color, bold=True, anchor="start"))
        out.append(_t(bar_x1 - 16, y + bar_h//2 + 6, label, FS_LABEL, DARK, anchor="end"))

    out.append(f'<rect x="{rx+20}" y="505" width="{rw-40}" height="42" rx="8" fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(_t(rx+rw//2, 534, "Inter-region ω = HIGH — Distinct microenvironments", FS_INSIGHT, RED_DK, bold=True))

    # ── Bottom conclusion bar ──
    out.append(f'<rect x="20" y="588" width="{VW-40}" height="44" rx="10" fill="{DARK}"/>')
    out.append(_t(VW//2, 618, "Baseline Rule: Intra-region ω < Inter-region ω  →  Cells differentiate LOCALLY, not by migrating", FS_BOTTOM, WHITE, bold=True))

    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════
#  SVG B — Slide 16: Anomaly (OPC breaks baseline)
# ═══════════════════════════════════════════════════════

def make_s16() -> str:
    """OPC anomaly: Intra-type (tiny) vs Inter-type (huge) contradiction."""
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW} {VH}">']
    out.append(f'<rect width="{VW}" height="{VH}" fill="{BG}" rx="6"/>')

    panel_h = 510
    bar_base_y = 460
    bar_max_scale = 1.4  # scale factor for bar heights in new VH

    # ── LEFT: Intra-type ω (OPC SVZ vs OPC Cortex) ──
    lx, lw = 20, 820
    out.append(f'<rect x="{lx}" y="90" width="{lw}" height="{panel_h}" rx="10" fill="{CARD_BG}" stroke="{CARD_STROKE}" stroke-width="2"/>')
    out.append(_t(lx+lw//2, 125, "Intra-type ω (OPC across regions)", FS_TITLE, DARK, bold=True))
    out.append(_t(lx+lw//2, 155, "OPC in SVZ vs OPC in Cortex — expected LOW divergence", FS_SUBTITLE, MUTED))

    bar_w = 120
    bar_h_max = 320 * bar_max_scale
    # SVZ bar (ω=3.2)
    svz_h = max(int(bar_h_max * 3.2 / 110.0), 15)
    out.append(f'<rect x="{lx+170}" y="{bar_base_y - svz_h}" width="{bar_w}" height="{svz_h}" rx="5" fill="{TEAL}"/>')
    out.append(_t(lx+170+bar_w//2, bar_base_y-svz_h-14, "ω=3.2", FS_VALUE, TEAL, bold=True))
    out.append(_t(lx+170+bar_w//2, bar_base_y+26, "SVZ", FS_LABEL, DARK))

    # Cortex bar (ω=22)
    ctx_h = int(bar_h_max * 22.0 / 110.0)
    out.append(f'<rect x="{lx+400}" y="{bar_base_y - ctx_h}" width="{bar_w}" height="{ctx_h}" rx="5" fill="{BLUE}"/>')
    out.append(_t(lx+400+bar_w//2, bar_base_y-ctx_h-14, "ω=22", FS_VALUE, BLUE, bold=True))
    out.append(_t(lx+400+bar_w//2, bar_base_y+26, "Cortex", FS_LABEL, DARK))

    out.append(_t(lx+650, 300, "Δω = 18.8", FS_LABEL, MUTED))
    out.append(_t(lx+650, 328, "(surprisingly SMALL)", FS_DIV_LBL, MUTED, italic=True))

    out.append(f'<rect x="{lx+20}" y="500" width="{lw-40}" height="42" rx="8" fill="#F0FDF4" stroke="#86EFAC" stroke-width="1.5"/>')
    out.append(_t(lx+lw//2, 529, "Intra-type ω = LOW → OPCs in SVZ and Cortex are HIGHLY similar → likely shared origin", FS_INSIGHT, GREEN, bold=True))

    # ── CENTER divider ──
    cx = lx + lw + 20
    out.append(f'<line x1="{cx}" y1="95" x2="{cx}" y2="560" stroke="#CBD5E1" stroke-width="2.5" stroke-dasharray="8,6"/>')
    out.append(_t(cx, 290, "≪", FS_DIVIDER, "#94A3B8", bold=True))
    out.append(_t(cx, 326, "vs", FS_DIV_LBL, "#94A3B8"))

    # ── RIGHT: Inter-type ω ──
    rx = cx + 20
    rw = VW - rx - 20
    out.append(f'<rect x="{rx}" y="90" width="{rw}" height="{panel_h}" rx="10" fill="{CARD_BG}" stroke="{CARD_STROKE}" stroke-width="2"/>')
    out.append(_t(rx+rw//2, 125, "Inter-type ω (within Cortex)", FS_TITLE, DARK, bold=True))
    out.append(_t(rx+rw//2, 155, "Astrocytes vs OPCs in SAME region — expected HIGH divergence", FS_SUBTITLE, MUTED))

    ast_h = int(bar_h_max * 107.5 / 110.0)
    out.append(f'<rect x="{rx+170}" y="{bar_base_y - ast_h}" width="{bar_w}" height="{ast_h}" rx="5" fill="{RED}"/>')
    out.append(_t(rx+170+bar_w//2, bar_base_y-ast_h-14, "ω=107.5", FS_VALUE, RED, bold=True))
    out.append(_t(rx+170+bar_w//2, bar_base_y+26, "Astrocytes", FS_LABEL, DARK))

    opc_h = int(bar_h_max * 22.0 / 110.0)
    out.append(f'<rect x="{rx+400}" y="{bar_base_y - opc_h}" width="{bar_w}" height="{opc_h}" rx="5" fill="{BLUE}"/>')
    out.append(_t(rx+400+bar_w//2, bar_base_y-opc_h-14, "ω=22", FS_VALUE, BLUE, bold=True))
    out.append(_t(rx+400+bar_w//2, bar_base_y+26, "OPCs", FS_LABEL, DARK))

    out.append(_t(rx+650, 300, "Δω = 85.5", FS_LABEL, MUTED))
    out.append(_t(rx+650, 328, "(extremely LARGE)", FS_DIV_LBL, MUTED, italic=True))

    out.append(f'<rect x="{rx+20}" y="500" width="{rw-40}" height="42" rx="8" fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(_t(rx+rw//2, 529, "Inter-type ω = HIGH → radically different in SAME region → independent lineages", FS_INSIGHT, RED_DK, bold=True))

    # ── CONTRADICTION bar ──
    out.append(f'<rect x="20" y="588" width="{VW-40}" height="44" rx="10" fill="{PURPLE}"/>')
    out.append(_t(VW//2, 618, "CONTRADICTION: Intra-type (22) ≪ Inter-type (107.5) → OPC similarity across regions EXCEEDS cell-type difference within region", FS_BOTTOM - 2, WHITE, bold=True))

    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════
#  SVG C — Slide 17: Brain migration hypothesis
# ═══════════════════════════════════════════════════════

def make_s17() -> str:
    """Brain diagram with regions, ω values, and migration arrow."""
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {VW} {VH}">']
    out.append(f'<rect width="{VW}" height="{VH}" fill="{BG}" rx="6"/>')

    # Brain outline — scaled up for taller canvas
    brain_path = ("M 180 200 "
                  "C 350 100, 700 80, 950 100 "
                  "C 1200 120, 1400 170, 1520 250 "
                  "C 1560 290, 1560 360, 1520 400 "
                  "C 1400 480, 1200 470, 1000 450 "
                  "C 800 430, 500 450, 300 420 "
                  "C 150 400, 120 350, 130 300 "
                  "C 140 230, 150 210, 180 200 Z")
    out.append(f'<path d="{brain_path}" fill="#F1F5F9" stroke="#CBD5E1" stroke-width="4" opacity="0.7"/>')

    # Regions
    # Cortex
    out.append('<ellipse cx="1150" cy="160" rx="200" ry="80" fill="#FEF2F2" stroke="#EF4444" stroke-width="3" opacity="0.55"/>')
    out.append(_t(1150, 138, "Cortex", FS_LABEL, RED_DK, bold=True))
    out.append(_t(1150, 162, "ω=107.5 Astrocytes | ω=22 OPCs", FS_SMALL, RED_DK))

    # SVZ
    out.append('<ellipse cx="870" cy="430" rx="90" ry="55" fill="#F0FDF4" stroke="#16A34A" stroke-width="3" opacity="0.6"/>')
    out.append(_t(870, 415, "SVZ", FS_LABEL, GREEN, bold=True))
    out.append(_t(870, 438, "ω=3.2 OPCs", FS_SMALL, GREEN))

    # Hippocampus
    out.append('<ellipse cx="780" cy="260" rx="80" ry="50" fill="#FEFCE8" stroke="#F59E0B" stroke-width="2.5" opacity="0.5"/>')
    out.append(_t(780, 248, "Hipp", FS_LABEL, "#92400E", bold=True))
    out.append(_t(780, 270, "ω=45", FS_SMALL, "#92400E"))

    # Striatum
    out.append('<ellipse cx="600" cy="300" rx="70" ry="40" fill="#F5F3FF" stroke="#7C3AED" stroke-width="2.5" opacity="0.5"/>')
    out.append(_t(600, 290, "Striatum", FS_LABEL, "#5B21B6", bold=True))
    out.append(_t(600, 312, "ω=38", FS_SMALL, "#5B21B6"))

    # Cerebellum
    out.append('<ellipse cx="320" cy="460" rx="125" ry="70" fill="#E0F2FE" stroke="#0EA5E9" stroke-width="3" opacity="0.55"/>')
    out.append(_t(320, 445, "Cerebellum", FS_LABEL, "#0369A1", bold=True))
    out.append(_t(320, 468, "ω=85 Bergmann glia", FS_SMALL, "#0369A1"))

    # Thalamus
    out.append('<ellipse cx="650" cy="400" rx="50" ry="32" fill="#FFF7ED" stroke="#F97316" stroke-width="2" opacity="0.45"/>')
    out.append(_t(650, 403, "Thal ω=40", FS_SMALL - 2, "#9A3412"))

    # Migration arrow (SVZ → Cortex)
    out.append('<defs>'
               '<filter id="glow"><feGaussianBlur stdDeviation="4" result="blur"/>'
               '<feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge></filter>'
               '<marker id="arrRed" markerWidth="16" markerHeight="12" refX="14" refY="6" orient="auto">'
               '<path d="M0,0 L16,6 L0,12" fill="#DC2626"/></marker>'
               '</defs>')
    out.append(f'<path d="M 890 400 C 970 310, 1050 250, 1120 210" '
               f'stroke="{RED}" stroke-width="5" fill="none" '
               f'stroke-dasharray="14,8" opacity="0.9" filter="url(#glow)" '
               f'marker-end="url(#arrRed)"/>')

    # Migration label pill
    out.append(f'<rect x="930" y="290" width="160" height="40" rx="20" fill="#FEF2F2" stroke="#FCA5A5" stroke-width="2"/>')
    out.append(_t(1010, 310, "OPC Migration", FS_LABEL, RED, bold=True))
    out.append(_t(1010, 326, "ω 3.2 → 22", FS_DIV_LBL, MUTED, italic=True))

    # ── Bottom conclusion ──
    out.append(f'<rect x="20" y="588" width="{VW-40}" height="44" rx="10" fill="{TEAL}"/>')
    out.append(_t(VW//2, 618, "CKI identifies migration via ω deviation: OPC SVZ→Cortex similarity exceeds inter-type divergence within Cortex", FS_BOTTOM - 2, WHITE))

    out.append('</svg>')
    return "\n".join(out)


# ═══════════════════════════════════════════════════════
#  SVG → PNG via Chrome headless
# ═══════════════════════════════════════════════════════

def svg_to_png(svg_path: Path, png_path: Path, view_w=2800, view_h=1050) -> Path:
    """Render SVG to high-res PNG using Chrome headless."""
    chrome_candidates = [
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    chrome = None
    for c in chrome_candidates:
        if Path(c).exists():
            chrome = c
            break
    if not chrome:
        print("  ERROR: No Chrome/Edge found")
        return None

    os.chdir(str(svg_path.parent))
    server = http.server.HTTPServer(("127.0.0.1", 18902), http.server.SimpleHTTPRequestHandler)
    t = threading.Thread(target=server.serve_forever, daemon=True)
    t.start()
    time.sleep(0.5)

    url = f"http://127.0.0.1:18902/{svg_path.name}"
    cmd = [
        chrome, "--headless", "--disable-gpu",
        "--force-device-scale-factor=2",
        f"--screenshot={png_path.resolve()}",
        f"--window-size={view_w},{view_h}",
        url
    ]
    subprocess.run(cmd, capture_output=True, timeout=30)
    server.shutdown()

    if png_path.exists():
        sz = png_path.stat().st_size // 1024
        print(f"    PNG: {png_path.name} ({sz} KB)")
        return png_path
    else:
        print(f"    ERROR: PNG not created")
        return None


# ═══════════════════════════════════════════════════════
#  PPT helpers — layout constants (v3.13)
# ═══════════════════════════════════════════════════════

# New aggressive layout: plot taller, cards shorter
IMG_LEFT, IMG_TOP = 0.5, 0.95
IMG_WIDTH, IMG_HEIGHT = 8.5, 3.2    # was 2.5

CARD_TOP = 4.25                      # was 3.70
CARD_H   = 1.05                      # was 1.50

C3_W, C3_GAP = 2.55, 0.30
C3_X = [IMG_LEFT, IMG_LEFT + C3_W + C3_GAP, IMG_LEFT + 2*(C3_W + C3_GAP)]
C2_W, C2_GAP = 4.10, 0.30
C2_X = [IMG_LEFT, IMG_LEFT + C2_W + C2_GAP]

# Theme
TITLE_FG  = (0x1E, 0x29, 0x3B)
BODY_FG   = (0x1E, 0x29, 0x3B)
CARD_BG_RGB = (0xFF, 0xFF, 0xFF)
C_TEAL    = (0x0D, 0x94, 0x88)
C_BLUE    = (0x25, 0x63, 0xEB)
C_RED     = (0xDC, 0x26, 0x26)
C_PURPLE  = (0x7C, 0x3A, 0xED)
C_DARK_BG = (0x0F, 0x17, 0x2A)
C_MUTED   = (0x64, 0x74, 0x8B)
C_WHITE   = (0xFF, 0xFF, 0xFF)
C_LIGHT   = (0xF1, 0xF5, 0xF9)

PNG_S15 = ASSETS / "_s15_wide.png"
PNG_S16 = ASSETS / "_s16_wide.png"
PNG_S17 = ASSETS / "_s17_wide.png"


def remove_all_pics(slide):
    spTree = slide.shapes._spTree
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    for pic in list(spTree.findall('.//p:pic', ns)):
        spTree.remove(pic)


def remove_non_titles(slide):
    """Keep only title (T<0.95'') and page number (T>4.8'', R>8.5'')."""
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        keep = False
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if shape.top < Inches(0.95):
                keep = True
            elif txt.isdigit() and len(txt) <= 2 and shape.top > Inches(4.8) and (shape.left + shape.width) > Inches(8.5):
                keep = True
        if not keep:
            spTree.remove(shape._element)


def add_card(slide, left, top, width, height, title, lines,
             title_color=None, bg_color=None, font_size=7, title_size=10):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*(bg_color or CARD_BG_RGB))
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*(title_color or C_BLUE))

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(*BODY_FG)


def rebuild_slide(slide, png_path, cards, label):
    remove_non_titles(slide)
    remove_all_pics(slide)
    slide.shapes.add_picture(str(png_path),
        Inches(IMG_LEFT), Inches(IMG_TOP),
        Inches(IMG_WIDTH), Inches(IMG_HEIGHT))
    n_cards = len(cards)
    for i, (title, lines, tc, bg) in enumerate(cards):
        x = C3_X[i] if n_cards == 3 else C2_X[i]
        w = C3_W if n_cards == 3 else C2_W
        fs = 6 if n_cards == 3 else 7
        ts = 9 if n_cards == 3 else 10
        add_card(slide, x, CARD_TOP, w, CARD_H, title, lines, tc, bg, fs, ts)
    print(f"    {label}: plot {IMG_WIDTH}×{IMG_HEIGHT}″ + {n_cards} cards ✓")


# ── Card content (trimmed for shorter cards) ──
#  EN
EN_S15 = [
    ("Intra-region ω", [
        "Same region, different cell types",
        "Microenvironment limits ω divergence",
        "Cortex: Astro(28), OPC(22), Neuron(18)"
    ], C_TEAL, None),
    ("Inter-region ω", [
        "Same type, different brain regions",
        "Distinct environments → higher ω",
        "OPC: SVZ(3.2)→Hipp(18)→Cortex(22)"
    ], C_RED, None),
    ("Baseline Rule", [
        "Intra < Inter across all regions",
        "Cells differentiate LOCALLY, not by migration"
    ], C_DARK_BG, C_LIGHT),
]
EN_S16 = [
    ("OPC: Breaking the Baseline", [
        "23.5% OPC pairs show strong ω anomaly",
        "ω(SVZ)=3.2, ω(Cortex)=22: Δω=18.8 (SMALL)",
        "Intra-type similarity > expected",
        "OPCs known to migrate in development/injury",
        "CKI discovers this without prior knowledge"
    ], C_RED, None),
    ("Astrocytes: High ω, Low Migration", [
        "Highest baseline ω=107.5 (strong identity)",
        "Only 14 migration candidates — very few",
        "Inter-type Δω=85.5 (extremely LARGE)",
        "High ω = local specialization, NOT migration",
        "CKI distinguishes both signals independently"
    ], C_TEAL, None),
]
EN_S17 = [
    ("Omega Gradient Across Brain", [
        "Cortex Astrocytes: ω=107.5 (highest specificity)",
        "Cerebellum Bergmann Glia: ω=85",
        "Hipp/Thal/Striatum: ω=38-45",
        "SVZ OPCs: ω=3.2 (least differentiated)"
    ], C_BLUE, None),
    ("OPC Migration Hypothesis", [
        "Intra-type: OPC(SVZ) vs OPC(Cortex) → ω≈22",
        "Inter-type: Astrocyte vs OPC(Cortex) → ω≈107.5",
        "OPC cross-region similarity > cross-type similarity",
        "→ SVZ → Cortex OPC migration is best hypothesis"
    ], C_RED, None),
]

#  ZH
ZH_S15 = [
    ("同脑区内 ω", [
        "相同脑区，不同细胞类型",
        "微环境限制表达差异 → ω 较低",
        "皮层: Astro(28), OPC(22), Neuron(18)"
    ], C_TEAL, None),
    ("跨脑区间 ω", [
        "相同类型，不同脑区",
        "不同微环境 → ω 更高",
        "OPC: SVZ(3.2)→Hipp(18)→Cortex(22)"
    ], C_RED, None),
    ("基线规律", [
        "同脑区 ω < 跨脑区 ω",
        "细胞在局部环境中分化，非迁移"
    ], C_DARK_BG, C_LIGHT),
]
ZH_S16 = [
    ("OPC：打破基线的特例", [
        "23.5% OPC配对显示强烈ω异常",
        "ω(SVZ)=3.2, ω(Cortex)=22: Δω=18.8",
        "同类间相似性高于预期",
        "OPC已知在发育中活跃迁移",
        "CKI在无先验知识下独立发现"
    ], C_RED, None),
    ("星形胶质细胞：高ω 低迁移", [
        "最高基线 ω=107.5（强细胞身份）",
        "仅14个迁移候选 — 数量极少",
        "跨类Δω=85.5（异常地大）",
        "高ω = 局部特异性分化，非迁移",
        "CKI独立区分两类信号"
    ], C_TEAL, None),
]
ZH_S17 = [
    ("各脑区 ω 梯度分布", [
        "皮层 Astrocytes: ω=107.5（最高）",
        "小脑 Bergmann Glia: ω=85",
        "海马/丘脑/纹状体: ω=38-45",
        "SVZ OPCs: ω=3.2（分化程度最低）"
    ], C_BLUE, None),
    ("OPC 迁移假说", [
        "同类比较：OPC(SVZ) vs OPC(Cortex) → ω≈22",
        "跨类比较：Astrocyte vs OPC（同皮层）→ ω≈107.5",
        "OPC跨脑区相似性 ＞ 皮层内跨类差异",
        "→ SVZ→皮层 OPC迁移是最合理假说"
    ], C_RED, None),
]


# ═══════════════════════════════════════════════════════
#  PART C: Species Generalization Slide
# ═══════════════════════════════════════════════════════

def add_simple_card(slide, x, y, w, h, title, tcolor, lines, btext):
    """3-area card: colored top bar | body lines | bottom note bar."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    CARD_BG_RGB = (0xFF, 0xFF, 0xFF)
    BODY_FG     = (0x1E, 0x29, 0x3B)
    MUTED        = (0x64, 0x74, 0x8B)
    # Card bg
    card = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*CARD_BG_RGB)
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(1.5)
    card.shadow.inherit = False

    # Title bar
    bar_h = 0.44
    bar = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*tcolor)
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_top = Inches(0.03)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Body
    body = slide.shapes.add_textbox(
        Inches(x + 0.08), Inches(y + bar_h + 0.05),
        Inches(w - 0.16), Inches(h - bar_h - 0.48))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(*BODY_FG)

    # Bottom note bar
    note_y = y + h - 0.44
    note = slide.shapes.add_shape(
        1, Inches(x + 0.06), Inches(note_y),
        Inches(w - 0.12), Inches(0.38))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    note.line.fill.background()
    note.shadow.inherit = False

    tf = note.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = btext
    run.font.size = Pt(7)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*MUTED)


def create_species_slide(prs, lang="EN"):
    """Create species generalization slide: 3-column card layout."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # ── Title ──
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.20), Inches(9.0), Inches(0.75))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    if lang == "EN":
        run.text = "Generalizing to Any Species"
    else:
        run.text = "泛化到任意物种"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    if lang == "EN":
        run2.text = "CKI is species-agnostic — all gene set detection is data-driven, no external reference required"
    else:
        run2.text = "CKI 与物种无关 — 基因集鉴定全部基于数据驱动，无需外部参考"
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(*C_MUTED)

    # ── Three-column cards ──
    col_w = 2.80
    col_h = 3.90
    col_gap = 0.25
    col_top = 1.20
    cols_x = [0.50, 0.50 + col_w + col_gap, 0.50 + 2*(col_w + col_gap)]

    if lang == "EN":
        cards = [
            ("Universal Detection", C_BLUE,
             ["Housekeeping: detection_rate (>90%) + CV (lowest 30%)",
              "Functional: HVG top 2,000 (excl. HK genes)",
              "Works for ANY species with RNA-seq data",
              "No external gene lists or curation required"],
             'compute(adata, species="any")  →  fully automatic'),
            ("Human / Mouse Bonus", C_TEAL,
             ["Built-in HRT Atlas (16 tissues, NAR 2021)",
              "Union strategy: detected ∪ reference (use_reference=True)",
              "Human & mouse get reference boost automatically",
              "Other species: purely data-driven, equally valid"],
             'compute(adata, species="human")  →  auto + HRT Atlas'),
            ("Custom Species", C_PURPLE,
             ["Minimal input: species=\"zebrafish\" → generic config",
              "No special setup: species=\"any_name\" → works immediately",
              "Advanced: custom HK reference via reference_path parameter",
              "Fully compatible with any organism"],
             'compute(adata, species="your_species")  →  results'),
        ]
    else:
        cards = [
            ("通用鉴定", C_BLUE,
             ["持家基因：detection_rate (>90%) + CV（最低 30%）",
              "功能基因：HVG top 2,000（排除持家基因）",
              "适用于任何有 RNA-seq 数据的物种",
              "无需外部基因列表或人工筛选"],
             'compute(adata, species="任意物种")  →  全自动完成'),
            ("人/小鼠 加成", C_TEAL,
             ["内置 HRT Atlas 参考集（16 组织，NAR 2021）",
              "并集策略：鉴定结果 ∪ 参考集（use_reference=True）",
              "人、小鼠自动获得参考集加持",
              "其他物种：纯数据驱动，同样有效"],
             'compute(adata, species="human")  →  自动 + HRT Atlas'),
            ("自定义物种", C_PURPLE,
             ["最少输入：species=\"zebrafish\" → 通用配置",
              "无需特殊设置：species=\"任意名称\" → 立即可用",
              "高级：通过 reference_path 提供自定义参考集",
              "完全适配任何生物"],
             'compute(adata, species="你的物种")  →  出结果'),
        ]

    for i, (title, tcolor, lines, btext) in enumerate(cards):
        x = cols_x[i]
        add_simple_card(slide, x, col_top, col_w, col_h,
                         title, tcolor, lines, btext)

    # ── Bottom bar ──
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(5.25), Inches(9.0), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*C_DARK_BG)
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    if lang == "EN":
        run.text = 'compute(adata, species="any_species", groupby="cell_type", group_a="X", group_b="Y")  →  fully automatic CKI analysis'
    else:
        run.text = 'compute(adata, species="任意物种", groupby="cell_type", group_a="X", group_b="Y")  →  全自动 CKI 分析'
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*C_WHITE)

    return slide# ═══════════════════════════════════════════════════════
#  Slide insertion helper
# ═══════════════════════════════════════════════════════

def move_slide_to(prs, from_idx, to_idx):
    """Move slide from from_idx to to_idx in the sldIdLst."""
    pres_elem = prs.part._element
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = pres_elem.find('.//p:sldIdLst', ns)
    sld_ids = list(sldIdLst.findall('p:sldId', ns))
    elem = sld_ids.pop(from_idx)
    sldIdLst.insert(to_idx, elem)
    print(f"    Moved slide {from_idx+1} → position {to_idx+1}")


# ═══════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("55_species_and_resize.py")
    print(f"  SVG viewBox: {VW}×{VH} ({VW/VH:.2f}:1)")
    print(f"  PPT plot: {IMG_WIDTH}×{IMG_HEIGHT}″ (was 2.5″)")
    print(f"  PPT cards: H={CARD_H}″ at T={CARD_TOP}″ (was H=1.50″, T=3.70″)")
    print("=" * 60)

    # ── Step 1: Generate SVGs ──
    print("\n[1/6] Generating taller SVGs (1700×640, 1.6× font scale)...")
    svg_map = {
        "s15": (ASSETS / "_s15_baseline_wide.svg", make_s15()),
        "s16": (ASSETS / "_s16_anomaly_wide.svg", make_s16()),
        "s17": (ASSETS / "_s17_migration_wide.svg", make_s17()),
    }
    for key, (path, svg) in svg_map.items():
        path.write_text(svg, encoding="utf-8")
        print(f"  {key}: {path.name} ({len(svg)//1024} KB)")

    # ── Step 2: SVG → PNG ──
    print("\n[2/6] Rendering high-res PNGs (Chrome 2800×1050, 2x scale)...")
    png_files = {}
    for key, (svg_path, _) in svg_map.items():
        png_path = ASSETS / f"_{key}_wide.png"
        result = svg_to_png(svg_path, png_path, view_w=2800, view_h=1050)
        if result:
            png_files[key] = result

    if len(png_files) < 3:
        print("  ERROR: Not all PNGs rendered.")
        return

    # ── Step 3: Update S15-17 in EN PPT ──
    print("\n[3/6] Rebuilding S15-17 in EN PPT...")
    prs_en = Presentation(str(PPT_EN))
    # S15-17 are at indices 15, 16, 17 after S7 insertion
    rebuild_slide(prs_en.slides[15], png_files["s15"], EN_S15, "S15 EN")
    rebuild_slide(prs_en.slides[16], png_files["s16"], EN_S16, "S16 EN")
    rebuild_slide(prs_en.slides[17], png_files["s17"], EN_S17, "S17 EN")
    prs_en.save(str(PPT_EN))
    print(f"  Saved: {PPT_EN.name}")

    # ── Step 4: Update S15-17 in ZH PPT ──
    print("\n[4/6] Rebuilding S15-17 in ZH PPT...")
    prs_zh = Presentation(str(PPT_ZH))
    rebuild_slide(prs_zh.slides[15], png_files["s15"], ZH_S15, "S15 ZH")
    rebuild_slide(prs_zh.slides[16], png_files["s16"], ZH_S16, "S16 ZH")
    rebuild_slide(prs_zh.slides[17], png_files["s17"], ZH_S17, "S17 ZH")
    prs_zh.save(str(PPT_ZH))
    print(f"  Saved: {PPT_ZH.name}")

    # ── Step 5: Create & insert Species slide ──
    print("\n[5/6] Creating species generalization slide (after S7)...")
    for ppt_path, lang in [(PPT_EN, "EN"), (PPT_ZH, "ZH")]:
        prs = Presentation(str(ppt_path))
        create_species_slide(prs, lang)
        n_slides = len(prs.slides)
        move_slide_to(prs, n_slides - 1, 7)  # after S7 (index 6), new S8 (index 7)
        prs.save(str(ppt_path))
        print(f"  {lang}: species slide inserted at position 8")

    # ── Step 6: Verify ──
    print(f"\n[6/6] VERIFICATION — EN Slide order:")
    prs = Presentation(str(PPT_EN))
    for i, slide in enumerate(prs.slides):
        titles = [s for s in slide.shapes if hasattr(s, 'text_frame') and s.top < Inches(1.1)]
        t = titles[0].text_frame.text.strip()[:75] if titles else '(no title)'
        has_pic = any(hasattr(s, 'image') for s in slide.shapes)
        note = " [PIC]" if has_pic else ""
        print(f"  S{i+1:2d}  {t}{note}")

    print(f"\n{'='*60}")
    print("ALL DONE — v3.13")
    print(f"  SVGs: 1700×640, fonts 1.6× larger")
    print(f"  Plot: 2.5″→3.2″  |  Cards: 1.50″→1.05″, 6-7pt")
    print(f"  New S8: Generalizing to Any Species (EN+ZH)")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
