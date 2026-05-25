"""
notebooks/55_species_and_resize.py  (PART C only)
Species generalization slide.
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Reuse theme from main script
TITLE_FG  = (0x1E, 0x29, 0x3B)
BODY_FG   = (0x1E, 0x29, 0x3B)
CARD_BG   = (0xFF, 0xFF, 0xFF)
BLUE      = (0x25, 0x63, 0xEB)
TEAL      = (0x0D, 0x94, 0x88)
RED       = (0xDC, 0x26, 0x26)
PURPLE    = (0x7C, 0x3A, 0xED)
AMBER     = (0xD9, 0x77, 0x06)
DARK_BG   = (0x0F, 0x17, 0x2A)
MUTED     = (0x64, 0x74, 0x8B)
WHITE     = (0xFF, 0xFF, 0xFF)
LIGHT     = (0xF1, 0xF5, 0xF9)


def add_simple_card(slide, x, y, w, h, title, tcolor, lines, btext):
    """3-area card: colored top bar | body lines | bottom note bar."""
    # Card bg
    card = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*CARD_BG)
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(1.5)
    card.shadow.inherit = False

    # Title bar
    bar_h = 0.44
    bar = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*tcolor)
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_top = Inches(0.03)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*WHITE)

    # Body
    body = slide.shapes.add_textbox(
        Inches(x + 0.08), Inches(y + bar_h + 0.05),
        Inches(w - 0.16), Inches(h - bar_h - 0.48))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(*BODY_FG)

    # Bottom note bar
    note_y = y + h - 0.44
    note = slide.shapes.add_shape(
        1, Inches(x + 0.06), Inches(note_y),
        Inches(w - 0.12), Inches(0.38))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    note.line.fill.background()
    note.shadow.inherit = False

    tf = note.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = btext
    run.font.size = Pt(7)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*MUTED)


def create_species_slide(prs, lang="EN"):
    """Create species generalization slide with 3-column card layout."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.18), Inches(9.0), Inches(0.72))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    if lang == "EN":
        run.text = "Generalizing to Any Species"
    else:
        run.text = "泛化到任意物种"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    if lang == "EN":
        run2.text = ("CKI is species-agnostic — all gene set detection is "
                     "data-driven, no external reference required")
    else:
        run2.text = "CKI 与物种无关 — 基因集鉴定全部基于数据驱动，无需外部参考"
    run2.font.size = Pt(10)
    run2.font.color.rgb = RGBColor(*MUTED)

    # 3 columns
    col_w   = 2.80
    col_h   = 3.90
    col_gap = 0.25
    col_y    = 1.12
    cols_x  = [0.50, 0.50 + col_w + col_gap, 0.50 + 2 * (col_w + col_gap)]

    if lang == "EN":
        cards = [
            ("Universal Detection", BLUE,
             ["Housekeeping: detection_rate (>90%) + CV (lowest 30%)",
              "Functional: HVG top 2,000 (excl. HK genes)",
              "Works for ANY species with RNA-seq data",
              "No external gene lists or curation required"],
             "compute(adata, species=\"any\")  →  fully automatic"),
            ("Human / Mouse Bonus", TEAL,
             ["Built-in HRT Atlas (16 tissues, NAR 2021)",
              "Union strategy: detected ∪ reference (use_reference=True)",
              "Human & mouse get reference boost automatically",
              "Other species: purely data-driven, equally valid"],
             "compute(adata, species=\"human\")  →  auto + HRT Atlas"),
            ("Custom Species", PURPLE,
             ["Minimal input: species=\"zebrafish\" → generic config",
              "No special setup: species=\"any_name\" → works immediately",
              "Advanced: custom HK reference via reference_path parameter",
              "Fully compatible with any organism"],
             "compute(adata, species=\"your_species\")  →  results"),
        ]
    else:
        cards = [
            ("通用鉴定", BLUE,
             ["持家基因：detection_rate (>90%) + CV（最低 30%）",
              "功能基因：HVG top 2,000（排除持家基因）",
              "适用于任何有 RNA-seq 数据的物种",
              "无需外部基因列表或人工筛选"],
             "compute(adata, species=\"任意物种\")  →  全自动完成"),
            ("人/小鼠 加成", TEAL,
             ["内置 HRT Atlas 参考集（16 组织，NAR 2021）",
              "并集策略：鉴定结果 ∪ 参考集（use_reference=True）",
              "人、小鼠自动获得参考集加持",
              "其他物种：纯数据驱动，同样有效"],
             "compute(adata, species=\"human\")  →  自动 + HRT Atlas"),
            ("自定义物种", PURPLE,
             ["最少输入：species=\"zebrafish\" → 通用配置",
              "无需特殊设置：species=\"任意名称\" → 立即可用",
              "高级：通过 reference_path 提供自定义参考集",
              "完全适配任何生物"],
             "compute(adata, species=\"你的物种\")  →  出结果"),
        ]

    for i, (title, tcolor, lines, btext) in enumerate(cards):
        add_simple_card(slide, cols_x[i], col_y, col_w, col_h,
                         title, tcolor, lines, btext)

    # Bottom bar
    bar = slide.shapes.add_shape(
        1, Inches(0.5), Inches(5.25), Inches(9.0), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*DARK_BG)
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    if lang == "EN":
        run.text = ('compute(adata, species="any_species", '
                     'groupby="cell_type", group_a="X", group_b="Y")  '
                     '→  fully automatic CKI analysis')
    else:
        run.text = ('compute(adata, species="任意物种", '
                     'groupby="cell_type", group_a="X", group_b="Y")  '
                     '→  全自动 CKI 分析')
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*WHITE)

    return slide
