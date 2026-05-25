"""
59_comprehensive_fix.py — 一次性根治 PPT 所有系统性缺陷

修复内容：
1. S3 重复图片清理
2. S7 底部栏溢出修复
3. S7-S21 页码全面修正
4. S16/S18/S19 SVG 中 107.5→121.8 数据更正 + PNG 重渲染
5. S18/S19 文本框数据更新
6. EN+ZH 两版同步
7. 最终审计
"""

import os, sys, subprocess, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

BASE = Path("results/figures_final")
EN_PPT = BASE / "CKI_Lecture_2026_v3.pptx"
ZH_PPT = BASE / "CKI_Lecture_2026_v3_ZH.pptx"

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

# ============================================================
# PHASE 1: Fix SVGs — replace 107.5 with 121.8
# ============================================================
print("=" * 60)
print("PHASE 1: Fix SVG data (107.5 → 121.8)")
print("=" * 60)

svg_fixes = {
    "_s16_anomaly_v3.svg": [
        ("ω (107.5)", "ω (121.8)"),
        ("ω=107.5", "ω=121.8"),
        ("ω = 107.5", "ω = 121.8"),
        ("107.5 (HIGH)", "121.8 (HIGH)"),
        ("107.5 Astro", "121.8 Astro"),
        (" 107.5 →", " 121.8 →"),
    ],
    "_s17_migration_v3.svg": [
        ("ω=107.5 Astro", "ω=121.8 Astro"),
        ("107.5 (highest", "121.8 (highest"),
    ],
}

for svg_name, replacements in svg_fixes.items():
    svg_path = BASE / svg_name
    content = svg_path.read_text(encoding="utf-8")
    for old, new in replacements:
        if old in content:
            content = content.replace(old, new)
            print(f"  {svg_name}: '{old}' → '{new}'")
        else:
            print(f"  {svg_name}: '{old}' NOT FOUND (skip)")
    svg_path.write_text(content, encoding="utf-8")

# Also fix the anomaly bar height if needed
svg_path = BASE / "_s16_anomaly_v3.svg"
content = svg_path.read_text(encoding="utf-8")
# Bar for Astrocyte: height=195 represented ω=107.5, now ω=121.8
# Scale = 195/107.5 * 121.8 = 220.9, round to 221
# y position was 142, new y = 142 + 195 - 221 = 116
old_bar = 'x="680" y="142" width="80" height="195"'
new_bar = 'x="680" y="116" width="80" height="221"'
if old_bar in content:
    content = content.replace(old_bar, new_bar)
    print(f"  _s16_anomaly_v3.svg: bar height 195→221 (y 142→116)")
svg_path.write_text(content, encoding="utf-8")

# Re-render PNGs via Chrome headless
print("\nRe-rendering PNGs...")
chrome = None
for candidate in [
    "C:/Program Files/Google/Chrome/Application/chrome.exe",
    "C:/Program Files (x86)/Google/Chrome/Application/chrome.exe",
    "C:/Program Files/Microsoft/Edge/Application/msedge.exe",
    "C:/Program Files (x86)/Microsoft/Edge/Application/msedge.exe",
]:
    if os.path.exists(candidate):
        chrome = candidate
        break

if chrome:
    png_map = {
        "_s16_anomaly_v3.svg": "_s16_anomaly_v3.png",
        "_s17_migration_v3.svg": "_s17_migration_v3.png",
    }
    for svg_name, png_name in png_map.items():
        svg_path = BASE / svg_name
        html_path = BASE / "_tmp_render.html"
        html = f"""<!DOCTYPE html>
<html><head><style>body{{margin:0;padding:0;overflow:hidden}}</style></head>
<body><object data="{svg_name}" type="image/svg+xml" width="1200" height="480"></object></body></html>"""
        html_path.write_text(html, encoding="utf-8")

        out_png = f"C:/Users/KnightZ/Desktop/_tmp_{png_name}"
        result = subprocess.run(
            [chrome, "--headless=new", "--disable-gpu",
             f"--screenshot={out_png}",
             "--window-size=1220,520",
             str(html_path)],
            capture_output=True, timeout=30, cwd=str(BASE)
        )
        if Path(out_png).exists():
            target = BASE / png_name
            shutil.copy2(out_png, str(target))
            Path(out_png).unlink()
            print(f"  {png_name}: {target.stat().st_size//1024}KB")
        else:
            print(f"  {png_name}: FAILED — {result.stderr.decode()[:200]}")
        html_path.unlink(missing_ok=True)
else:
    print("  WARNING: no Chrome/Edge found, PNGs not re-rendered")

# ============================================================
# PHASE 2: Fix PPTs
# ============================================================
print("\n" + "=" * 60)
print("PHASE 2: Fix PPT slides")
print("=" * 60)

for tag, pptx_path in [("EN", EN_PPT), ("ZH", ZH_PPT)]:
    prs = Presentation(str(pptx_path))
    print(f"\n--- {tag} ---")

    # ----------------------------------------------------------
    # 2a. Fix S3 duplicate images (12 PIC → 4 PIC)
    # ----------------------------------------------------------
    slide = prs.slides[2]  # S3 (0-indexed)
    spTree = slide._element.find(".//p:spTree", NS)
    if spTree is not None:
        pics = spTree.findall("p:pic", NS)
        print(f"  S3: {len(pics)} pic elements found")
        # Keep first 4 unique pics, remove the rest
        # Each set of 4 is at same positions, keep first set
        if len(pics) > 4:
            for pic in pics[4:]:
                spTree.remove(pic)
            print(f"  S3: removed {len(pics)-4} duplicate pics")

    # ----------------------------------------------------------
    # 2b. Fix S7 bottom bar overflow
    # ----------------------------------------------------------
    slide = prs.slides[6]  # S7 (0-indexed)
    for s in slide.shapes:
        if hasattr(s, "text_frame"):
            txt = s.text_frame.text.strip()
            if "Input (adata + species)" in txt and "Auto" in txt:
                # Expand height to fit text
                s.height = Inches(0.50)
                s.top = Inches(6.55)
                print(f"  S7: expanded bottom bar H→0.50\"")
                break

    # ----------------------------------------------------------
    # 2c. Fix page numbers
    # ----------------------------------------------------------
    # Slide layout (0-indexed):
    #   S1  = Title (no #)
    #   S2  = page "2"  ✓
    #   S3  = page "3"  ✓
    #   S4  = page "4"  ✓
    #   S5  = page "5"  ✓
    #   S6  = page "6"  ✓
    #   S7  = [NEW] Automated Gene Set Detection — no #
    #   S8  = [NEW] Generalizing to Any Species — no #
    #   S9  = was S7, shows "7" → should be "9"
    #   S10 = was S8, shows "8" → should be "10"
    #   ...
    #   S21 = was S19, shows "19" → should be "21"
    #   S22 = Thank You (no #)
    #
    # Also fix brain slides (S16-S19):
    #   S16 = brain regional, shows "16" ✓ (was 14, +2 = 16)
    #   S17 = baseline, shows "15" → should be "17"
    #   S18 = extreme OPC, shows "16" → should be "18"
    #   S19 = migration, shows "17" → should be "19"

    # Add page numbers to S7 and S8
    for slide_idx, page_num in [(6, "7"), (7, "8")]:
        slide = prs.slides[slide_idx]
        page_box = slide.shapes.add_textbox(
            Inches(9.2), Inches(7.05), Inches(0.5), Inches(0.25))
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = page_num
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        print(f"  S{slide_idx+1}: added page number '{page_num}'")

    # Fix page numbers on S9-S21 (0-indexed 8-20)
    # Old page numbers: 7,8,9,10,11,12,13,14,15,16,17,18,19
    # New page numbers: 9,10,11,12,13,14,15,16,17,18,19,20,21
    # But brain slides S16-S19 (idx 15-18) have wrong existing numbers:
    #   S16(idx15) shows "16" — already correct, skip
    #   S17(idx16) shows "15" → should be "17"
    #   S18(idx17) shows "16" → should be "18"
    #   S19(idx18) shows "17" → should be "19"

    page_map = {
        8:  "9",    # S9  (was S7, shows "7")
        9:  "10",   # S10 (was S8, shows "8")
        10: "11",   # S11 (was S9, shows "9")
        11: "12",   # S12 (was S10, shows "10")
        12: "13",   # S13 (was S11, shows "11")
        13: "14",   # S14 (was S12, shows "12")
        14: "15",   # S15 (was S13, shows "13")
        15: "16",   # S16 (was S14, shows "16") — already correct
        16: "17",   # S17 (was S15, shows "15")
        17: "18",   # S18 (was S16, shows "16")
        18: "19",   # S19 (was S17, shows "17")
        19: "20",   # S20 (was S18, shows "18")
        20: "21",   # S21 (was S19, shows "19")
    }

    for slide_idx, new_num in page_map.items():
        slide = prs.slides[slide_idx]
        # Find and update page number text
        for s in slide.shapes:
            if hasattr(s, "text_frame"):
                txt = s.text_frame.text.strip()
                # Match shapes that are just a number (page numbers)
                if txt in {"7","8","9","10","11","12","13","14","15","16","17","18","19","20","21"}:
                    # Verify it's actually a page number (small font, bottom-right)
                    try:
                        if s.left/914400 > 8.5 and s.top/914400 > 5.0:
                            # Update text
                            for p in s.text_frame.paragraphs:
                                for r in p.runs:
                                    r.text = new_num
                            print(f"  S{slide_idx+1}: page {txt} → {new_num}")
                            break
                    except:
                        pass

    # ----------------------------------------------------------
    # 2d. Fix S18/S19 text box data
    # ----------------------------------------------------------
    for slide_idx in [17, 18]:  # S18, S19 (0-indexed)
        slide = prs.slides[slide_idx]
        for s in slide.shapes:
            if hasattr(s, "text_frame"):
                txt = s.text_frame.text
                if "107.5" in txt:
                    new_txt = txt.replace("107.5", "121.8")
                    # Update all runs
                    for p in s.text_frame.paragraphs:
                        for r in p.runs:
                            if "107.5" in (r.text or ""):
                                r.text = r.text.replace("107.5", "121.8")
                    print(f"  S{slide_idx+1}: replaced 107.5→121.8 in text box")

                # Fix Δω calculation: 85.5 → 105.8
                if "Δω=85.5" in txt:
                    for p in s.text_frame.paragraphs:
                        for r in p.runs:
                            if "Δω=85.5" in (r.text or ""):
                                r.text = r.text.replace("Δω=85.5", "Δω=105.8")
                    print(f"  S{slide_idx+1}: replaced Δω=85.5→105.8")

    # ----------------------------------------------------------
    # 2e. Replace S18/S19 PNG images with corrected versions
    # ----------------------------------------------------------
    png_map_slides = {
        17: "_s16_anomaly_v3.png",   # S18 (0-indexed 17)
        18: "_s17_migration_v3.png",  # S19 (0-indexed 18)
    }

    for slide_idx, png_name in png_map_slides.items():
        png_path = BASE / png_name
        if not png_path.exists():
            print(f"  S{slide_idx+1}: {png_name} not found, skip image replace")
            continue

        slide = prs.slides[slide_idx]
        # Find existing picture
        for s in slide.shapes:
            if hasattr(s, "image"):
                old_left = s.left
                old_top = s.top
                old_width = s.width
                old_height = s.height
                # Remove old pic
                s._element.getparent().remove(s._element)
                # Add new pic
                slide.shapes.add_picture(
                    str(png_path), old_left, old_top, old_width, old_height)
                print(f"  S{slide_idx+1}: replaced PIC with corrected {png_name}")
                break

    # Save
    prs.save(str(pptx_path))
    print(f"  {tag}: saved")

# ============================================================
# PHASE 3: Final audit
# ============================================================
print("\n" + "=" * 60)
print("PHASE 3: Final audit")
print("=" * 60)

for tag, pptx_path in [("EN", EN_PPT), ("ZH", ZH_PPT)]:
    prs = Presentation(str(pptx_path))
    issues = []

    for idx, slide in enumerate(prs.slides):
        # Check for duplicate images
        pics = slide._element.findall(".//p:pic", NS)
        if idx == 2:  # S3
            if len(pics) > 4:
                issues.append(f"S{idx+1}: {len(pics)} pic elements (should be 4)")

        # Check page numbers
        page_num = None
        for s in slide.shapes:
            if hasattr(s, "text_frame"):
                txt = s.text_frame.text.strip()
                if txt.isdigit() and s.left/914400 > 8.5 and s.top/914400 > 5.0:
                    page_num = int(txt)

        expected_nums = {
            0: None, 1: 2, 2: 3, 3: 4, 4: 5, 5: 6, 6: 7, 7: 8,
            8: 9, 9: 10, 10: 11, 11: 12, 12: 13, 13: 14, 14: 15,
            15: 16, 16: 17, 17: 18, 18: 19, 19: 20, 20: 21, 21: None
        }
        expected = expected_nums.get(idx)
        if expected is not None and page_num != expected:
            issues.append(f"S{idx+1}: page #{page_num}, expected {expected}")

        # Check for 107.5 in S18/S19
        if idx in [17, 18]:
            for s in slide.shapes:
                if hasattr(s, "text_frame") and "107.5" in s.text_frame.text:
                    issues.append(f"S{idx+1}: still contains '107.5'")

        # Check overflow
        for s in slide.shapes:
            t = s.top / 914400
            h = s.height / 914400
            if t + h > 7.6:
                issues.append(f"S{idx+1}: overflow bottom T={t:.1f}+H={h:.1f}")

    if issues:
        print(f"  {tag}: ⚠ {len(issues)} issues")
        for iss in issues:
            print(f"    {iss}")
    else:
        print(f"  {tag}: ✓ CLEAN")

print("\nDone.")
