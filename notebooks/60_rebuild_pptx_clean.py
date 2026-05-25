"""
60_rebuild_pptx_clean.py — 从损坏的PPTX重建干净版本
使用python-pptx API完整重建，确保PowerPoint兼容性。
处理EN和ZH两个版本。
"""

import os, copy, io
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pathlib import Path

BASE = Path("results/figures_final")
EN_SRC = BASE / "CKI_Lecture_2026_v3.pptx"
ZH_SRC = BASE / "CKI_Lecture_2026_v3_ZH.pptx"
EN_OUT = BASE / "CKI_Lecture_2026_v4.pptx"
ZH_OUT = BASE / "CKI_Lecture_2026_v4_ZH.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def clean_slide_xml(sp_tree_xml):
    """Clean up potentially problematic XML attributes in spTree elements."""
    # Remove empty name="" attributes from cNvPr elements
    for elem in sp_tree_xml.iter():
        # Clear empty name attributes that might cause issues
        if elem.get("name") == "":
            elem.attrib.pop("name", None)
        # Clear empty descr attributes
        if elem.get("descr") == "":
            elem.attrib.pop("descr", None)
    return sp_tree_xml


def rebuild_pptx(src_path, dst_path):
    """Rebuild a PPTX from source, creating clean XML."""
    src = Presentation(str(src_path))
    
    # Read source as zip for images
    with zipfile.ZipFile(str(src_path), 'r') as z:
        src_names = set(z.namelist())
    
    # Create new blank presentation
    dst = Presentation()
    dst.slide_width = src.slide_width
    dst.slide_height = src.slide_height
    
    # Remove default blank slide (clean up completely)
    sldIdLst = dst.part._element.find(qn('p:sldIdLst'))
    while sldIdLst is not None and len(sldIdLst) > 0:
        for child in list(sldIdLst):
            rid = child.get(qn('r:id'))
            if rid:
                try:
                    dst.part.drop_rel(rid)
                except Exception:
                    pass
            sldIdLst.remove(child)
        # Re-fetch in case the element was replaced
        sldIdLst = dst.part._element.find(qn('p:sldIdLst'))
    
    # Also clean up any orphan slide relationships in presentation.xml.rels
    pres_rels = dst.part.rels
    orphan_rids = []
    for rid, rel in pres_rels.items():
        if 'slide' in rel.reltype and 'slideLayout' not in rel.reltype and 'slideMaster' not in rel.reltype:
            orphan_rids.append(rid)
    for rid in orphan_rids:
        try:
            dst.part.drop_rel(rid)
        except Exception:
            pass
    
    import tempfile
    
    for slide_idx, src_slide in enumerate(src.slides):
        # Get source slide's spTree
        src_spTree = src_slide._element.find(f'{{{P_NS}}}cSld/{{{P_NS}}}spTree')
        if src_spTree is None:
            print(f"  WARNING: Slide {slide_idx+1} has no spTree, skipping")
            continue
        
        # Collect ALL image blobs referenced by this slide
        src_rels = src_slide.part.rels
        img_blobs = {}  # {old_rid: (blob, content_type)}
        for rid, rel in src_rels.items():
            if 'image' in rel.reltype:
                try:
                    img_blobs[rid] = (rel.target_part.blob, rel.target_part.content_type)
                except Exception:
                    pass
        
        # Deep copy spTree and clean
        new_spTree = copy.deepcopy(src_spTree)
        new_spTree = clean_slide_xml(new_spTree)
        
        # Add blank slide with layout
        layout = dst.slide_layouts[0]
        dst_slide = dst.slides.add_slide(layout)
        
        # Register images with the new slide and get new rIds
        rid_map = {}  # {old_rid: new_rid}
        for old_rid, (blob, content_type) in img_blobs.items():
            # Determine file extension from content type
            ext = 'png'
            if 'jpeg' in content_type or 'jpg' in content_type:
                ext = 'jpg'
            elif 'gif' in content_type:
                ext = 'gif'
            elif 'bmp' in content_type:
                ext = 'bmp'
            elif 'tiff' in content_type:
                ext = 'tiff'
            
            # Write blob to temp file
            with tempfile.NamedTemporaryFile(
                suffix=f'.{ext}', delete=False
            ) as tmp:
                tmp.write(blob)
                tmp_path = tmp.name
            
            try:
                image_part, new_rid = dst_slide.part.get_or_add_image_part(tmp_path)
                rid_map[old_rid] = new_rid
            except Exception as e:
                print(f"    WARNING: Slide {slide_idx+1}: image {old_rid}: {e}")
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass
        
        # Update blip r:embed references in the new spTree
        for blip in new_spTree.iter(f'{{{A_NS}}}blip'):
            old_rid = blip.get(qn('r:embed'))
            if old_rid and old_rid in rid_map:
                blip.set(qn('r:embed'), rid_map[old_rid])
        
        # Replace the empty spTree with our cleaned one
        dst_spTree = dst_slide._element.find(f'{{{P_NS}}}cSld/{{{P_NS}}}spTree')
        cSld = dst_slide._element.find(f'{{{P_NS}}}cSld')
        
        if dst_spTree is not None:
            cSld.remove(dst_spTree)
        
        # Copy slide background if present
        src_bg = src_slide._element.find(f'{{{P_NS}}}cSld/{{{P_NS}}}bg')
        if src_bg is not None:
            new_bg = copy.deepcopy(src_bg)
            cSld.insert(0, new_bg)
        
        # Add cleaned spTree with updated image references
        cSld.append(new_spTree)
        
        # Remove notes slide relationship if any
        slide_rels = dst_slide.part.rels
        notes_to_drop = []
        for rid, rel in slide_rels.items():
            if 'notesSlide' in rel.reltype:
                notes_to_drop.append(rid)
        for rid in notes_to_drop:
            dst_slide.part.drop_rel(rid)
    
    # Final cleanup: remove any orphan slide relationships
    pres_sldIdLst = dst.part._element.find(qn('p:sldIdLst'))
    if pres_sldIdLst is not None:
        valid_rids = set()
        for child in pres_sldIdLst:
            rid = child.get(qn('r:id'))
            if rid:
                valid_rids.add(rid)
        
        pres_rels = dst.part.rels
        orphan_rids = []
        for rid, rel in pres_rels.items():
            if 'slide' in rel.reltype and 'slideLayout' not in rel.reltype and 'slideMaster' not in rel.reltype:
                if rid not in valid_rids:
                    orphan_rids.append(rid)
        for rid in orphan_rids:
            try:
                dst.part.drop_rel(rid)
            except Exception:
                pass
    
    # Save
    dst.save(str(dst_path))
    return len(dst.slides)


def verify_pptx(path):
    """Verify the rebuilt PPTX."""
    print(f"\n  Verifying {path.name}...")
    
    # ZIP check
    with zipfile.ZipFile(str(path), 'r') as z:
        bad = z.testzip()
        print(f"    ZIP: {'OK' if not bad else f'CORRUPT at {bad}'}")
        
        # XML parse check
        xml_ok = True
        for name in z.namelist():
            if name.endswith('.xml') or name.endswith('.rels'):
                try:
                    etree.fromstring(z.read(name))
                except Exception as e:
                    print(f"    XML ERROR: {name}: {e}")
                    xml_ok = False
        if xml_ok:
            print(f"    All XML: OK")
        
        # No notes slides
        notes = [n for n in z.namelist() if 'notesSlide' in n or 'notesMaster' in n]
        print(f"    Notes slides: {len(notes)} (should be 0)")
    
    # python-pptx check
    try:
        prs = Presentation(str(path))
        print(f"    python-pptx: {len(prs.slides)} slides OK")
        
        # Check no slide has notes
        has_notes = False
        for slide in prs.slides:
            if slide.has_notes_slide:
                has_notes = True
        print(f"    Slides with notes: {'Yes' if has_notes else 'None (good)'}")
        
    except Exception as e:
        print(f"    python-pptx ERROR: {e}")
    
    print(f"    File size: {os.path.getsize(str(path)):,} bytes")


# ============================================================
# MAIN
# ============================================================
if __name__ == "__main__":
    for tag, src, dst in [
        ("EN", EN_SRC, EN_OUT),
        ("ZH", ZH_SRC, ZH_OUT),
    ]:
        print(f"\n{'='*60}")
        print(f"Rebuilding {tag} version...")
        print(f"  Source: {src.name}")
        print(f"  Output: {dst.name}")
        
        n = rebuild_pptx(src, dst)
        print(f"  {n} slides rebuilt")
        
        verify_pptx(dst)
    
    print(f"\n{'='*60}")
    print("Done! New files:")
    print(f"  EN: {EN_OUT}")
    print(f"  ZH: {ZH_OUT}")
    print(f"\nPlease open these v4 files in PowerPoint to test.")
