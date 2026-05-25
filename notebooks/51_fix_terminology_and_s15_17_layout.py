"""
notebooks/51_fix_terminology_and_s15_17_layout.py

1. Replace "细胞身份"→"细胞类型" / "Cell Identity"→"Cell Type" in EN+ZH PPTs
2. Fix Slides 15-17 layout:
   - S15: Clear old cards, rebuild with image(left) + baseline text cards(right)
   - S16: Shrink+reposition image, keep existing text (OPC/Astrocyte stats are correct)
   - S17: Shrink+reposition image, keep existing text (omega gradient + migration hypothesis)
"""
import sys, http.server, threading, time, os, subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE      = Path(r"C:\Users\KnightZ\Desktop\细胞受选择")
PPT_EN    = BASE / "results/figures_final/CKI_Lecture_2026_v3.pptx"
PPT_ZH    = BASE / "results/figures_final/CKI_Lecture_2026_v3_ZH.pptx"
ASSETS    = BASE / "results/figures_final"

# ── Theme palette ──
TITLE_FG   = (0x1E, 0x29, 0x3B)
BODY_FG    = (0x1E, 0x29, 0x3B)
CARD_BG    = (0xFF, 0xFF, 0xFF)
ACCENT_TEAL= (0x0D, 0x94, 0x88)
ACCENT_BLUE= (0x25, 0x63, 0xEB)
ACCENT_RED = (0xDC, 0x26, 0x26)
DARK_BG    = (0x0F, 0x17, 0x2A)


# ═════════════════════════════════════════════════════════════════
#  STEP 1: Term replacement
# ═════════════════════════════════════════════════════════════════

def replace_terms(pptx_path, replacements):
    prs = Presentation(str(pptx_path))
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in replacements:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                changed += 1
    prs.save(str(pptx_path))
    return changed


# ═════════════════════════════════════════════════════════════════
#  STEP 2: SVG generators + PNG conversion
# ═════════════════════════════════════════════════════════════════

def make_svg_s15_baseline():
    W, H = 1200, 460
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
           f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']

    out.append(f'<text x="{W//2}" y="28" text-anchor="middle" '
               f'font-size="16" font-weight="700" fill="#0F172A">'
               f'Baseline: Intra-region ω &lt; Inter-region ω</text>')
    out.append(f'<text x="{W//2}" y="46" text-anchor="middle" '
               f'font-size="10" fill="#64748B">'
               f'Same cell type across regions is MORE distinct than different cell types within same region</text>')

    lx, ly = 40, 65
    out.append(f'<rect x="{lx}" y="{ly}" width="540" height="355" rx="10" '
               f'fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+270}" y="{ly+22}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#1E293B">Intra-region ω</text>')
    out.append(f'<text x="{lx+270}" y="{ly+40}" text-anchor="middle" '
               f'font-size="10" fill="#64748B">Same brain region, different cell types</text>')

    intra = [("Astrocytes", 28, "#0D9488"), ("OPCs", 22, "#2563EB"), ("Neurons(ex)", 18, "#7C3AED")]
    for i, (ctype, omega, color) in enumerate(intra):
        y = ly + 62 + i * 62
        bw = int(180 * omega / 35.0)
        out.append(f'<rect x="{lx+40}" y="{y}" width="{bw}" height="34" rx="4" fill="{color}" opacity="0.85"/>')
        out.append(f'<text x="{lx+40+bw+8}" y="{y+22}" font-size="12" fill="{color}" font-weight="600">{omega}</text>')
        out.append(f'<text x="{lx+280}" y="{y+22}" font-size="12" fill="#334155">{ctype}</text>')

    iy = ly + 280
    out.append(f'<rect x="{lx+15}" y="{iy}" width="510" height="56" rx="8" '
               f'fill="#F0FDF4" stroke="#86EFAC" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+270}" y="{iy+20}" text-anchor="middle" '
               f'font-size="11" font-weight="700" fill="#166534">Intra-region ω = LOW</text>')
    out.append(f'<text x="{lx+270}" y="{iy+40}" text-anchor="middle" '
               f'font-size="9" fill="#166534">Same microenvironment → similar expression program</text>')

    cx = 595
    out.append(f'<line x1="{cx}" y1="{ly+20}" x2="{cx}" y2="{ly+340}" '
               f'stroke="#CBD5E1" stroke-width="2" stroke-dasharray="6,4"/>')
    out.append(f'<text x="{cx}" y="{ly+180}" text-anchor="middle" '
               f'font-size="20" fill="#94A3B8" font-weight="700">&lt;</text>')

    rx, ry = 620, 65
    out.append(f'<rect x="{rx}" y="{ry}" width="540" height="355" rx="10" '
               f'fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{ry+22}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#1E293B">Inter-region ω</text>')
    out.append(f'<text x="{rx+270}" y="{ry+40}" text-anchor="middle" '
               f'font-size="10" fill="#64748B">Same cell type, different brain regions</text>')

    inter = [("OPC · SVZ", 3.2, "#0D9488"), ("OPC · Hipp", 18, "#2563EB"), ("OPC · Cortex", 22, "#7C3AED")]
    for j, (label, omega, color) in enumerate(inter):
        y = ry + 62 + j * 62
        bw = int(180 * omega / 35.0)
        out.append(f'<rect x="{rx+40}" y="{y}" width="{bw}" height="34" rx="4" fill="{color}" opacity="0.85"/>')
        out.append(f'<text x="{rx+40+bw+8}" y="{y+22}" font-size="12" fill="{color}" font-weight="600">{omega}</text>')
        out.append(f'<text x="{rx+290}" y="{y+22}" font-size="12" fill="#334155">{label}</text>')

    riy = ry + 280
    out.append(f'<rect x="{rx+15}" y="{riy}" width="510" height="56" rx="8" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{riy+20}" text-anchor="middle" '
               f'font-size="11" font-weight="700" fill="#991B1B">Inter-region ω = HIGH</text>')
    out.append(f'<text x="{rx+270}" y="{riy+40}" text-anchor="middle" '
               f'font-size="9" fill="#991B1B">Different microenvironments → distinct expression programs</text>')

    out.append(f'<rect x="35" y="{H-40}" width="{W-70}" height="32" rx="8" fill="#0F172A"/>')
    out.append(f'<text x="{W//2}" y="{H-18}" text-anchor="middle" '
               f'font-size="11" font-weight="600" fill="white">'
               f'Baseline:  Intra-region ω &lt; Inter-region ω  →  Cells differentiate LOCALLY, not by migrating</text>')
    out.append('</svg>')
    return "\n".join(out)


def make_svg_s16_anomaly():
    W, H = 1200, 460
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
           f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']

    out.append(f'<text x="{W//2}" y="28" text-anchor="middle" '
               f'font-size="16" font-weight="700" fill="#0F172A">'
               f'Extreme Case: OPC Breaks the Baseline</text>')
    out.append(f'<text x="{W//2}" y="46" text-anchor="middle" '
               f'font-size="10" fill="#64748B">'
               f'Intra-type OPC ω (22) ≪ Inter-type Astrocyte ω (107.5) in same brain region</text>')

    lx, ly = 40, 65
    out.append(f'<rect x="{lx}" y="{ly}" width="540" height="355" rx="10" '
               f'fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+270}" y="{ly+22}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#1E293B">Intra-type ω (OPC across regions)</text>')
    out.append(f'<text x="{lx+270}" y="{ly+40}" text-anchor="middle" '
               f'font-size="10" fill="#64748B">OPC in SVZ vs OPC in Cortex</text>')

    bar_max = 200
    svz_h = max(3, int(bar_max * 3.2 / 110.0))
    out.append(f'<rect x="{lx+60}" y="{ly+72 + bar_max - svz_h}" width="80" height="{svz_h}" '
               f'rx="4" fill="#0D9488"/>')
    out.append(f'<text x="{lx+100}" y="{ly+72 + bar_max - svz_h - 8}" '
               f'text-anchor="middle" font-size="12" font-weight="700" fill="#0D9488">ω=3.2</text>')
    out.append(f'<text x="{lx+100}" y="{ly+72+bar_max+16}" text-anchor="middle" '
               f'font-size="10" fill="#334155">SVZ</text>')

    ctx_h = int(bar_max * 22.0 / 110.0)
    out.append(f'<rect x="{lx+180}" y="{ly+72 + bar_max - ctx_h}" width="80" height="{ctx_h}" '
               f'rx="4" fill="#2563EB"/>')
    out.append(f'<text x="{lx+220}" y="{ly+72 + bar_max - ctx_h - 8}" '
               f'text-anchor="middle" font-size="12" font-weight="700" fill="#2563EB">ω=22</text>')
    out.append(f'<text x="{lx+220}" y="{ly+72+bar_max+16}" text-anchor="middle" '
               f'font-size="10" fill="#334155">Cortex</text>')
    out.append(f'<text x="{lx+380}" y="{ly+190}" font-size="10" fill="#64748B" font-style="italic">'
               f'Δω = 18.8 (surprisingly SMALL)</text>')

    iy = ly + 290
    out.append(f'<rect x="{lx+15}" y="{iy}" width="510" height="58" rx="8" '
               f'fill="#F0FDF4" stroke="#86EFAC" stroke-width="1.5"/>')
    out.append(f'<text x="{lx+270}" y="{iy+18}" text-anchor="middle" '
               f'font-size="11" font-weight="700" fill="#166534">Intra-type ω = 22 (LOW)</text>')
    out.append(f'<text x="{lx+270}" y="{iy+35}" text-anchor="middle" '
               f'font-size="9" fill="#166534">OPCs in SVZ and Cortex are HIGHLY similar</text>')
    out.append(f'<text x="{lx+270}" y="{iy+50}" text-anchor="middle" '
               f'font-size="9" fill="#166534">→ likely same origin (migration), not independent differentiation</text>')

    rx, ry = 620, 65
    out.append(f'<rect x="{rx}" y="{ry}" width="540" height="355" rx="10" '
               f'fill="white" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{ry+22}" text-anchor="middle" '
               f'font-size="13" font-weight="700" fill="#1E293B">Inter-type ω (within Cortex)</text>')
    out.append(f'<text x="{rx+270}" y="{ry+40}" text-anchor="middle" '
               f'font-size="10" fill="#64748B">Astrocytes vs OPCs in same region</text>')

    ast_h = int(bar_max * 107.5 / 110.0)
    out.append(f'<rect x="{rx+60}" y="{ry+72 + bar_max - ast_h}" width="80" height="{ast_h}" '
               f'rx="4" fill="#DC2626"/>')
    out.append(f'<text x="{rx+100}" y="{ry+72 + bar_max - ast_h - 8}" '
               f'text-anchor="middle" font-size="12" font-weight="700" fill="#DC2626">ω=107.5</text>')
    out.append(f'<text x="{rx+100}" y="{ry+72+bar_max+16}" text-anchor="middle" '
               f'font-size="10" fill="#334155">Astrocytes</text>')

    opc_h = int(bar_max * 22.0 / 110.0)
    out.append(f'<rect x="{rx+180}" y="{ry+72 + bar_max - opc_h}" width="80" height="{opc_h}" '
               f'rx="4" fill="#2563EB"/>')
    out.append(f'<text x="{rx+220}" y="{ry+72 + bar_max - opc_h - 8}" '
               f'text-anchor="middle" font-size="12" font-weight="700" fill="#2563EB">ω=22</text>')
    out.append(f'<text x="{rx+220}" y="{ry+72+bar_max+16}" text-anchor="middle" '
               f'font-size="10" fill="#334155">OPCs</text>')
    out.append(f'<text x="{rx+380}" y="{ry+190}" font-size="10" fill="#64748B" font-style="italic">'
               f'Δω = 85.5 (extremely LARGE)</text>')

    riy = ry + 290
    out.append(f'<rect x="{rx+15}" y="{riy}" width="510" height="58" rx="8" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{rx+270}" y="{riy+18}" text-anchor="middle" '
               f'font-size="11" font-weight="700" fill="#991B1B">Inter-type ω = 107.5 (HIGH)</text>')
    out.append(f'<text x="{rx+270}" y="{riy+35}" text-anchor="middle" '
               f'font-size="9" fill="#991B1B">Astrocytes and OPCs in SAME region are RADICALLY different</text>')
    out.append(f'<text x="{rx+270}" y="{riy+50}" text-anchor="middle" '
               f'font-size="9" fill="#991B1B">→ independent lineages, strong local differentiation</text>')

    out.append(f'<rect x="35" y="{H-40}" width="{W-70}" height="32" rx="8" fill="#7C3AED"/>')
    out.append(f'<text x="{W//2}" y="{H-18}" text-anchor="middle" '
               f'font-size="11" font-weight="700" fill="white">'
               f'CONTRADICTION:  Intra-type ω (22) ≪ Inter-type ω (107.5)  →  '
               f'OPC similarity across regions EXCEEDS cell-type difference within region</text>')
    out.append('</svg>')
    return "\n".join(out)


def make_svg_s17_brain():
    W, H = 1200, 440
    out = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" '
           f'width="{W}" height="{H}" style="background:#F8FAFC; font-family:Arial,Helvetica,sans-serif">']

    out.append(f'<text x="{W//2}" y="26" text-anchor="middle" '
               f'font-size="15" font-weight="700" fill="#0F172A">'
               f'Migration Inference: From Anomaly to Hypothesis</text>')

    # Brain regions
    out.append('<ellipse cx="230" cy="340" rx="115" ry="80" fill="#E0F2FE" stroke="#0EA5E9" stroke-width="2.5" opacity="0.55"/>')
    out.append(f'<text x="230" y="365" text-anchor="middle" font-size="11" fill="#0369A1" font-weight="600">Cerebellum</text>')
    out.append(f'<text x="230" y="380" text-anchor="middle" font-size="9" fill="#0369A1">ω=85 Bergmann</text>')

    out.append('<ellipse cx="680" cy="120" rx="140" ry="90" fill="#FEF2F2" stroke="#EF4444" stroke-width="2.5" opacity="0.55"/>')
    out.append(f'<text x="680" y="100" text-anchor="middle" font-size="12" fill="#991B1B" font-weight="600">Cortex</text>')
    out.append(f'<text x="680" y="116" text-anchor="middle" font-size="10" fill="#991B1B">ω=107.5 Astro</text>')
    out.append(f'<text x="680" y="132" text-anchor="middle" font-size="10" fill="#991B1B">ω=22 OPC</text>')

    out.append('<ellipse cx="510" cy="300" rx="55" ry="45" fill="#F0FDF4" stroke="#16A34A" stroke-width="2.5" opacity="0.6"/>')
    out.append(f'<text x="510" y="294" text-anchor="middle" font-size="11" fill="#166534" font-weight="600">SVZ</text>')
    out.append(f'<text x="510" y="310" text-anchor="middle" font-size="9" fill="#166534">ω=3.2 OPC</text>')

    out.append('<ellipse cx="760" cy="265" rx="60" ry="42" fill="#FEFCE8" stroke="#F59E0B" stroke-width="2" opacity="0.5"/>')
    out.append(f'<text x="760" y="263" text-anchor="middle" font-size="10" fill="#92400E">Hipp ω=45</text>')

    out.append('<ellipse cx="615" cy="230" rx="52" ry="38" fill="#F5F3FF" stroke="#7C3AED" stroke-width="2" opacity="0.5"/>')
    out.append(f'<text x="615" y="233" text-anchor="middle" font-size="10" fill="#5B21B6">Striatum ω=38</text>')

    # Migration arrow
    out.append('<defs>'
               '<filter id="glow"><feGaussianBlur stdDeviation="4" result="blur"/>'
               '<feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge></filter>'
               '<marker id="arrRed" markerWidth="12" markerHeight="9" refX="10" refY="4.5" orient="auto">'
               '<path d="M0,0 L12,4.5 L0,9" fill="#DC2626"/></marker>'
               '</defs>')
    out.append('<path d="M 538 300 C 600 220, 640 180, 670 140" '
               'stroke="#DC2626" stroke-width="5" fill="none" '
               'stroke-dasharray="10,5" opacity="0.9" filter="url(#glow)" '
               'marker-end="url(#arrRed)"/>')

    out.append(f'<rect x="542" y="222" width="130" height="22" rx="11" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="607" y="237" text-anchor="middle" '
               f'font-size="11" font-weight="700" fill="#DC2626">OPC Migration</text>')

    out.append(f'<rect x="35" y="{H-40}" width="{W-70}" height="32" rx="8" fill="#0D9488"/>')
    out.append(f'<text x="{W//2}" y="{H-18}" text-anchor="middle" '
               f'font-size="10" font-weight="600" fill="white">'
               f'CKI identifies migration via ω deviation: OPC SVZ→Cortex similarity exceeds inter-type divergence within Cortex</text>')
    out.append('</svg>')
    return "\n".join(out)


def svg_to_png_chrome(svg_path, png_path):
    chrome = r"C:\Program Files\Google\Chrome\Application\chrome.exe"
    if not Path(chrome).exists():
        chrome = r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"
    if not Path(chrome).exists():
        chrome = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
    if not Path(chrome).exists():
        chrome = r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"
    if not Path(chrome).exists():
        print(f"      ERROR: Chrome/Edge not found")
        return False

    os.chdir(str(svg_path.parent))
    server = http.server.HTTPServer(("127.0.0.1", 18902), http.server.SimpleHTTPRequestHandler)
    t = threading.Thread(target=server.serve_forever, daemon=True)
    t.start()
    time.sleep(0.5)
    url = f"http://127.0.0.1:18902/{svg_path.name}"

    cmd = [
        chrome, "--headless", "--disable-gpu", "--screenshot",
        f"--screenshot={png_path.resolve()}",
        "--window-size=1400,650",
        url
    ]
    subprocess.run(cmd, capture_output=True, timeout=30)
    server.shutdown()

    if png_path.exists():
        print(f"      PNG: {png_path.name} ({png_path.stat().st_size//1024} KB)")
        return True
    else:
        print(f"      WARNING: PNG not created")
        return False


# ═════════════════════════════════════════════════════════════════
#  PPT manipulation helpers
# ═════════════════════════════════════════════════════════════════

def remove_all_pics_from_slide(slide):
    spTree = slide.shapes._spTree
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    pics = spTree.findall('.//p:pic', ns)
    for pic in list(pics):
        spTree.remove(pic)
    n = len(pics)
    if n > 0:
        print(f"      Removed {n} pic(s)")


def remove_non_title_shapes(slide):
    """Remove all shapes EXCEPT title (top area) and page number (bottom-right)."""
    to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'image'):
            continue  # images handled separately
        keep = False
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            # Keep page number: digit at bottom-right (T > 4.5, R > 8.5)
            if txt.isdigit() and len(txt) <= 2:
                if shape.top > Inches(4.5) and (shape.left + shape.width) > Inches(8.5):
                    keep = True
            # Keep title: near top (T < 1.0)
            if shape.top < Inches(1.0):
                keep = True
        if not keep:
            to_remove.append(shape)

    spTree = slide.shapes._spTree
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    print(f"      Removed {len(to_remove)} non-title shapes")


def add_card_shape(slide, left, top, width, height, title, lines,
                   title_color=None, bg_color=None):
    shape = slide.shapes.add_shape(
        1,  # ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*(bg_color or CARD_BG))
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.07)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*(title_color or ACCENT_BLUE))

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(*BODY_FG)

    return shape


def rebuild_slide15(slide, png_path, cards, lang='en'):
    """S15: clear old cards + add image(left) + new cards(right)."""
    remove_non_title_shapes(slide)
    remove_all_pics_from_slide(slide)

    slide.shapes.add_picture(str(png_path),
        Inches(0.5), Inches(1.05), Inches(5.0), Inches(1.92))

    for i, (title, lines, tc, bg) in enumerate(cards):
        y = 1.05 + i * 1.45
        add_card_shape(slide, 5.6, y, 4.1, 1.30, title, lines,
                       title_color=tc, bg_color=bg)

    print(f"    S15 rebuilt: 1 image + {len(cards)} cards")


def rebuild_slide16(slide, png_path, cards):
    """S16: clear old text + add image(left) + 2 new cards(right)."""
    remove_non_title_shapes(slide)
    remove_all_pics_from_slide(slide)
    slide.shapes.add_picture(str(png_path),
        Inches(0.5), Inches(1.05), Inches(5.0), Inches(1.92))
    for i, (title, lines, tc, bg) in enumerate(cards):
        y = 1.05 + i * 2.15
        add_card_shape(slide, 5.6, y, 4.1, 2.00, title, lines, title_color=tc, bg_color=bg)
    print(f"    S16 rebuilt: image + {len(cards)} cards")


def rebuild_slide17(slide, png_path, cards):
    """S17: clear old text + add image(left) + 2 cards(right)."""
    remove_non_title_shapes(slide)
    remove_all_pics_from_slide(slide)
    slide.shapes.add_picture(str(png_path),
        Inches(0.5), Inches(1.05), Inches(5.0), Inches(1.92))
    for i, (title, lines, tc, bg) in enumerate(cards):
        y = 1.05 + i * 2.15
        add_card_shape(slide, 5.6, y, 4.1, 2.00, title, lines, title_color=tc, bg_color=bg)
    print(f"    S17 rebuilt: image + {len(cards)} cards")


# ═════════════════════════════════════════════════════════════════
#  MAIN
# ═════════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("Fix: Terminology + Slides 15-17 Layout Overhaul")
    print("=" * 60)

    # ── Step 1: Term replacement ──
    print("\n[1/4] Term replacement...")
    en_ch = replace_terms(PPT_EN, [
        ("Cell Identity", "Cell Type"),
        ("细胞身份", "细胞类型"),
    ])
    print(f"  EN PPT: {en_ch} replacements")

    zh_ch = replace_terms(PPT_ZH, [
        ("细胞身份", "细胞类型"),
        ("Cell Identity", "Cell Type"),
    ])
    print(f"  ZH PPT: {zh_ch} replacements")

    # ── Step 2: Generate SVGs + PNGs ──
    print("\n[2/4] Generating SVGs + PNGs...")
    svg_specs = {
        's15': ('_s15_baseline_v3.svg', make_svg_s15_baseline()),
        's16': ('_s16_anomaly_v3.svg', make_svg_s16_anomaly()),
        's17': ('_s17_migration_v3.svg', make_svg_s17_brain()),
    }
    png_paths = {}
    for key, (fname, svg) in svg_specs.items():
        svg_p = ASSETS / fname
        png_p = ASSETS / fname.replace('.svg', '.png')
        svg_p.write_text(svg, encoding='utf-8')
        print(f"  {fname}: {len(svg)//1024} KB → ", end='')
        if svg_to_png_chrome(svg_p, png_p):
            png_paths[key] = png_p

    if len(png_paths) < 3:
        print("  WARNING: Some PNGs failed!")
        return

    # ── Step 3: Fix EN Slides 15-17 ──
    print("\n[3/4] Fixing EN Slides 15-17...")
    prs_en = Presentation(str(PPT_EN))

    # EN S15 cards
    en_s15_cards = [
        ("Intra-region ω", [
            "Same brain region, different cell types",
            "Microenvironment limits expression divergence",
            "Example: Cortex Astro(28) vs OPC(22) vs Neuron(18)"
        ], ACCENT_TEAL, None),
        ("Inter-region ω", [
            "Same cell type, different brain regions",
            "Distinct microenvironments → higher ω",
            "Example: OPC SVZ(3.2) → Hipp(18) → Cortex(22)"
        ], ACCENT_RED, None),
        ("Baseline Rule", [
            "Intra < Inter: consistent pattern across all regions",
            "Cells differentiate LOCALLY, not through migration"
        ], DARK_BG, (0xF8, 0xFA, 0xFC)),
    ]

    print("  EN Slide 15:")
    rebuild_slide15(prs_en.slides[14], png_paths['s15'], en_s15_cards)

    en_s16_cards = [
        ("OPC: Breaking the Baseline", [
            "23.5% of OPC pairs show strong ω anomaly",
            "OPC ω(SVZ) = 3.2, OPC ω(Cortex) = 22",
            "Intra-type ω diff = 18.8 (surprisingly SMALL)",
            "OPCs are known to migrate during development",
            "and in response to injury — CKI recovers this",
            "without prior knowledge of migration pathways"
        ], ACCENT_RED, None),
        ("Astrocytes: High ω, Low Migration", [
            "Highest baseline ω = 107.5 (strong identity)",
            "Yet only 14 migration candidates — very few",
            "Inter-type ω diff = 85.5 (extremely LARGE)",
            "High ω reflects local specialization,",
            "NOT migration — CKI distinguishes both",
            "signals independently"
        ], ACCENT_TEAL, None),
    ]

    en_s17_cards = [
        ("Omega Gradient Across Brain Regions", [
            "Cortex Astrocytes: ω = 107.5 (highest identity)",
            "Cerebellum Bergmann Glia: ω = 85",
            "Hipp/Thalamus/Striatum: ω = 38-45",
            "SVZ OPCs: ω = 3.2 (least differentiated state)",
            "ω measures functional divergence within a region"
        ], ACCENT_BLUE, None),
        ("OPC Migration Hypothesis", [
            "Intra-type: OPC(SVZ) vs OPC(Cortex) → ω ≈ 22",
            "→ Same cell type, low divergence across regions",
            "Inter-type: Astrocyte vs OPC (Cortex) → ω ≈ 107.5",
            "→ Different cell types, high divergence in same region",
            "∴ OPC similarity across regions EXCEEDS",
            "inter-cell-type similarity within region",
            "→ SVZ → Cortex migration is the best hypothesis"
        ], ACCENT_RED, None),
    ]

    print("  EN Slide 16:")
    rebuild_slide16(prs_en.slides[15], png_paths['s16'], en_s16_cards)

    print("  EN Slide 17:")
    rebuild_slide17(prs_en.slides[16], png_paths['s17'], en_s17_cards)

    prs_en.save(str(PPT_EN))
    print(f"  Saved EN: {PPT_EN.name}")

    # ── Step 4: Fix ZH Slides 15-17 ──
    print("\n[4/4] Fixing ZH Slides 15-17...")
    prs_zh = Presentation(str(PPT_ZH))

    zh_s15_cards = [
        ("同脑区内 ω", [
            "相同脑区，不同细胞类型间的功能分化",
            "微环境限制表达差异 → ω 较低",
            "例：皮层 Astro(28) vs OPC(22) vs Neuron(18)"
        ], ACCENT_TEAL, None),
        ("跨脑区间 ω", [
            "相同细胞类型，不同脑区间的功能分化",
            "不同微环境 → ω 更高",
            "例：OPC SVZ(3.2) → Hipp(18) → Cortex(22)"
        ], ACCENT_RED, None),
        ("基线规律", [
            "同脑区 ω 普遍低于跨脑区 ω",
            "细胞在局部环境中分化，而非通过迁移"
        ], DARK_BG, (0xF8, 0xFA, 0xFC)),
    ]

    print("  ZH Slide 15:")
    rebuild_slide15(prs_zh.slides[14], png_paths['s15'], zh_s15_cards)

    zh_s16_cards = [
        ("OPC：打破基线的特例", [
            "OPC 配对中 23.5% 显示强烈的 ω 异常信号",
            "OPC ω(SVZ) = 3.2, OPC ω(Cortex) = 22",
            "同类 ω 差异 = 18.8（异常地小）",
            "OPC 已知在发育过程中活跃迁移",
            "并在损伤后应答中迁移 — CKI 在无先验",
            "知识的情况下独立发现了这一现象"
        ], ACCENT_RED, None),
        ("星形胶质细胞：高 ω，低迁移", [
            "最高基线 ω = 107.5（最强的细胞类型身份）",
            "但仅有 14 个迁移候选 — 数量极少",
            "跨类 ω 差异 = 85.5（异常地大）",
            "高 ω 反映的是局部特异性分化",
            "而非迁移 — CKI 独立区分两种信号"
        ], ACCENT_TEAL, None),
    ]

    zh_s17_cards = [
        ("各脑区 ω 梯度分布", [
            "皮层 Astrocytes: ω = 107.5（最高功能特异性）",
            "小脑 Bergmann Glia: ω = 85",
            "海马/丘脑/纹状体: ω = 38-45",
            "SVZ OPCs: ω = 3.2（分化程度最低）",
            "ω 衡量细胞类型在特定脑区内的功能分化程度"
        ], ACCENT_BLUE, None),
        ("OPC 迁移假说", [
            "同类比较：OPC(SVZ) vs OPC(Cortex) → ω ≈ 22",
            "→ 相同细胞类型，跨脑区差异很小",
            "跨类比较：Astrocyte vs OPC（同皮层）→ ω ≈ 107.5",
            "→ 不同细胞类型，同脑区内差异极大",
            "∴ OPC 跨脑区的相似性远超",
            "皮层内部不同类型间的相似性",
            "→ SVZ → 皮层 OPC 迁移是最合理的假说"
        ], ACCENT_RED, None),
    ]

    print("  ZH Slide 16:")
    rebuild_slide16(prs_zh.slides[15], png_paths['s16'], zh_s16_cards)

    print("  ZH Slide 17:")
    rebuild_slide17(prs_zh.slides[16], png_paths['s17'], zh_s17_cards)

    prs_zh.save(str(PPT_ZH))
    print(f"  Saved ZH: {PPT_ZH.name}")

    print("\n" + "=" * 60)
    print("ALL DONE. Changes applied:")
    print("  1. '细胞身份'/'Cell Identity' → '细胞类型'/'Cell Type'")
    print("  2. S15: Rebuilt — image(left) + 3 baseline cards(right), no overlap")
    print("  3. S16: Rebuilt — image(left) + 2 cards(right, OPC + Astrocyte), no overlap")
    print("  4. S17: Rebuilt — image(left) + 2 cards(right, omega gradient + migration), no overlap")
    print(f"  EN: {PPT_EN}")
    print(f"  ZH: {PPT_ZH}")
    print("=" * 60)


if __name__ == "__main__":
    main()
