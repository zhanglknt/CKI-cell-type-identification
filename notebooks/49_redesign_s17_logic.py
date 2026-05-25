"""
Redesign Slide 17 v2: make OPC migration inference logic explicit.

Core logic:
  Intra-type omega (same cell type across regions) << Inter-type omega (different cell types in same region)
  → Migration is the most parsimonious explanation.

Visual: SVG brain schematic (upper 3/4) + comparison panel (lower 1/4).
"""

import os
import subprocess
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

BASE = r"C:\Users\KnightZ\Desktop\细胞受选择\results\figures_final"
PPT_EN = f"{BASE}/CKI_Lecture_2026_v3.pptx"
PPT_ZH = f"{BASE}/CKI_Lecture_2026_v3_ZH.pptx"

# Taller canvas to fit comparison panel
PNG_W, PNG_H = 2700, 560


def omega_color(omega):
    """Blue(222deg) → Red(0deg) via log scale."""
    import math
    lo, hi = 3, 110
    t = (math.log(max(omega, lo)) - math.log(lo)) / (math.log(hi) - math.log(lo))
    t = max(0, min(1, t))
    h = 222 - 220 * t
    s, l = 72, 48
    s /= 100; l /= 100
    c = (1 - abs(2*l - 1)) * s
    x = c * (1 - abs((h/60) % 2 - 1))
    m = l - c/2
    if h < 60:    r, g, b = c, x, 0
    elif h < 120: r, g, b = x, c, 0
    elif h < 180: r, g, b = 0, c, x
    elif h < 240: r, g, b = 0, x, c
    elif h < 300: r, g, b = x, 0, c
    else:         r, g, b = c, 0, x
    r, g, b = int((r+m)*255), int((g+m)*255), int((b+m)*255)
    return f"#{r:02X}{g:02X}{b:02X}"


def make_svg():
    """Generate SVG with brain + bottom comparison panel."""
    out = []
    out.append(f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {PNG_W} {PNG_H}">')
    out.append('<defs>')
    # Brain gradient
    out.append('<linearGradient id="brainGrad" x1="0" y1="0" x2="1" y2="0">')
    out.append('  <stop offset="0%" stop-color="#E8EDF4"/>')
    out.append('  <stop offset="100%" stop-color="#F5F7FA"/>')
    out.append('</linearGradient>')
    # Glow for migration arrow
    out.append('<filter id="glow" x="-20%" y="-20%" width="140%" height="140%">')
    out.append('  <feGaussianBlur stdDeviation="6" result="blur"/>')
    out.append('  <feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge>')
    out.append('</filter>')
    out.append('<filter id="shadow" x="-10%" y="-10%" width="120%" height="120%">')
    out.append('  <feDropShadow dx="1.5" dy="1.5" stdDeviation="2" flood-opacity="0.18"/>')
    out.append('</filter>')
    out.append('<marker id="arrowRed" markerWidth="14" markerHeight="9" refX="12" refY="4.5" orient="auto" markerUnits="userSpaceOnUse">')
    out.append('  <path d="M 0,0 L 14,4.5 L 0,9 z" fill="#DC2626"/>')
    out.append('</marker>')
    # Background gradient for comparison panel
    out.append('<linearGradient id="cmpGrad" x1="0" y1="0" x2="1" y2="0">')
    out.append('  <stop offset="0%" stop-color="#F0F9FF"/>')
    out.append('  <stop offset="100%" stop-color="#FFF7ED"/>')
    out.append('</linearGradient>')
    out.append('</defs>')

    # Background
    out.append(f'<rect width="{PNG_W}" height="{PNG_H}" fill="#F8FAFC" rx="10"/>')

    # ================================================================
    # BRAIN OUTLINE (top 420px area)
    # ================================================================
    brain_path = (
        'M 70,230 '
        'C 70,140 110,65 210,45 '
        'C 380,18 720,5 1300,12 '
        'C 1700,18 2050,30 2320,45 '
        'C 2480,60 2590,95 2630,135 '
        'C 2660,168 2665,205 2645,245 '
        'C 2600,298 2500,336 2340,358 '
        'C 2120,385 1750,395 1350,395 '
        'C 1000,395 700,385 400,368 '
        'C 230,350 120,330 75,290 '
        'Z'
    )
    out.append(f'<path d="{brain_path}" fill="url(#brainGrad)" stroke="#94A3B8" stroke-width="2.2"/>')

    # Internal anatomical references
    out.append('<path d="M 320,72 C 600,48 1050,44 1600,55" '
               'fill="none" stroke="#CBD5E1" stroke-width="1" stroke-dasharray="6,5" opacity="0.6"/>')
    out.append('<path d="M 350,168 C 550,148 850,140 1150,145 C 1380,150 1580,155 1720,162" '
               'fill="none" stroke="#CBD5E1" stroke-width="1.3" stroke-dasharray="6,5" opacity="0.55"/>')

    # ================================================================
    # REGIONS (compact for the taller layout)
    # ================================================================
    REGIONS = [
        ("Cortex",     1850, 92, 60, 107.5, "top",    1850, 14,
         [("Astrocytes \u03c9=107.5", "#DC2626"), ("OPCs \u03c9=22", "#D97706")]),
        ("SVZ",         980, 220, 34, 3.2, "bottom", 980, 290,
         [("OPCs \u03c9=3.2", "#2563EB")]),
        ("Hippocampus",1600, 138, 40, 45, "right",  1710, 130,
         [("Neurons \u03c9=45", "#A855F7")]),
        ("Striatum",   1120, 275, 34, 38, "bottom",  1120, 340,
         [("Neurons \u03c9=38", "#8B5CF6")]),
        ("Thalamus",   1320, 248, 34, 42, "bottom",  1320, 315,
         [("Neurons \u03c9=42", "#A855F7")]),
        ("Cerebellum", 2450, 230, 50, 85, "top",    2450, 155,
         [("Bergmann Glia \u03c9=85", "#EF4444")]),
    ]

    for name, cx, cy, r, omega, align, lx, ly, details in REGIONS:
        color = omega_color(omega)
        out.append(f'<circle cx="{cx}" cy="{cy}" r="{r+4}" fill="none" '
                   f'stroke="{color}" stroke-width="2.5" opacity="0.22"/>')
        out.append(f'<circle cx="{cx}" cy="{cy}" r="{r}" fill="{color}" '
                   f'fill-opacity="0.12" stroke="{color}" stroke-width="2.2" filter="url(#shadow)"/>')

        omega_str = f"\u03c9={omega}"
        fs = 16 if omega < 100 else 14
        out.append(f'<text x="{cx}" y="{cy+1}" text-anchor="middle" dominant-baseline="middle" '
                   f'font-family="Arial,Helvetica,sans-serif" font-size="{fs}" '
                   f'font-weight="800" fill="{color}">{omega_str}</text>')

        # Region label
        anchor = "middle" if align in ("top", "bottom") else ("start" if lx > cx else "end")
        out.append(f'<text x="{lx}" y="{ly}" text-anchor="{anchor}" '
                   f'font-family="Arial,Helvetica,sans-serif" font-size="15" '
                   f'font-weight="700" fill="{color}">{name}</text>')

        # Cell type detail lines
        for di, (dtxt, dcol) in enumerate(details):
            if align == "top":
                dx, dy, da = cx, ly + 22 + di * 17, "middle"
            elif align == "bottom":
                dx, dy, da = cx, ly + 22 + di * 17, "middle"
            else:
                dx, dy, da = lx + 6, cy + (di - len(details)/2 + 0.5) * 15, "start"
            if 0 < dy < 420:
                out.append(f'<text x="{dx}" y="{dy}" text-anchor="{da}" '
                           f'font-family="Arial,Helvetica,sans-serif" font-size="11" '
                           f'fill="{dcol}" font-style="italic">{dtxt}</text>')

    # ================================================================
    # OPC MIGRATION ARROW
    # ================================================================
    svz_x, svz_y = 980, 220
    ctx_x, ctx_y = 1850, 92
    mig_path = (f'M {svz_x+26},{svz_y-20} '
                f'C {svz_x+150},{110} {ctx_x-200},{40} {ctx_x-38},{ctx_y-30}')
    out.append(f'<path d="{mig_path}" fill="none" stroke="#DC2626" stroke-width="5" '
               f'marker-end="url(#arrowRed)" filter="url(#glow)" opacity="0.92"/>')
    mig_path2 = (f'M {svz_x+26},{svz_y-8} '
                 f'C {svz_x+170},{140} {ctx_x-220},{70} {ctx_x-38},{ctx_y-18}')
    out.append(f'<path d="{mig_path2}" fill="none" stroke="#EF4444" stroke-width="2.5" '
               f'opacity="0.5" stroke-linecap="round"/>')

    # Migration pill
    lbl_x = (svz_x + ctx_x) / 2 + 80
    lbl_y = (svz_y + ctx_y) / 2 - 85
    pill_w, pill_h = 148, 28
    out.append(f'<rect x="{lbl_x - pill_w/2:.0f}" y="{lbl_y - pill_h/2:.0f}" '
               f'width="{pill_w}" height="{pill_h}" rx="14" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{lbl_x}" y="{lbl_y + 1}" text-anchor="middle" dominant-baseline="middle" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="13" '
               f'font-weight="800" fill="#DC2626">OPC Migration</text>')

    # Omega gradient annotation on arrow
    out.append(f'<text x="{lbl_x + 115}" y="{lbl_y - 5}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10" '
               f'fill="#94A3B8" font-style="italic">\u03c9 3\u219222</text>')

    # ================================================================
    # COMPARISON PANEL (y=430 to y=555)
    # ================================================================
    panel_y, panel_h = 432, 122
    out.append(f'<rect x="60" y="{panel_y}" width="{PNG_W-120}" height="{panel_h}" '
               f'fill="url(#cmpGrad)" rx="10" stroke="#E2E8F0" stroke-width="1.5"/>')
    out.append(f'<text x="{PNG_W/2}" y="{panel_y+22}" text-anchor="middle" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="13" '
               f'font-weight="700" fill="#475569">'
               f'\u03c9 Inference Logic: Migration Hypothesis</text>')

    # ---- LEFT SIDE: Intra-type (low omega → migration) ----
    left_x = 120
    left_w = 1160
    # Label
    out.append(f'<text x="{left_x}" y="{panel_y+48}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="12" '
               f'font-weight="700" fill="#2563EB">Intra-type \u03c9</text>')
    out.append(f'<text x="{left_x}" y="{panel_y+66}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10.5" '
               f'fill="#64748B">Same cell type (OPC) across brain regions</text>')

    # Bar: OPC comparison bar (narrow, blue)
    bar_start = left_x + 20
    bar_end = bar_start + 300
    bar_mid = panel_y + 90
    bar_h = 22

    # Intrinsic low omega: just show narrow bars
    # SVZ OPC omega base (small bar)
    out.append(f'<rect x="{bar_start}" y="{bar_mid}" width="22" height="{bar_h}" '
               f'fill="#2563EB" rx="3"/>')
    out.append(f'<text x="{bar_start+1}" y="{bar_mid+bar_h+15}" text-anchor="middle" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10" '
               f'fill="#2563EB" font-weight="700">SVZ\n\u03c9=3.2</text>')

    # Arrow between
    arr_x = bar_start + 45
    arr_w = 200
    out.append(f'<line x1="{arr_x}" y1="{bar_mid+bar_h/2}" x2="{arr_x+arr_w}" y2="{bar_mid+bar_h/2}" '
               f'stroke="#94A3B8" stroke-width="1.5" marker-end="url(#arrowRed)"/>')
    out.append(f'<text x="{arr_x+arr_w/2}" y="{bar_mid-6}" text-anchor="middle" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10" '
               f'fill="#64748B" font-style="italic">\u03c9 \u2248 22</text>')

    # Cortex OPC (slightly wider bar indicating higher omega)
    ctx_opc_x = arr_x + arr_w + 20
    ctx_opc_w = 48
    out.append(f'<rect x="{ctx_opc_x}" y="{bar_mid}" width="{ctx_opc_w}" height="{bar_h}" '
               f'fill="#D97706" rx="3"/>')
    out.append(f'<text x="{ctx_opc_x+ctx_opc_w/2}" y="{bar_mid+bar_h+15}" text-anchor="middle" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10" '
               f'fill="#D97706" font-weight="700">Cortex\n\u03c9=22</text>')

    # Conclusion arrow
    concl_x = ctx_opc_x + ctx_opc_w + 30
    out.append(f'<text x="{concl_x}" y="{bar_mid+bar_h/2+1}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10.5" '
               f'fill="#059669" font-weight="700">\u2192 LOW divergence</text>')
    out.append(f'<text x="{concl_x}" y="{bar_mid+bar_h/2+17}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10.5" '
               f'fill="#059669" font-weight="700">\u2192 COMMON ORIGIN</text>')
    out.append(f'<text x="{concl_x}" y="{bar_mid+bar_h/2+33}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10.5" '
               f'fill="#059669" font-weight="700">\u2192 MIGRATION</text>')

    # ---- RIGHT SIDE: Inter-type (high omega → distinct) ----
    right_x = 1390
    right_w = 1160
    out.append(f'<text x="{right_x}" y="{panel_y+48}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="12" '
               f'font-weight="700" fill="#DC2626">Inter-type \u03c9</text>')
    out.append(f'<text x="{right_x}" y="{panel_y+66}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10.5" '
               f'fill="#64748B">Different cell types within the same region (Cortex)</text>')

    # Bars showing multiple cell types with wide gaps
    rbar_start = right_x + 20
    rbar_mid = bar_mid  # same y

    # Astrocyte (very wide bar = very different)
    astro_w = 130
    out.append(f'<rect x="{rbar_start}" y="{rbar_mid}" width="{astro_w}" height="{bar_h}" '
               f'fill="#DC2626" rx="3"/>')
    out.append(f'<text x="{rbar_start+2}" y="{rbar_mid+bar_h+15}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10" '
               f'fill="#DC2626" font-weight="700">Astrocyte \u03c9=107.5</text>')

    # Gap between bars
    gap_x = rbar_start + astro_w + 20

    # OPC in cortex
    opc2_x = gap_x + 50
    out.append(f'<rect x="{opc2_x}" y="{rbar_mid}" width="45" height="{bar_h}" '
               f'fill="#D97706" rx="3"/>')
    out.append(f'<text x="{opc2_x}" y="{rbar_mid+bar_h+15}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10" '
               f'fill="#D97706" font-weight="700">OPC \u03c9=22</text>')

    # Omega difference annotation
    diff_label_x = gap_x + 10
    out.append(f'<text x="{diff_label_x}" y="{rbar_mid-6}" text-anchor="middle" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="11" '
               f'fill="#DC2626" font-weight="800">\u03c9 \u2248 107</text>')

    # Right conclusion
    rconcl_x = opc2_x + 45 + 30
    out.append(f'<text x="{rconcl_x}" y="{rbar_mid+rbar_mid/2+1}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10.5" '
               f'fill="#DC2626" font-weight="700">\u2192 HIGH divergence</text>')
    out.append(f'<text x="{rconcl_x}" y="{rbar_mid+rbar_mid/2+17}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10.5" '
               f'fill="#DC2626" font-weight="700">\u2192 DISTINCT LINEAGES</text>')
    out.append(f'<text x="{rconcl_x}" y="{rbar_mid+rbar_mid/2+33}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10.5" '
               f'fill="#DC2626" font-weight="700">\u2192 LOCAL DIFFERENTIATION</text>')

    out.append('</svg>')
    return '\n'.join(out)


def svg_to_png(svg_path, png_path):
    """Convert SVG to hi-res PNG via Chrome headless screenshot."""
    html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0}}body{{
display:flex;justify-content:center;align-items:center;
background:#F8FAFC;width:{PNG_W}px;height:{PNG_H}px;}}
svg{{width:{PNG_W}px;height:{PNG_H}px;display:block;}}</style></head>
<body>
{open(svg_path, encoding='utf-8').read()}
</body></html>"""

    html_path = svg_path.replace('.svg', '.html')
    with open(html_path, 'w', encoding='utf-8') as f:
        f.write(html)

    chrome_paths = [
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    chrome = None
    for p in chrome_paths:
        if os.path.exists(p):
            chrome = p
            break
    if not chrome:
        raise RuntimeError("No browser found")

    file_url = f"file:///{html_path.replace(chr(92), '/')}"
    cmd = [
        chrome, '--headless=new', '--disable-gpu', '--no-sandbox',
        f'--window-size={PNG_W},{PNG_H}',
        f'--screenshot={png_path}', file_url
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
    if result.returncode != 0:
        print(f"  Browser stderr: {result.stderr[:500]}")

    os.remove(html_path)
    return os.path.exists(png_path) and os.path.getsize(png_path) > 5000


def replace_brain_and_text(pptx_path, png_path, is_zh=False):
    """Replace brain image + rewrite card texts with clear inference logic."""
    prs = Presentation(pptx_path)
    slide = prs.slides[16]

    # Remove old picture, add new one
    old_pic = slide.shapes[11]
    left, top, width, height = old_pic.left, old_pic.top, old_pic.width, old_pic.height
    old_pic._element.getparent().remove(old_pic._element)
    slide.shapes.add_picture(png_path, left, top, width, height)

    # Left card: describe the omega landscape across brain regions
    left_tf = slide.shapes[6].text_frame; left_tf.clear()
    if is_zh:
        left_lines = [
            "\u03c9 衡量细胞类型在特定脑区中的功能分化程度。",
            "\u03c9 越高 = 该区域细胞功能越特化、越偏离基线。",
            "皮层 Astrocytes: \u03c9=107.5（极高分化, 最独特的胶质身份）",
            "小脑 Bergmann Glia: \u03c9=85（独立聚类, 辐射胶质特征）",
            "海马/丘脑/纹状体: \u03c9=38\u201345（中等神经元特征）",
            "SVZ OPCs: \u03c9=3.2（极低分化, 接近基线状态）",
        ]
    else:
        left_lines = [
            "\u03c9 measures functional divergence of a cell type within a brain region.",
            "Higher \u03c9 = greater specialization away from baseline expression.",
            "Cortex Astrocytes: \u03c9=107.5 (highest divergence \u2014 most distinct glial identity)",
            "Cerebellum Bergmann Glia: \u03c9=85 (unique radial glia cluster)",
            "Hippocampus/Thalamus/Striatum: \u03c9=38\u201345 (intermediate neuronal signature)",
            "SVZ OPCs: \u03c9=3.2 (near-baseline \u2014 least differentiated state)",
        ]
    for i, ln in enumerate(left_lines):
        p = left_tf.paragraphs[0] if i == 0 else left_tf.add_paragraph()
        r = p.add_run(); r.text = ln; r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(0x33, 0x41, 0x54); r.font.name = "Arial"
    # Add bold labels for key lines
    for idx in [2, 3, 5]:
        run = left_tf.paragraphs[idx].runs[0]
        run.font.bold = True

    # Right card: explain the inference logic (intra-type vs inter-type)
    right_tf = slide.shapes[10].text_frame; right_tf.clear()
    if is_zh:
        right_lines = [
            "核心推理: Intra-type \u03c9 \u226a Inter-type \u03c9 \u2192 迁移假说",
            "",
            "Intra-type: OPC(SVZ) vs OPC(皮层) \u2192 \u03c9\u224822",
            "  \u2192 同一细胞类型, 跨脑区功能分化很低",
            "  \u2192 功能高度相似 \u2192 共同起源",
            "",
            "Inter-type: Astrocyte vs OPC(同皮层内) \u2192 \u03c9\u2248107.5",
            "  \u2192 不同细胞类型, 同一位置功能分化极高",
            "  \u2192 功能截然不同 \u2192 独立谱系",
            "",
            "\u2234 OPC在SVZ和皮层的相似度远超皮层内不同细胞间的相似度",
            "\u2192 OPC自SVZ向皮层迁移, 是解释该数据最简洁的假说",
        ]
    else:
        right_lines = [
            "Key inference: Intra-type \u03c9 \u226a Inter-type \u03c9 \u2192 Migration hypothesis",
            "",
            "Intra-type: OPC(SVZ) vs OPC(Cortex) \u2192 \u03c9\u224822",
            "  \u2192 Same cell type, low divergence across brain regions",
            "  \u2192 High functional similarity \u2192 common origin",
            "",
            "Inter-type: Astrocyte vs OPC (within Cortex) \u2192 \u03c9\u2248107.5",
            "  \u2192 Different cell types, high divergence within same region",
            "  \u2192 Radically different function \u2192 independent lineages",
            "",
            "\u2234 OPC similarity across SVZ and Cortex far exceeds",
            "   inter-cell-type similarity within Cortex alone.",
            "\u2192 OPC migration from SVZ to Cortex is the most parsimonious hypothesis.",
        ]
    for i, ln in enumerate(right_lines):
        p = right_tf.paragraphs[0] if i == 0 else right_tf.add_paragraph()
        r = p.add_run(); r.text = ln; r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(0x33, 0x41, 0x54); r.font.name = "Arial"
    # Bold the key conclusion lines
    for idx_text in right_lines:
        if idx_text.startswith("\u2234") or idx_text.startswith("\u2192") or idx_text.startswith("Key"):
            pass  # will bold the paragraph header
    # Bold first line (title)
    right_tf.paragraphs[0].runs[0].font.bold = True
    right_tf.paragraphs[0].runs[0].font.size = Pt(11)
    right_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
    # Bold intra/inter labels
    for pi in [2, 7]:
        right_tf.paragraphs[pi].runs[0].font.bold = True
    # Bold conclusion
    for pi in range(len(right_lines) - 2, len(right_lines)):
        right_tf.paragraphs[pi].runs[0].font.bold = True

    prs.save(pptx_path)
    print(f"  Saved: {os.path.basename(pptx_path)}")


def main():
    print("=== Step 1: Generate SVG ===")
    svg_content = make_svg()
    svg_path = f"{BASE}/_s17_brain_v2.svg"
    with open(svg_path, 'w', encoding='utf-8') as f:
        f.write(svg_content)
    print(f"  SVG: {len(svg_content):,} chars")

    print("\n=== Step 2: SVG -> PNG ===")
    png_path = f"{BASE}/_s17_brain_v2.png"
    ok = svg_to_png(svg_path, png_path)
    if ok:
        print(f"  PNG: {os.path.getsize(png_path)/1024:.1f} KB")
    else:
        print("  FAILED"); return

    print("\n=== Step 3: EN PPT ===")
    replace_brain_and_text(PPT_EN, png_path, is_zh=False)

    print("\n=== Step 4: ZH PPT ===")
    replace_brain_and_text(PPT_ZH, png_path, is_zh=True)

    print("\n=== DONE ===")


if __name__ == "__main__":
    main()
