"""
notebooks/52_rebalance_s15_17.py

Rebalance Slides 15-17 layout:
  Before: image (left, 5") + cards (right column) → awkward empty space below image
  After:  image (wide top, 8.5") + cards (horizontal row below) → balanced, no dead space

Applies to both EN (CKI_Lecture_2026_v3.pptx) and ZH (CKI_Lecture_2026_v3_ZH.pptx).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE      = Path(r"C:\Users\KnightZ\Desktop\细胞受选择")
PPT_EN    = BASE / "results/figures_final/CKI_Lecture_2026_v3.pptx"
PPT_ZH    = BASE / "results/figures_final/CKI_Lecture_2026_v3_ZH.pptx"
ASSETS    = BASE / "results/figures_final"

# ── Theme ──
TITLE_FG   = (0x1E, 0x29, 0x3B)
BODY_FG    = (0x1E, 0x29, 0x3B)
CARD_BG    = (0xFF, 0xFF, 0xFF)
ACCENT_TEAL= (0x0D, 0x94, 0x88)
ACCENT_BLUE= (0x25, 0x63, 0xEB)
ACCENT_RED = (0xDC, 0x26, 0x26)
DARK_BG    = (0x0F, 0x17, 0x2A)

# ── Layout constants ──
IMG_LEFT   = 0.5
IMG_TOP    = 1.05
IMG_WIDTH  = 8.5    # was 5.0 — stretched wide
IMG_HEIGHT = 2.0    # was 1.92 — slightly taller to fill top half

CARDS_TOP  = 3.25   # below image
CARDS_BOT  = 5.20   # near slide bottom
CARDS_H    = CARDS_BOT - CARDS_TOP  # ≈ 1.95"

# 3-card row (Slide 15): 3 columns with gaps
C3_W = 2.55
C3_GAP = 0.30
C3_X = [IMG_LEFT, IMG_LEFT + C3_W + C3_GAP, IMG_LEFT + 2*(C3_W + C3_GAP)]

# 2-card row (Slides 16-17)
C2_W = 4.10
C2_GAP = 0.30
C2_X = [IMG_LEFT, IMG_LEFT + C2_W + C2_GAP]

# PNG paths (already generated from 51_ script)
PNG_S15 = ASSETS / "_s15_baseline_v3.png"
PNG_S16 = ASSETS / "_s16_anomaly_v3.png"
PNG_S17 = ASSETS / "_s17_migration_v3.png"


# ═════════════════════════════════════════════════════════════════
#  Helpers
# ═════════════════════════════════════════════════════════════════

def remove_all_pics_from_slide(slide):
    spTree = slide.shapes._spTree
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    pics = spTree.findall('.//p:pic', ns)
    for pic in list(pics):
        spTree.remove(pic)
    return len(pics)


def remove_non_title_shapes(slide):
    """Remove shapes except title (T < 1.0") and page number (bottom-right)."""
    to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'image'):
            continue
        keep = False
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt.isdigit() and len(txt) <= 2:
                if shape.top > Inches(4.5) and (shape.left + shape.width) > Inches(8.5):
                    keep = True
            if shape.top < Inches(1.0):
                keep = True
        if not keep:
            to_remove.append(shape)
    spTree = slide.shapes._spTree
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    return len(to_remove)


def add_card_shape(slide, left, top, width, height, title, lines,
                   title_color=None, bg_color=None, font_size=9):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*(bg_color or CARD_BG))
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*(title_color or ACCENT_BLUE))

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(*BODY_FG)

    return shape


def rebuild_slide(slide, png_path, cards, slide_label):
    """Generic: clear old → add wide image top → N cards horizontal row below."""
    n_removed_txt = remove_non_title_shapes(slide)
    n_removed_pic = remove_all_pics_from_slide(slide)
    print(f"    Cleared: {n_removed_txt} text + {n_removed_pic} pic(s)")

    slide.shapes.add_picture(str(png_path),
        Inches(IMG_LEFT), Inches(IMG_TOP),
        Inches(IMG_WIDTH), Inches(IMG_HEIGHT))

    for i, (title, lines, tc, bg) in enumerate(cards):
        x = C3_X[i] if len(cards) == 3 else C2_X[i]
        w = C3_W if len(cards) == 3 else C2_W
        fs = 8 if len(cards) == 3 else 9  # smaller font for 3 narrow columns
        add_card_shape(slide, x, CARDS_TOP, w, CARDS_H,
                       title, lines, title_color=tc, bg_color=bg, font_size=fs)

    print(f"    {slide_label}: image(top, {IMG_WIDTH}x{IMG_HEIGHT}\") + {len(cards)} cards(row) ✓")


# ═════════════════════════════════════════════════════════════════
#  Card content — width-adjusted for new layout
# ═════════════════════════════════════════════════════════════════

# ── EN ──
EN_S15 = [
    ("Intra-region ω", [
        "Same region, different types",
        "Microenvironment limits divergence",
        "Cortex: Astro(28), OPC(22), Neuron(18)"
    ], ACCENT_TEAL, None),
    ("Inter-region ω", [
        "Same type, different regions",
        "Distinct environments → higher ω",
        "OPC: SVZ(3.2)→Hipp(18)→Cortex(22)"
    ], ACCENT_RED, None),
    ("Baseline Rule", [
        "Intra < Inter: consistent across regions",
        "Cells differentiate LOCALLY, not via migration"
    ], DARK_BG, (0xF1, 0xF5, 0xF9)),
]

EN_S16 = [
    ("OPC: Breaking the Baseline", [
        "23.5% of OPC pairs show strong ω anomaly",
        "OPC ω(SVZ) = 3.2, OPC ω(Cortex) = 22",
        "Intra-type ω diff = 18.8 (surprisingly SMALL)",
        "OPCs are known to migrate during development",
        "and in response to injury — CKI recovers this",
        "without prior knowledge of migration pathways"
    ], ACCENT_RED, None),
    ("Astrocytes: High ω, Low Migration", [
        "Highest baseline ω = 107.5 (strong identity)",
        "Yet only 14 migration candidates — very few",
        "Inter-type ω diff = 85.5 (extremely LARGE)",
        "High ω reflects local specialization,",
        "NOT migration — CKI distinguishes both",
        "signals independently"
    ], ACCENT_TEAL, None),
]

EN_S17 = [
    ("Omega Gradient Across Brain Regions", [
        "Cortex Astrocytes: ω = 107.5 (highest identity)",
        "Cerebellum Bergmann Glia: ω = 85",
        "Hipp/Thalamus/Striatum: ω = 38-45",
        "SVZ OPCs: ω = 3.2 (least differentiated)",
        "ω measures functional divergence within a region"
    ], ACCENT_BLUE, None),
    ("OPC Migration Hypothesis", [
        "Intra-type: OPC(SVZ) vs OPC(Cortex) → ω ≈ 22",
        "→ Same type, low divergence across regions",
        "Inter-type: Astrocyte vs OPC (Cortex) → ω ≈ 107.5",
        "→ Different types, high divergence, same region",
        "∴ OPC similarity across regions EXCEEDS",
        "inter-cell-type similarity within region",
        "→ SVZ → Cortex migration is best hypothesis"
    ], ACCENT_RED, None),
]

# ── ZH ──
ZH_S15 = [
    ("同脑区内 ω", [
        "相同脑区，不同细胞类型",
        "微环境限制表达差异 → ω 较低",
        "皮层: Astro(28), OPC(22), Neuron(18)"
    ], ACCENT_TEAL, None),
    ("跨脑区间 ω", [
        "相同细胞类型，不同脑区",
        "不同微环境 → ω 更高",
        "OPC: SVZ(3.2)→Hipp(18)→Cortex(22)"
    ], ACCENT_RED, None),
    ("基线规律", [
        "同脑区 ω < 跨脑区 ω",
        "细胞在局部环境中分化，非迁移"
    ], DARK_BG, (0xF1, 0xF5, 0xF9)),
]

ZH_S16 = [
    ("OPC：打破基线的特例", [
        "OPC配对中23.5%显示强烈ω异常信号",
        "OPC ω(SVZ) = 3.2, OPC ω(Cortex) = 22",
        "同类ω差异 = 18.8（异常地小）",
        "OPC已知在发育过程中活跃迁移",
        "并在损伤后应答中迁移 — CKI在无",
        "先验知识下独立发现这一现象"
    ], ACCENT_RED, None),
    ("星形胶质细胞：高 ω，低迁移", [
        "最高基线 ω = 107.5（强细胞类型身份）",
        "但仅14个迁移候选 — 数量极少",
        "跨类ω差异 = 85.5（异常地大）",
        "高ω反映局部特异性分化",
        "而非迁移 — CKI独立区分两信号"
    ], ACCENT_TEAL, None),
]

ZH_S17 = [
    ("各脑区 ω 梯度分布", [
        "皮层 Astrocytes: ω = 107.5（最高特异性）",
        "小脑 Bergmann Glia: ω = 85",
        "海马/丘脑/纹状体: ω = 38-45",
        "SVZ OPCs: ω = 3.2（分化程度最低）",
        "ω 衡量特定脑区内功能分化程度"
    ], ACCENT_BLUE, None),
    ("OPC 迁移假说", [
        "同类比较：OPC(SVZ) vs OPC(Cortex) → ω≈22",
        "→ 相同类型，跨脑区差异很小",
        "跨类比较：Astrocyte vs OPC（同皮层）→ ω≈107.5",
        "→ 不同类型，同脑区内差异极大",
        "∴ OPC跨脑区相似性远超",
        "皮层内部不同类型间相似性",
        "→ SVZ→皮层 OPC迁移是最合理假说"
    ], ACCENT_RED, None),
]


# ═════════════════════════════════════════════════════════════════
#  MAIN
# ═════════════════════════════════════════════════════════════════

def main():
    for ppt_path, label, s15, s16, s17 in [
        (PPT_EN, "EN", EN_S15, EN_S16, EN_S17),
        (PPT_ZH, "ZH", ZH_S15, ZH_S16, ZH_S17),
    ]:
        print(f"\n{'='*60}")
        print(f"Processing {label}: {ppt_path.name}")
        print(f"{'='*60}")

        prs = Presentation(str(ppt_path))

        print(f"  Slide 15:")
        rebuild_slide(prs.slides[14], PNG_S15, s15, f"S15 {label}")

        print(f"  Slide 16:")
        rebuild_slide(prs.slides[15], PNG_S16, s16, f"S16 {label}")

        print(f"  Slide 17:")
        rebuild_slide(prs.slides[16], PNG_S17, s17, f"S17 {label}")

        prs.save(str(ppt_path))
        print(f"  Saved: {ppt_path.name}")

    print(f"\n{'='*60}")
    print("ALL DONE. New layout for S15-S17:")
    print(f"  Image: top, {IMG_WIDTH}\" wide × {IMG_HEIGHT}\" high — fills width, no dead space")
    print(f"  Cards: horizontal row below, T={CARDS_TOP}\", H={CARDS_H:.2f}\"")
    print(f"  S15: 3 columns ({C3_W}\" each)")
    print(f"  S16/S17: 2 columns ({C2_W}\" each)")
    print(f"  EN: {PPT_EN}")
    print(f"  ZH: {PPT_ZH}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
