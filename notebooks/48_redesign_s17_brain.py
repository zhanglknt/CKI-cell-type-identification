"""
Redesign Slide 17: Brain Regional Specialization & OPC Migration Map.

Design v2: taller brain outline, larger text, clearer migration arrow.
Target image area: 8138160×1402080 EMU (8.9" × 1.53") in PPT
Output PNG at 2x for crispness: 1780×306 px (200 DPI equivalent)

v2.0 - improved layout, bigger fonts, better proportions
"""

import os
import subprocess
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor


BASE = r"C:\Users\KnightZ\Desktop\细胞受选择\results\figures_final"
PPT_EN = f"{BASE}/CKI_Lecture_2026_v3.pptx"
PPT_ZH = f"{BASE}/CKI_Lecture_2026_v3_ZH.pptx"

# Target dimensions — match PPT area ratio ~5.82:1
PNG_W, PNG_H = 2700, 464


def omega_color(omega):
    """Map omega value to hue from blue(222°) to red(0°), logarithmic scale."""
    import math
    lo, hi = 3, 110
    t = (math.log(max(omega, lo)) - math.log(lo)) / (math.log(hi) - math.log(lo))
    t = max(0, min(1, t))
    h = 222 - 220 * t  # blue → red
    s, l = 72, 48
    return hsl2hex(h, s, l)


def hsl2hex(h, s, l):
    s /= 100; l /= 100
    c = (1 - abs(2*l - 1)) * s
    x = c * (1 - abs((h/60) % 2 - 1))
    m = l - c/2
    if h < 60:    r, g, b = c, x, 0
    elif h < 120: r, g, b = x, c, 0
    elif h < 180: r, g, b = 0, c, x
    elif h < 240: r, g, b = 0, x, c
    elif h < 300: r, g, b = x, 0, c
    else:         r, g, b = c, 0, x
    r, g, b = int((r+m)*255), int((g+m)*255), int((b+m)*255)
    return f"#{r:02X}{g:02X}{b:02X}"


def make_svg():
    """Generate the brain schematic SVG."""
    out = []
    out.append(f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {PNG_W} {PNG_H}">')
    out.append('<defs>')
    out.append('<linearGradient id="brainGrad" x1="0" y1="0" x2="1" y2="0">')
    out.append('  <stop offset="0%" stop-color="#E8EDF4"/>')
    out.append('  <stop offset="100%" stop-color="#F5F7FA"/>')
    out.append('</linearGradient>')
    # Glow for migration arrow
    out.append('<filter id="glow" x="-20%" y="-20%" width="140%" height="140%">')
    out.append('  <feGaussianBlur stdDeviation="6" result="blur"/>')
    out.append('  <feMerge><feMergeNode in="blur"/><feMergeNode in="SourceGraphic"/></feMerge>')
    out.append('</filter>')
    out.append('<filter id="shadow" x="-10%" y="-10%" width="120%" height="120%">')
    out.append('  <feDropShadow dx="1.5" dy="1.5" stdDeviation="2" flood-opacity="0.18"/>')
    out.append('</filter>')
    # Arrow markers
    out.append('<marker id="arrowRed" markerWidth="14" markerHeight="9" refX="12" refY="4.5" orient="auto" markerUnits="userSpaceOnUse">')
    out.append('  <path d="M 0,0 L 14,4.5 L 0,9 z" fill="#DC2626"/>')
    out.append('</marker>')
    out.append('</defs>')

    # Background
    out.append(f'<rect width="{PNG_W}" height="{PNG_H}" fill="#F8FAFC" rx="10"/>')

    # ================================================================
    # BRAIN OUTLINE (mouse sagittal view) — use full vertical space
    # Brain occupies roughly (60,15) → (2640,445)
    # ================================================================
    brain_path = (
        'M 70,280 '
        'C 70,180 110,95 210,65 '       # olfactory → frontal
        'C 380,28 720,12 1300,20 '      # frontal cortex top
        'C 1700,26 2050,40 2320,58 '    # posterior cortex
        'C 2480,78 2590,118 2630,165 ' # cerebellum top
        'C 2660,205 2665,245 2645,290 ' # cerebellum back
        'C 2600,350 2500,395 2340,420 ' # brainstem bottom
        'C 2120,450 1750,458 1350,458 ' # ventral midline
        'C 1000,458 700,450 400,430 '   # ventral anterior
        'C 230,408 120,385 75,340 '     # olfactory base
        'Z'
    )
    out.append(f'<path d="{brain_path}" fill="url(#brainGrad)" stroke="#94A3B8" stroke-width="2.2"/>')

    # Subtle internal structures
    # Corpus callosum
    out.append('<path d="M 320,88 C 600,62 1050,56 1600,68" '
               'fill="none" stroke="#CBD5E1" stroke-width="1" stroke-dasharray="6,5" opacity="0.6"/>')
    # Lateral ventricle hint
    out.append('<path d="M 350,195 C 550,172 850,162 1150,168 C 1380,173 1580,178 1720,188" '
               'fill="none" stroke="#CBD5E1" stroke-width="1.3" stroke-dasharray="6,5" opacity="0.55"/>')
    # Third ventricle hint
    out.append('<line x1="1120" y1="255" x2="1140" y2="355" '
               'stroke="#CBD5E1" stroke-width="1" stroke-dasharray="4,5" opacity="0.5"/>')

    # ================================================================
    # REGION DATA
    # (name, cx, cy, radius, omega, label_pos_x, label_pos_y, cell_detail)
    # ================================================================
    REGIONS = [
        ("Cortex",     1850, 108, 62, 107.5,
         "top",    1850, 25,
         [("Astrocytes ω=107.5", "#DC2626"), ("OPCs ω=22", "#D97706")]),
        ("SVZ",         980, 248, 36, 3.2,
         "bottom", 980, 330,
         [("OPCs ω=3.2", "#2563EB")]),
        ("Hippocampus",1600, 162, 42, 45,
         "right",   1710, 155,
         [("Neurons ω=45", "#A855F7")]),
        ("Striatum",   1120, 310, 36, 38,
         "bottom",  1120, 390,
         [("Neurons ω=38", "#8B5CF6")]),
        ("Thalamus",   1320, 278, 36, 42,
         "bottom",  1320, 360,
         [("Neurons ω=42", "#A855F7")]),
        ("Cerebellum", 2450, 268, 52, 85,
         "top",    2450, 190,
         [("Bergmann Glia ω=85", "#EF4444")]),
    ]

    # Draw region circles + labels
    for name, cx, cy, r, omega, align, lx, ly, details in REGIONS:
        color = omega_color(omega)

        # Outer ring (subtle)
        out.append(f'<circle cx="{cx}" cy="{cy}" r="{r+5}" fill="none" '
                   f'stroke="{color}" stroke-width="2.5" opacity="0.25"/>')
        # Main circle with shadow
        out.append(f'<circle cx="{cx}" cy="{cy}" r="{r}" fill="{color}" '
                   f'fill-opacity="0.12" stroke="{color}" stroke-width="2.2" filter="url(#shadow)"/>')

        # Omega value inside circle
        omega_str = f"ω={omega}"
        fs = 17 if omega < 100 else 15
        out.append(f'<text x="{cx}" y="{cy+1}" text-anchor="middle" dominant-baseline="middle" '
                   f'font-family="Arial,Helvetica,sans-serif" font-size="{fs}" '
                   f'font-weight="800" fill="{color}">{omega_str}</text>')

        # Connector line to region name
        if align == "top":
            conn_y1 = cy - r - 5
            conn_y2 = ly + 14 if ly > cy else ly + 14
            out.append(f'<line x1="{cx}" y1="{conn_y1}" x2="{lx}" y2="{conn_y2-2}" '
                       f'stroke="{color}" stroke-width="1.2" stroke-dasharray="4,4" opacity="0.5"/>')
            anchor = "middle"; dy_adj = 0
        elif align == "bottom":
            conn_y1 = cy + r + 5
            conn_y2 = ly - 14
            out.append(f'<line x1="{cx}" y1="{conn_y1}" x2="{lx}" y2="{conn_y2+2}" '
                       f'stroke="{color}" stroke-width="1.2" stroke-dasharray="4,4" opacity="0.5"/>')
            anchor = "middle"; dy_adj = 0
        else:  # right/left
            if lx > cx:  # right side
                conn_x1 = cx + r + 5; conn_x2 = lx - 2
                out.append(f'<line x1="{conn_x1}" y1="{cy}" x2="{conn_x2}" y2="{ly}" '
                           f'stroke="{color}" stroke-width="1.2" stroke-dasharray="4,4" opacity="0.5"/>')
                anchor = "start"
            else:
                anchor = "end"
            dy_adj = 4

        # Region name
        out.append(f'<text x="{lx}" y="{ly+dy_adj}" text-anchor="{anchor}" '
                   f'font-family="Arial,Helvetica,sans-serif" font-size="16" '
                   f'font-weight="700" fill="{color}">{name}</text>')

        # Cell type detail annotations
        for di, (dtxt, dcol) in enumerate(details):
            if align == "top":
                dx = cx; dy = ly + 30 + di * 18
                da = "middle"
            elif align == "bottom":
                dx = cx; dy = ly + 30 + di * 18
                da = "middle"
            else:  # right
                dx = lx + 8; dy = cy + (di - len(details)/2 + 0.5) * 16
                da = "start"

            # Only show if within canvas bounds roughly
            if 0 < dy < PNG_H - 5 and 0 < dx < PNG_W:
                out.append(f'<text x="{dx}" y="{dy}" text-anchor="{da}" '
                           f'font-family="Arial,Helvetica,sans-serif" font-size="11.5" '
                           f'fill="{dcol}" font-style="italic">{dtxt}</text>')

    # ================================================================
    # OPC MIGRATION ARROW: SVZ (980,248) → Cortex (1850,108)
    # ================================================================
    svz_x, svz_y = 980, 248
    ctx_x, ctx_y = 1850, 108

    # Curved path going upward through the middle of brain
    ctrl1_x = svz_x + 150
    ctrl1_y = 140   # pull curve up
    ctrl2_x = ctx_x - 200
    ctrl2_y = 50    # pull curve up near cortex

    # Main thick migration line
    mig_path = (f'M {svz_x+28},{svz_y-22} '
                f'C {ctrl1_x},{ctrl1_y} {ctrl2_x},{ctrl2_y} {ctx_x-40},{ctx_y-32}')
    out.append(f'<path d="{mig_path}" fill="none" stroke="#DC2626" stroke-width="5" '
               f'marker-end="url(#arrowRed)" filter="url(#glow)" opacity="0.92"/>')

    # Parallel thinner line for visual weight
    mig_path2 = (f'M {svz_x+28},{svz_y-10} '
                 f'C {ctrl1_x+30},{ctrl1_y+30} {ctrl2_x-30},{ctrl2_y+30} {ctx_x-40},{ctx_y-20}')
    out.append(f'<path d="{mig_path2}" fill="none" stroke="#EF4444" stroke-width="2.5" '
               f'opacity="0.5" stroke-linecap="round"/>')

    # Migration label pill (centered on curve)
    lbl_x = (svz_x + ctx_x) / 2 + 80
    lbl_y = (svz_y + ctx_y) / 2 - 90
    pill_w, pill_h = 148, 30
    out.append(f'<rect x="{lbl_x - pill_w/2:.0f}" y="{lbl_y - pill_h/2:.0f}" '
               f'width="{pill_w}" height="{pill_h}" rx="15" '
               f'fill="#FEF2F2" stroke="#FCA5A5" stroke-width="1.5"/>')
    out.append(f'<text x="{lbl_x}" y="{lbl_y + 1}" text-anchor="middle" dominant-baseline="middle" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="14" '
               f'font-weight="800" fill="#DC2626">OPC Migration</text>')

    # Small "omega gradient" annotation on the arrow
    grad_txt_x = lbl_x + 110
    grad_txt_y = lbl_y - 5
    out.append(f'<text x="{grad_txt_x}" y="{grad_txt_y}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="10" '
               f'fill="#94A3B8" font-style="italic">\u03c9 3\u219222</text>')

    # ================================================================
    # LEGEND (bottom-left)
    # ================================================================
    lg_x, lg_y = 70, 440
    out.append(f'<text x="{lg_x}" y="{lg_y+3}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="11" '
               f'fill="#64748B" font-weight="600">Low \u03c9</text>')

    # Color bar
    bar_x = lg_x + 50
    bar_w = 120
    for i in range(bar_w):
        t = i / (bar_w - 1)
        color = omega_color(3 + t * (110 - 3))
        rx = bar_x + i * 1
        out.append(f'<rect x="{rx}" y="{lg_y-4}" width="1" height="12" fill="{color}" />')

    out.append(f'<text x="{lg_x + bar_x + bar_w + 12}" y="{lg_y+3}" text-anchor="start" '
               f'font-family="Arial,Helvetica,sans-serif" font-size="11" '
               f'fill="#64748B" font-weight="600">High \u03c9</text>')

    out.append('</svg>')
    return '\n'.join(out)


def svg_to_png(svg_path, png_path):
    """Convert SVG to high-res PNG via Chrome headless."""
    html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0}}body{{
display:flex;justify-content:center;align-items:center;
background:#F8FAFC;width:{PNG_W}px;height:{PNG_H}px;}}
svg{{width:{PNG_W}px;height:{PNG_H}px;display:block;}}</style></head>
<body>
{open(svg_path, encoding='utf-8').read()}
</body></html>"""

    html_path = svg_path.replace('.svg', '.html')
    with open(html_path, 'w', encoding='utf-8') as f:
        f.write(html)

    chrome_paths = [
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    chrome = None
    for p in chrome_paths:
        if os.path.exists(p):
            chrome = p
            break

    if not chrome:
        raise RuntimeError("No browser found")

    file_url = f"file:///{html_path.replace(chr(92), '/')}"
    cmd = [
        chrome, '--headless=new', '--disable-gpu', '--no-sandbox',
        f'--window-size={PNG_W},{PNG_H}',
        f'--screenshot={png_path}', file_url
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
    if result.returncode != 0:
        print(f"  Browser stderr: {result.stderr[:500]}")

    os.remove(html_path)

    if os.path.exists(png_path) and os.path.getsize(png_path) > 5000:
        return True
    return False


def replace_brain_image_and_text(pptx_path, png_path):
    """Replace brain pic + update card text on Slide 17."""
    prs = Presentation(pptx_path)
    slide = prs.slides[16]

    old_pic = slide.shapes[11]
    left, top, width, height = old_pic.left, old_pic.top, old_pic.width, old_pic.height
    old_pic._element.getparent().remove(old_pic._element)
    slide.shapes.add_picture(png_path, left, top, width, height)

    # Update left card (shape [6])
    left_tf = slide.shapes[6].text_frame
    left_tf.clear()
    left_lines = [
        "Astrocytes in cortex show the highest omega (107.5) — most distinct glial identity.",
        "Cerebellar Bergmann glia form a unique cluster (omega = 85).",
        "Thalamus (omega = 42), striatum (omega = 38), hippocampus (omega = 45): intermediate neuronal signatures.",
    ]
    for i, ln in enumerate(left_lines):
        p = left_tf.paragraphs[0] if i == 0 else left_tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(10); run.font.color.rgb = RGBColor(0x33, 0x41, 0x54)
        run.font.name = "Arial"

    # Update right card (shape [10])
    right_tf = slide.shapes[10].text_frame
    right_tf.clear()
    right_lines = [
        "OPCs originate in SVZ (omega = 3.2), then migrate radially toward cortex.",
        "By reaching cortex, OPCs mature to omega = 22 — a 7-fold functional differentiation increase.",
        "CKI captures this spatial trajectory as an omega gradient along the migration path.",
        "This gradient generates testable hypotheses linking location to differentiation state.",
    ]
    for i, ln in enumerate(right_lines):
        p = right_tf.paragraphs[0] if i == 0 else right_tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(10); run.font.color.rgb = RGBColor(0x33, 0x41, 0x54)
        run.font.name = "Arial"

    prs.save(pptx_path)
    print(f"  Saved: {os.path.basename(pptx_path)}")


def main():
    import shutil

    print("=== Step 1: Generate SVG === ")
    svg_content = make_svg()
    svg_path = f"{BASE}/_s17_brain_new.svg"
    with open(svg_path, 'w', encoding='utf-8') as f:
        f.write(svg_content)
    print(f"  SVG: {len(svg_content)} chars")

    print("\n=== Step 2: Convert SVG → PNG ===")
    png_path = f"{BASE}/_s17_brain_new.png"
    ok = svg_to_png(svg_path, png_path)
    if ok:
        print(f"  PNG: {os.path.getsize(png_path)/1024:.1f} KB")
    else:
        print("  FAILED"); return

    print("\n=== Step 3: EN PPT ===")
    replace_brain_image_and_text(PPT_EN, png_path)

    print("\n=== Step 4: ZH PPT ===")
    # ZH version needs Chinese text updates
    prs_zh = Presentation(PPT_ZH)
    slide_zh = prs_zh.slides[16]
    old_pic = slide_zh.shapes[11]
    left, top, w, h = old_pic.left, old_pic.top, old_pic.width, old_pic.height
    old_pic._element.getparent().remove(old_pic._element)
    slide_zh.shapes.add_picture(png_path, left, top, w, h)

    # ZH left card
    left_tf = slide_zh.shapes[6].text_frame; left_tf.clear()
    zh_left = [
        "皮层星形胶质细胞omega最高(107.5)，具有最独特的胶质身份。",
        "小脑Bergmann胶质细胞形成独立聚类(omega = 85)。",
        "丘脑(omega = 42)、纹状体(omega = 38)、海马(omega = 45)：中等神经元特征。",
    ]
    for i, ln in enumerate(zh_left):
        p = left_tf.paragraphs[0] if i==0 else left_tf.add_paragraph()
        r = p.add_run(); r.text = ln; r.font.size=Pt(10)
        r.font.color.rgb = RGBColor(0x33,0x41,0x54); r.font.name="Arial"

    # ZH right card
    right_tf = slide_zh.shapes[10].text_frame; right_tf.clear()
    zh_right = [
        "OPC起源于室管膜下区(SVZ, omega=3.2)，然后向皮层迁移辐射。",
        "到达皮层后，OPC成熟至omega = 22——功能性分化增加7倍。",
        "CKI将此空间轨迹捕获为沿迁移路径的omega梯度。",
        "该梯度生成可验证的假说，将位置与分化状态关联起来。",
    ]
    for i, ln in enumerate(zh_right):
        p = right_tf.paragraphs[0] if i==0 else right_tf.add_paragraph()
        r = p.add_run(); r.text = ln; r.font.size=Pt(10)
        r.font.color.rgb = RGBColor(0x33,0x41,0x54); r.font.name="Arial"

    prs_zh.save(PPT_ZH)
    print(f"  Saved ZH: {os.path.basename(PPT_ZH)}")

    # Clean up temp files
    for ext in ['_s17fix.pptx', '_s17fix_zh.pptx']:
        fp = f"{BASE}/{ext}"
        if os.path.exists(fp):
            try: os.remove(fp)
            except: pass

    print("\n=== DONE ===")


if __name__ == "__main__":
    main()
