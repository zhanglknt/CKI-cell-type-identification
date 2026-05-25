"""
notebooks/54_gene_sets_and_rebalance.py

Two tasks:
  A) Rebalance S15-17: increase plot height 2.0→2.5″, decrease card height 1.95→1.50″
  B) Insert new S7: Automated Gene Set Detection (HK + Functional genes)

Applies to both EN and ZH PPTs.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

BASE   = Path(r"C:\Users\KnightZ\Desktop\细胞受选择")
ASSETS = BASE / "results/figures_final"
PPT_EN = ASSETS / "CKI_Lecture_2026_v3.pptx"
PPT_ZH = ASSETS / "CKI_Lecture_2026_v3_ZH.pptx"

# ── Theme ──
TITLE_FG  = (0x1E, 0x29, 0x3B)
BODY_FG   = (0x1E, 0x29, 0x3B)
CARD_BG   = (0xFF, 0xFF, 0xFF)
TEAL      = (0x0D, 0x94, 0x88)
BLUE      = (0x25, 0x63, 0xEB)
RED       = (0xDC, 0x26, 0x26)
PURPLE    = (0x7C, 0x3A, 0xED)
GREEN     = (0x05, 0x96, 0x69)
AMBER     = (0xD9, 0x77, 0x06)
DARK_BG   = (0x0F, 0x17, 0x2A)
MUTED     = (0x64, 0x74, 0x8B)
WHITE     = (0xFF, 0xFF, 0xFF)
LIGHT_SLATE = (0xF1, 0xF5, 0xF9)

# ═══════════════════════════════════════════════════════
#  PART A: Rebalance S15-17
# ═══════════════════════════════════════════════════════

# New layout
IMG_LEFT, IMG_TOP = 0.5, 1.05
IMG_WIDTH, IMG_HEIGHT = 8.5, 2.5    # was 2.0

CARDS_TOP = 3.70                      # was 3.25
CARDS_H   = 1.50                      # was 1.95

C3_W, C3_GAP = 2.55, 0.30
C3_X = [IMG_LEFT, IMG_LEFT + C3_W + C3_GAP, IMG_LEFT + 2*(C3_W + C3_GAP)]
C2_W, C2_GAP = 4.10, 0.30
C2_X = [IMG_LEFT, IMG_LEFT + C2_W + C2_GAP]

PNG_S15 = ASSETS / "_s15_wide.png"
PNG_S16 = ASSETS / "_s16_wide.png"
PNG_S17 = ASSETS / "_s17_wide.png"


def remove_all_pics(slide):
    spTree = slide.shapes._spTree
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    for pic in list(spTree.findall('.//p:pic', ns)):
        spTree.remove(pic)


def remove_non_titles(slide):
    """Keep: title (T<1.0″) and page number (T>4.5″, R>8.5″). Remove rest."""
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        keep = False
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if shape.top < Inches(1.0):
                keep = True
            elif txt.isdigit() and len(txt) <= 2 and shape.top > Inches(4.5) and (shape.left + shape.width) > Inches(8.5):
                keep = True
        if not keep:
            spTree.remove(shape._element)


def add_card(slide, left, top, width, height, title, lines,
             title_color=None, bg_color=None, font_size=8):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*(bg_color or CARD_BG))
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*(title_color or BLUE))

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(*BODY_FG)


def rebuild_slide(slide, png_path, cards, label):
    remove_non_titles(slide)
    remove_all_pics(slide)
    slide.shapes.add_picture(str(png_path),
        Inches(IMG_LEFT), Inches(IMG_TOP),
        Inches(IMG_WIDTH), Inches(IMG_HEIGHT))
    for i, (title, lines, tc, bg) in enumerate(cards):
        x = C3_X[i] if len(cards) == 3 else C2_X[i]
        w = C3_W if len(cards) == 3 else C2_W
        fs = 7 if len(cards) == 3 else 8
        add_card(slide, x, CARDS_TOP, w, CARDS_H, title, lines, tc, bg, fs)
    print(f"    {label}: plot {IMG_WIDTH}×{IMG_HEIGHT}″ + {len(cards)} cards ✓")


# Card content (same as before, adjusted for cards using fewer lines as needed)
# ── EN ──
EN_S15 = [
    ("Intra-region ω", [
        "Same region, different types",
        "Microenvironment limits divergence",
        "Cortex: Astro(28), OPC(22), Neuron(18)"
    ], TEAL, None),
    ("Inter-region ω", [
        "Same type, different regions",
        "Distinct environments → higher ω",
        "OPC: SVZ(3.2)→Hipp(18)→Cortex(22)"
    ], RED, None),
    ("Baseline Rule", [
        "Intra < Inter across regions",
        "Cells differentiate LOCALLY"
    ], DARK_BG, LIGHT_SLATE),
]
EN_S16 = [
    ("OPC: Breaking the Baseline", [
        "23.5% OPC pairs show strong ω anomaly",
        "OPC ω(SVZ)=3.2, OPC ω(Cortex)=22",
        "Intra-type Δω=18.8 (surprisingly SMALL)",
        "OPCs migrate during development/injury",
        "CKI recovers this without prior knowledge"
    ], RED, None),
    ("Astrocytes: High ω, Low Migration", [
        "Highest baseline ω=107.5 (strong identity)",
        "Only 14 migration candidates",
        "Inter-type Δω=85.5 (extremely LARGE)",
        "High ω = local specialization, NOT migration",
        "CKI distinguishes both signals independently"
    ], TEAL, None),
]
EN_S17 = [
    ("Omega Gradient Across Brain", [
        "Cortex Astrocytes: ω=107.5 (highest)",
        "Cerebellum Bergmann Glia: ω=85",
        "Hipp/Thal/Striatum: ω=38-45",
        "SVZ OPCs: ω=3.2 (least differentiated)",
        "ω measures functional divergence within region"
    ], BLUE, None),
    ("OPC Migration Hypothesis", [
        "Intra-type: OPC(SVZ) vs OPC(Cortex) → ω≈22",
        "→ Same type, low divergence across regions",
        "Inter-type: Astrocyte vs OPC(Cortex) → ω≈107.5",
        "→ Different types, high divergence, same region",
        "OPC cross-region similarity > cross-type similarity",
        "→ SVZ → Cortex migration is best hypothesis"
    ], RED, None),
]

# ── ZH ──
ZH_S15 = [
    ("同脑区内 ω", [
        "相同脑区，不同细胞类型",
        "微环境限制表达差异 → ω 较低",
        "皮层: Astro(28), OPC(22), Neuron(18)"
    ], TEAL, None),
    ("跨脑区间 ω", [
        "相同细胞类型，不同脑区",
        "不同微环境 → ω 更高",
        "OPC: SVZ(3.2)→Hipp(18)→Cortex(22)"
    ], RED, None),
    ("基线规律", [
        "同脑区 ω < 跨脑区 ω",
        "细胞在局部环境中分化，非迁移"
    ], DARK_BG, LIGHT_SLATE),
]
ZH_S16 = [
    ("OPC：打破基线的特例", [
        "23.5% OPC配对显示强烈ω异常",
        "OPC ω(SVZ)=3.2, OPC ω(Cortex)=22",
        "同类Δω=18.8（异常地小）",
        "OPC已知在发育中活跃迁移",
        "CKI在无先验知识下独立发现"
    ], RED, None),
    ("星形胶质细胞：高ω 低迁移", [
        "最高基线 ω=107.5（强细胞身份）",
        "仅14个迁移候选 — 数量极少",
        "跨类Δω=85.5（异常地大）",
        "高ω反映局部特异性分化，非迁移",
        "CKI独立区分两信号"
    ], TEAL, None),
]
ZH_S17 = [
    ("各脑区 ω 梯度分布", [
        "皮层 Astrocytes: ω=107.5（最高特异性）",
        "小脑 Bergmann Glia: ω=85",
        "海马/丘脑/纹状体: ω=38-45",
        "SVZ OPCs: ω=3.2（分化程度最低）",
        "ω 衡量特定脑区内功能分化程度"
    ], BLUE, None),
    ("OPC 迁移假说", [
        "同类比较：OPC(SVZ) vs OPC(Cortex) → ω≈22",
        "→ 相同类型，跨脑区差异很小",
        "跨类比较：Astrocyte vs OPC（同皮层）→ ω≈107.5",
        "→ 不同类型，同脑区内差异极大",
        "OPC跨脑区相似性 ＞ 皮层内部不同类型间相似性",
        "→ SVZ→皮层 OPC迁移是最合理假说"
    ], RED, None),
]


# ═══════════════════════════════════════════════════════
#  PART B: Gene Sets Slide
# ═══════════════════════════════════════════════════════

def create_gene_sets_slide(prs, lang="EN"):
    """Create a new slide describing HK + Functional gene set automation."""
    # Use default layout (only layout available)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # ── Title ──
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9.0), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    if lang == "EN":
        run.text = "Automated Gene Set Detection"
    else:
        run.text = "基因集自动鉴定"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TITLE_FG)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    if lang == "EN":
        run2.text = "No manual gene list needed — fully data-driven"
    else:
        run2.text = "无需手动指定基因列表 — 完全数据驱动"
    run2.font.size = Pt(12)
    run2.font.color.rgb = RGBColor(*MUTED)

    # ── Two large cards ──
    card_top   = 1.20
    card_h     = 4.30
    card_w     = 4.20
    left_x     = 0.50
    right_x    = 5.00
    card_gap   = 0.30

    # Card title bar height
    bar_h = 0.45

    # ── LEFT CARD: Housekeeping Genes ──
    if lang == "EN":
        hk_title = "Housekeeping Genes"
        hk_sub   = "Stably expressed · Low variation · Across cell types"
        hk_items = [
            ("detection_rate", "Expressed in > 90% of cells\nDefault method (robust to noise)", GREEN),
            ("cv", "Coefficient of variation in lowest 30%\nCaptures expression stability", BLUE),
            ("combined", "Both criteria simultaneously\nHighest stringency (default)", PURPLE),
            ("HRT Atlas", "Built-in reference from 16 human tissues\nUnion with detected set for completeness", AMBER),
        ]
        hk_bottom = "→ Fully automatic; defaults: combined + HRT Atlas reference"
    else:
        hk_title = "持家基因鉴定"
        hk_sub   = "稳定表达 · 低变异度 · 跨细胞类型"
        hk_items = [
            ("detection_rate", ">90% 细胞中检出\n默认方法（对噪声稳健）", GREEN),
            ("cv", "表达变异系数最低 30%\n捕捉表达稳定性", BLUE),
            ("combined", "同时满足以上两条\n最高严格度（默认）", PURPLE),
            ("HRT Atlas", "内置 16 种人组织参考集\n与检出集取并集，确保完整", AMBER),
        ]
        hk_bottom = "→ 全自动；默认：combined + HRT Atlas 参考集"

    add_gene_card(slide, left_x, card_top, card_w, card_h, bar_h,
                  hk_title, hk_sub, hk_items, hk_bottom, TEAL)

    # ── RIGHT CARD: Functional Genes ──
    if lang == "EN":
        fg_title = "Functional / Variable Genes"
        fg_sub   = "Highly variable · Distinguish cell types"
        fg_items = [
            ("hvg (default)", "Scanpy highly variable genes\nTop 2,000 · Exclude HK genes\nCaptures biological variability", RED),
            ("markers", "Per-cluster differential expression\nTop marker genes per cluster\nAlternative when clusters known", BLUE),
        ]
        fg_bottom = "→ Auto-detected; defaults: hvg top 2,000 (excl. HK)"
    else:
        fg_title = "功能 / 可变基因"
        fg_sub   = "高变异 · 区分细胞类型"
        fg_items = [
            ("hvg（默认）", "Scanpy 高变异基因\nTop 2,000 · 排除持家基因\n捕捉生物学变异", RED),
            ("markers", "逐 cluster 差异表达\n每 cluster 取 top marker\n已知 cluster 时的替代方案", BLUE),
        ]
        fg_bottom = "→ 自动鉴定；默认：hvg top 2,000（排除持家基因）"

    add_gene_card(slide, right_x, card_top, card_w, card_h, bar_h,
                  fg_title, fg_sub, fg_items, fg_bottom, RED)

    # ── Bottom section: Workflow summary ──
    if lang == "EN":
        summary = (
            "Input (adata + species)  →  Auto-detect HK genes  →  Auto-detect functional genes  "
            "→  Compute CKI (Kn, Kf, ω)  —  All automated in cki.compute()"
        )
    else:
        summary = (
            "输入 (adata + species)  →  自动鉴定持家基因  →  自动鉴定功能基因  "
            "→  计算 CKI (Kn, Kf, ω)  —  cki.compute() 全自动完成"
        )

    # Bottom bar
    bar = slide.shapes.add_shape(
        1, Inches(0.5), Inches(5.70), Inches(9.0), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*DARK_BG)
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = summary
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*WHITE)

    # ── Page number ──
    return slide


def add_gene_card(slide, x, y, w, h, bar_h, title, subtitle, items, bottom_note, accent_color):
    """Create a card with colored title bar, content rows, and bottom note."""
    # Card background
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*CARD_BG)
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(1.5)
    card.shadow.inherit = False

    # Title bar
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*accent_color)
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*WHITE)

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + bar_h + 0.08), Inches(w - 0.30), Inches(0.35))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(*MUTED)
    run.font.italic = True

    # Content items — each is a small row with colored label + description
    item_y_start = y + bar_h + 0.50
    item_h = 0.72

    for i, (label, desc, color) in enumerate(items):
        iy = item_y_start + i * (item_h + 0.06)

        # Color indicator dot
        dot = slide.shapes.add_shape(
            9, Inches(x + 0.15), Inches(iy + 0.05),
            Inches(0.18), Inches(0.18))  # oval
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(*color)
        dot.line.fill.background()

        # Label
        lb = slide.shapes.add_textbox(
            Inches(x + 0.40), Inches(iy), Inches(w - 0.55), Inches(0.22))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*color)

        # Description
        db = slide.shapes.add_textbox(
            Inches(x + 0.40), Inches(iy + 0.22), Inches(w - 0.55), Inches(0.45))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.space_before = Pt(0)
        run = p.add_run()
        desc_short = desc.split("\n")[0]
        run.text = desc_short
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(*BODY_FG)
        # Second line if exists
        if "\n" in desc:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(0)
            run2 = p2.add_run()
            run2.text = desc.split("\n", 1)[1]
            run2.font.size = Pt(8)
            run2.font.color.rgb = RGBColor(*MUTED)

    # Bottom note
    note_y = y + h - 0.45
    note = slide.shapes.add_shape(
        1, Inches(x + 0.10), Inches(note_y), Inches(w - 0.20), Inches(0.35))
    note.fill.solid()
    if accent_color == TEAL:
        note.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    else:
        note.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2)
    note.line.fill.background()
    note.shadow.inherit = False

    tf = note.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = bottom_note
    run.font.size = Pt(7.5)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*MUTED)


# ═══════════════════════════════════════════════════════
#  Slide insertion helper
# ═══════════════════════════════════════════════════════

def insert_slide_at(prs, slide_idx, target_pos):
    """
    Move a slide (by its index in the current list) to target_pos (0-indexed).
    Manipulates the <p:sldIdLst> XML directly.
    """
    # Access the sldIdLst in presentation.xml
    pres_elem = prs.part._element
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
          'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

    sldIdLst = pres_elem.find('.//p:sldIdLst', ns)
    if sldIdLst is None:
        raise ValueError("sldIdLst not found")

    # Get all sldId elements
    sld_ids = list(sldIdLst.findall('p:sldId', ns))

    # Move the element at slide_idx to target_pos
    elem = sld_ids.pop(slide_idx)
    sldIdLst.insert(target_pos, elem)

    print(f"    Slide moved from position {slide_idx+1} → {target_pos+1}")


# ═══════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════

def main():
    # ──── PART A: Rebalance S15-17 ────
    print("=" * 60)
    print("PART A: Rebalance S15-17 plots/cards")
    print(f"  Plot: {IMG_WIDTH}×{IMG_HEIGHT}″ (was 2.0″)")
    print(f"  Cards: T={CARDS_TOP}″, H={CARDS_H}″ (was T=3.25″, H=1.95″)")
    print(f"  Card fonts: S15=7pt, S16/S17=8pt")
    print("=" * 60)

    for ppt_path, label, s15, s16, s17 in [
        (PPT_EN, "EN", EN_S15, EN_S16, EN_S17),
        (PPT_ZH, "ZH", ZH_S15, ZH_S16, ZH_S17),
    ]:
        print(f"\n--- {label}: {ppt_path.name} ---")
        prs = Presentation(str(ppt_path))

        # S15-17 are at indices 14, 15, 16 (0-based)
        rebuild_slide(prs.slides[14], PNG_S15, s15, f"S15 {label}")
        rebuild_slide(prs.slides[15], PNG_S16, s16, f"S16 {label}")
        rebuild_slide(prs.slides[16], PNG_S17, s17, f"S17 {label}")

        prs.save(str(ppt_path))
        print(f"  Saved.")

    # ──── PART B: Insert Gene Sets Slide ────
    print(f"\n{'='*60}")
    print("PART B: Create & insert Gene Sets slide (after S6)")
    print("=" * 60)

    for ppt_path, lang in [(PPT_EN, "EN"), (PPT_ZH, "ZH")]:
        print(f"\n--- {lang}: {ppt_path.name} ---")
        prs = Presentation(str(ppt_path))

        # Create new slide (appended at end)
        create_gene_sets_slide(prs, lang)
        print(f"  Gene sets slide created (at end, pos {len(prs.slides)})")

        # Move from last position to position 6 (after S6, 0-indexed: S6 is index 5)
        insert_slide_at(prs, len(prs.slides) - 1, 6)

        prs.save(str(ppt_path))
        print(f"  Saved.")

    # ──── Verify ────
    print(f"\n{'='*60}")
    print("VERIFICATION — Slide order:")
    prs = Presentation(str(PPT_EN))
    for i, slide in enumerate(prs.slides):
        titles = [s for s in slide.shapes if hasattr(s, 'text_frame') and s.top < Inches(1.2)]
        title_text = titles[0].text_frame.text.strip()[:70] if titles else '(no title)'
        # Check for pic
        has_pic = any(hasattr(s, 'image') for s in slide.shapes)
        pic_note = " [PIC]" if has_pic else ""
        print(f"  S{i+1:2d}  {title_text}{pic_note}")

    print(f"\nALL DONE.")
    print(f"  Plot height: 2.0″ → 2.5″")
    print(f"  Card height: 1.95″ → 1.50″")
    print(f"  New S7: Automated Gene Set Detection (EN+ZH)")


if __name__ == "__main__":
    main()
