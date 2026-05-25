"""
Fix overlapping text on S7, S11, S16 in both EN and ZH PPTs.
"""

from pptx import Presentation
from pptx.util import Inches, Emu
from pathlib import Path

BASE = Path("results/figures_final")

# ============================================================
# S7 fix: "→ Fully automatic..." overlaps bottom of last HK item
# Move the arrow text and the bottom bar down by 0.20"
# ============================================================

def fix_s7(slide):
    """Move arrow text and bottom bar down slightly."""
    for s in slide.shapes:
        t = s.top / 914400
        l = s.left / 914400
        h = s.height / 914400
        w = s.width / 914400

        # Find arrow text boxes (→ Fully automatic... or → 全自动...)
        if hasattr(s, "text_frame"):
            txt = s.text_frame.text.strip()
            if txt.startswith("→"):
                # Move from ~5.05 to ~5.25
                s.top = Inches(5.25)

    # Move bottom bar from ~5.70 to ~5.90
    for s in slide.shapes:
        t = s.top / 914400
        if t > 5.60 and t < 5.80 and hasattr(s, "text_frame"):
            txt = s.text_frame.text.strip()
            if "Input" in txt or "输入" in txt:
                s.top = Inches(5.90)


# ============================================================
# S11 fix: legend "■ k_n (JS-HK) ■ k_f..." overlaps "k_n: 22"
# Move legend down by 0.50"
# ============================================================

def fix_s11(slide):
    """Move legend line down."""
    for s in slide.shapes:
        t = s.top / 914400
        if hasattr(s, "text_frame"):
            txt = s.text_frame.text.strip()
            if "■" in txt and ("k_n" in txt or "k_f" in txt):
                s.top = Inches(4.55)


# ============================================================
# S16 fix: right column cards overlap bar values
# Move cards right by 0.50", narrow slightly
# Move gradient summary down
# ============================================================

def fix_s16(slide):
    """Adjust right-column cards, gradient summary, and page number."""
    for s in slide.shapes:
        t = s.top / 914400
        l = s.left / 914400
        w = s.width / 914400

        # Page number "16": move down
        if hasattr(s, "text_frame"):
            txt = s.text_frame.text.strip()
            if txt == "16":
                s.top = Inches(5.60)

    for s in slide.shapes:
        t = s.top / 914400
        l = s.left / 914400
        w = s.width / 914400

        if hasattr(s, "text_frame"):
            txt = s.text_frame.text.strip()

            # Right column cards: move further right (clear of 121.8 at R=6.18)
            if txt in ("What the Gradient Tells Us", "梯度的含义", "Key Insight", "关键洞察"):
                s.left = Inches(6.30)
                s.width = Inches(3.20)
                continue

            # Right column card bodies
            if "Astrocytes show highest" in txt or "星形胶质细胞" in txt:
                s.left = Inches(6.30)
                s.width = Inches(3.20)
                continue

            if "Glial cells are not" in txt or "胶质细胞并非" in txt:
                s.left = Inches(6.30)
                s.width = Inches(3.20)
                continue

            # Gradient summary: move down slightly
            if "7.6x gradient" in txt or "7.6x " in txt:
                s.top = Inches(5.00)

            # Bottom summary bar: move down
            if "omega range:" in txt or "omega 范围" in txt:
                s.top = Inches(5.35)


# ============================================================
# Apply to both PPTs
# ============================================================

for tag, pptx_name in [
    ("EN", "CKI_Lecture_2026_v3.pptx"),
    ("ZH", "CKI_Lecture_2026_v3_ZH.pptx"),
]:
    path = BASE / pptx_name
    prs = Presentation(str(path))

    fix_s7(prs.slides[6])    # S7 (0-indexed)
    fix_s11(prs.slides[10])  # S11
    fix_s16(prs.slides[15])  # S16

    prs.save(str(path))
    print(f"{tag}: saved {pptx_name}")

# ============================================================
# Verify
# ============================================================

print()
print("=== Verification ===")
for tag, pptx_name in [
    ("EN", "CKI_Lecture_2026_v3.pptx"),
    ("ZH", "CKI_Lecture_2026_v3_ZH.pptx"),
]:
    prs = Presentation(str(BASE / pptx_name))
    issues = 0
    for idx in [6, 10, 15]:
        slide = prs.slides[idx]
        text_boxes = [
            (s.top / 914400, s.left / 914400,
             s.height / 914400, s.width / 914400,
             s.text_frame.text.strip()[:40])
            for s in slide.shapes
            if hasattr(s, "text_frame") and s.text_frame.text.strip()
        ]
        slide_issues = 0
        for i in range(len(text_boxes)):
            for j in range(i + 1, len(text_boxes)):
                t1, l1, h1, w1, txt1 = text_boxes[i]
                t2, l2, h2, w2, txt2 = text_boxes[j]
                if l1 < l2 + w2 and l1 + w1 > l2 and t1 < t2 + h2 and t1 + h1 > t2:
                    ol = max(l1, l2)
                    or_ = min(l1 + w1, l2 + w2)
                    ot = max(t1, t2)
                    ob = min(t1 + h1, t2 + h2)
                    ow, oh = or_ - ol, ob - ot
                    if ow > 0.1 and oh > 0.1:
                        slide_issues += 1
        status = "✗" if slide_issues else "✓"
        print(f"  {tag} S{idx+1}: {status} ({slide_issues} overlaps)")
        issues += slide_issues

    print(f"  {tag} total issues: {issues}")
