"""
62_fix_s18_s19_png.py - S18/S19 破损PNG修复（HTTP服务器方式 + 注入PPTX）
Chrome headless + HTTP服务器渲染SVG→PNG (2x DPI)，然后注入v4 PPTX。
"""
import subprocess, os, sys, threading, time, http.server, socketserver
from pathlib import Path

CHROME = r"C:\Program Files\Google\Chrome\Application\chrome.exe"
BASE = Path(r"C:\Users\KnightZ\Desktop\细胞受选择\results\figures_final")

TMP_DIR = BASE / "_tmp_http_render"
TMP_DIR.mkdir(exist_ok=True)

# ============================================================
# Step 1: 创建HTML文件（包装SVG）
# ============================================================
print("=" * 60)
print("Step 1: Creating HTML wrappers for SVGs")
print("=" * 60)

svg_pairs = [
    ("_s16_anomaly_v3.svg", "_s16_anomaly_v3.png", "S18 chart"),
    ("_s17_migration_v3.svg", "_s17_migration_v3.png", "S19 chart"),
]

html_files = []
for svg_name, png_name, label in svg_pairs:
    svg_path = BASE / svg_name
    if not svg_path.exists():
        print(f"  SKIP: {svg_name} not found")
        continue
    svg = svg_path.read_text(encoding="utf-8")
    html = (
        '<!DOCTYPE html><html><head><meta charset="utf-8">'
        '<style>body{margin:0;padding:0;background:white;overflow:hidden;}</style>'
        '</head><body>' + svg + '</body></html>'
    )
    html_path = TMP_DIR / f"{svg_name}.html"
    html_path.write_text(html, encoding="utf-8")
    html_files.append((html_path, png_name, label))
    print(f"  {label}: {html_path.name} ({len(html)} chars)")

# ============================================================
# Step 2: 启动HTTP服务器
# ============================================================
print("\n" + "=" * 60)
print("Step 2: Starting HTTP server")
print("=" * 60)

PORT = 18848
os.chdir(str(TMP_DIR))

class QuietHTTPHandler(http.server.SimpleHTTPRequestHandler):
    def log_message(self, format, *args):
        pass  # suppress logs

httpd = socketserver.TCPServer(("", PORT), QuietHTTPHandler)
server_thread = threading.Thread(target=httpd.serve_forever, daemon=True)
server_thread.start()
time.sleep(0.5)
print(f"  HTTP server running on port {PORT}")

# ============================================================
# Step 3: Chrome headless 截图
# ============================================================
print("\n" + "=" * 60)
print("Step 3: Rendering PNGs via Chrome headless")
print("=" * 60)

WIDTH, HEIGHT = 1220, 520

for html_path, png_name, label in html_files:
    html_name = html_path.name
    url = f"http://localhost:{PORT}/{html_name}"
    
    png_out = BASE / png_name
    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        f"--window-size={WIDTH},{HEIGHT}",
        f"--screenshot={png_out}",
        "--force-device-scale-factor=2",
        url,
    ]
    
    print(f"  {label}: {url}")
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
    
    if png_out.exists():
        sz = png_out.stat().st_size
        print(f"    -> {png_name}: {sz:,} bytes ({sz//1024} KB)")
    else:
        print(f"    FAILED: {png_name}")
        print(f"    stderr: {result.stderr[:300]}")

# ============================================================
# Step 4: 关闭服务器
# ============================================================
print("\n" + "=" * 60)
print("Step 4: Shutting down HTTP server")
print("=" * 60)
httpd.shutdown()
server_thread.join(timeout=2)
# Clean up temp dir
import shutil
shutil.rmtree(str(TMP_DIR), ignore_errors=True)
print("  Done.")

# ============================================================
# Step 5: 注入PNG到PPTX v4文件
# ============================================================
print("\n" + "=" * 60)
print("Step 5: Injecting PNGs into PPTX v4 files")
print("=" * 60)

from pptx import Presentation

EN_PPT = BASE / "CKI_Lecture_2026_v4.pptx"
ZH_PPT = BASE / "CKI_Lecture_2026_v4_ZH.pptx"

# S18 = slide 18 (0-indexed 17), uses _s16_anomaly_v3.png
# S19 = slide 19 (0-indexed 18), uses _s17_migration_v3.png
slide_png_map = {
    17: "_s16_anomaly_v3.png",   # S18
    18: "_s17_migration_v3.png",  # S19
}

for tag, pptx_path in [("EN", EN_PPT), ("ZH", ZH_PPT)]:
    print(f"\n  --- {tag} ---")
    prs = Presentation(str(pptx_path))
    
    for slide_idx, png_name in slide_png_map.items():
        png_path = BASE / png_name
        if not png_path.exists():
            print(f"    S{slide_idx+1}: {png_name} not found, skip")
            continue
        
        slide = prs.slides[slide_idx]
        found = False
        for s in slide.shapes:
            if hasattr(s, "image"):
                old_left = s.left
                old_top = s.top
                old_width = s.width
                old_height = s.height
                # Remove old pic element
                s._element.getparent().remove(s._element)
                # Add new pic at same position/size
                slide.shapes.add_picture(
                    str(png_path), old_left, old_top, old_width, old_height)
                print(f"    S{slide_idx+1}: replaced PIC with {png_name} ({png_path.stat().st_size//1024} KB)")
                found = True
                break
        
        if not found:
            print(f"    S{slide_idx+1}: no image shape found to replace")
    
    prs.save(str(pptx_path))
    print(f"    {tag}: saved ({pptx_path.stat().st_size//1024} KB)")

print("\n" + "=" * 60)
print("DONE. All S18/S19 PNGs re-rendered and injected.")
print("=" * 60)
